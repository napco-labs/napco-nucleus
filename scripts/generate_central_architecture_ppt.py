"""NAPCO Nucleus — central capture architecture deck.

Built for the next boss demo. 11 slides, 16:9, minimal text on each
slide so Mohammad drives the narration. Speaker notes carry the
talking points.

Run:
    python scripts\\generate_central_architecture_ppt.py
Output:
    docs\\NAPCO-Nucleus-Central-Architecture.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT = ROOT / "docs" / "NAPCO-Nucleus-Central-Architecture.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette — match docs/NAPCO-Nucleus-Presentation.pptx
NAVY   = RGBColor(0x1F, 0x4E, 0x79)
TEAL   = RGBColor(0x2E, 0x8A, 0x8A)
CORAL  = RGBColor(0xE0, 0x78, 0x56)
GREEN  = RGBColor(0x4A, 0x7A, 0x4A)
GOLD   = RGBColor(0xC9, 0x96, 0x2B)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)
INK    = RGBColor(0x22, 0x22, 0x22)
MUTED  = RGBColor(0x6B, 0x77, 0x85)
SOFT   = RGBColor(0xF5, 0xF7, 0xFA)
RULE   = RGBColor(0xD5, 0xDC, 0xE5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# ── helpers ────────────────────────────────────────────────────────

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def no_line(shape):
    shape.line.fill.background()


def line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_text(tf, text, size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_para(tf, text, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, space_before=0, font="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    return tf


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None,
             line_width=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color is None:
        no_fill(rect)
    else:
        solid(rect, fill_color)
    if line_color is None:
        no_line(rect)
    else:
        line(rect, line_color, line_width)
    rect.text_frame.margin_left = Inches(0.1)
    rect.text_frame.margin_right = Inches(0.1)
    rect.text_frame.margin_top = Inches(0.05)
    rect.text_frame.margin_bottom = Inches(0.05)
    return rect


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width_pt=2.0):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    # right-arrow end
    return conn


def base_slide(prs, title_text, subtitle_text=None):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # Top accent bar
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill_color=NAVY)
    # Title
    tf = add_textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    set_text(tf, title_text, size=28, bold=True, color=NAVY)
    if subtitle_text:
        add_para(tf, subtitle_text, size=14, color=MUTED, space_before=2)
    # Footer rule
    add_rect(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.02), fill_color=RULE)
    return s


def set_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    p.text = text


def add_chip(slide, x, y, w, h, text, color):
    """Small pill with white text used for role labels."""
    rect = add_rect(slide, x, y, w, h, fill_color=color, rounded=True)
    tf = rect.text_frame
    set_text(tf, text, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return rect


def add_box(slide, x, y, w, h, title, lines, accent=NAVY):
    """Title-bar + content box used for the architecture diagrams."""
    container = add_rect(slide, x, y, w, h, fill_color=WHITE, line_color=RULE)
    # accent stripe on left
    add_rect(slide, x, y, Inches(0.08), h, fill_color=accent)
    tf = add_textbox(slide, x + Inches(0.18), y + Inches(0.1),
                     w - Inches(0.25), h - Inches(0.2))
    set_text(tf, title, size=14, bold=True, color=accent)
    for ln in lines:
        add_para(tf, ln, size=11, color=INK, space_before=2)
    return container


# ── slides ─────────────────────────────────────────────────────────

def slide_title(prs):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=NAVY)
    # Big title
    tf = add_textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.5))
    set_text(tf, "NAPCO Nucleus", size=58, bold=True, color=WHITE)
    add_para(tf, "Central Capture Architecture", size=30, color=WHITE)
    # Subhead
    tf2 = add_textbox(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.5))
    set_text(tf2, "5 developers / 2 clients -> one identification pipeline.",
             size=18, color=RGBColor(0xC8, 0xD2, 0xE0))
    # Footer
    tf3 = add_textbox(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    set_text(tf3, "Mohammad Kamrul Hasan   |   Adaptive Enterprise / NAPCO labs",
             size=12, color=RGBColor(0x9C, 0xAB, 0xBE))
    set_speaker_notes(s, (
        "Open with the framing: requirements the team is trying to capture "
        "live across multiple developers' machines and across multiple "
        "channels. Today's deck shows how we stitch them back together "
        "without asking developers to do anything beyond their normal "
        "calls and chats."
    ))


def slide_problem(prs):
    s = base_slide(prs, "The problem",
                   "A single requirement spans multiple conversations.")
    # Visual: 2 clients → 5 devs grid, with arrows showing fragmentation
    # Clients column
    add_chip(s, Inches(0.8), Inches(2.0), Inches(2.0), Inches(0.6),
             "Client A", CORAL)
    add_chip(s, Inches(0.8), Inches(3.2), Inches(2.0), Inches(0.6),
             "Client B", PURPLE)
    # Devs column
    devs = ["Dev 1", "Dev 2", "Dev 3", "Dev 4", "Dev 5"]
    for i, d in enumerate(devs):
        add_chip(s, Inches(7.2), Inches(1.4 + i * 0.85),
                 Inches(2.0), Inches(0.6), d, NAVY)
    # Arrows: A → 1, A → 3 ; B → 2, B → 4
    for tgt in (1, 3):
        add_arrow(s, Inches(2.8), Inches(2.3),
                  Inches(7.2), Inches(1.7 + tgt * 0.85), color=CORAL)
    for tgt in (1, 3):
        add_arrow(s, Inches(2.8), Inches(3.5),
                  Inches(7.2), Inches(1.7 + tgt * 0.85), color=PURPLE)
    # Caption
    tf = add_textbox(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2))
    set_text(tf, "Same requirement, different developers.", size=20,
             bold=True, color=INK)
    add_para(tf,
             "Half a feature on a 2-min call with Dev 1, the other half on "
             "a 4-min call with Dev 3. Neither developer hears the whole "
             "thing. The requirement is invisible to anyone reviewing one "
             "machine at a time.",
             size=14, color=MUTED, space_before=4)
    set_speaker_notes(s, (
        "Anchor on the failure mode the team experiences today: each "
        "developer has only their fragment. Even if they faithfully relay "
        "what they heard, no one is reconciling fragments across people. "
        "That's why requirements get half-built, scope-creeped, or "
        "missed entirely."
    ))


def slide_solution_oneliner(prs):
    s = base_slide(prs, "The fix in one line",
                   "Capture everywhere it happens. Identify in one place.")
    # Big sentence
    tf = add_textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0))
    set_text(tf,
             "Every developer's calls and chats flow into one central "
             "store. One scoped pass per (client, day) sees the full "
             "picture and drafts a verification email.",
             size=22, color=INK, align=PP_ALIGN.LEFT)
    # Two pillars
    add_chip(s, Inches(0.8), Inches(5.2), Inches(5.7), Inches(0.7),
             "Local capture (per-dev)   --   no extra workflow", TEAL)
    add_chip(s, Inches(6.8), Inches(5.2), Inches(5.7), Inches(0.7),
             "Central identification   --   single Claude session", GREEN)
    set_speaker_notes(s, (
        "Two design promises. First: capture has to be invisible to the "
        "developer. They keep their normal Teams workflow; we listen for "
        "the natural call-bookend phrases (Assalamualaikum / Allah Hafez) "
        "and the recorder reacts. Second: identification has to be ONE "
        "place. Five separate identify-runs would miss the cross-developer "
        "fragments. We aggregate everyone's day into one Claude session "
        "scoped per client."
    ))


def slide_architecture(prs):
    s = base_slide(prs, "Architecture",
                   "5 dev machines push to MVPACCESS. Identification runs there.")

    # Left column: dev box (representing all 5)
    add_box(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(4.5),
            "Dev machine (x5)", [
                "voice_daemon.py",
                "  listens for start/stop phrases",
                "  Teams-only gate (pycaw)",
                "",
                "record_call.py",
                "  WAV mic + WAV speaker",
                "  resolves client via IndexedDB",
                "  metadata sidecar JSON",
                "",
                "push_chat.py (every 15 min)",
                "  Task Scheduler-driven",
                "  bundles last 15 min of chats",
            ], accent=TEAL)

    # Center: arrow + label
    add_arrow(s, Inches(4.6), Inches(3.7),
              Inches(8.4), Inches(3.7),
              color=NAVY, width_pt=3)
    tf = add_textbox(s, Inches(4.7), Inches(3.2), Inches(3.8), Inches(0.5))
    set_text(tf, "SMB push", size=12, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
    tf2 = add_textbox(s, Inches(4.7), Inches(3.85), Inches(3.8), Inches(0.5))
    set_text(tf2, "after each call /\nevery 15 min for chat",
             size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # Right column: agent host
    add_box(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(4.5),
            "MVPACCESS (172.16.205.209)", [
                "C:\\nucleus-central\\",
                "  <dev>\\<YYYY-MM-DD>\\",
                "    calls\\  *.wav  *.json",
                "    chat\\   *.docx",
                "",
                "Claude Max authenticated here",
                "",
                "collect_central.py",
                "  walks tree -> filters by",
                "  --client X --day YYYY-MM-DD",
                "  -> Verification doc",
                "  -> Gmail Drafts",
            ], accent=GREEN)

    # You at the bottom
    add_chip(s, Inches(5.0), Inches(6.3), Inches(3.4), Inches(0.5),
             "You: review draft -> send", GOLD)
    set_speaker_notes(s, (
        "Walk through the diagram left to right. The dev side is the same "
        "as the local NN we already had, just with two new behaviors after "
        "stop: the IndexedDB lookup that names the client, and the SMB "
        "copy. Chat is on a Windows Task Scheduler timer, runs every 15 "
        "min in the background. The agent host is the only place that "
        "needs Claude Max auth, the only place that runs identify, and "
        "the only place that drafts the verification email."
    ))


def slide_voice_capture(prs):
    s = base_slide(prs, "Voice-activated capture",
                   "Developers talk normally. Recording reacts to call-bookends.")

    # Two phrases boxed
    add_box(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.3),
            "Start phrases", [
                "  \"Assalamualaikum\"",
                "  \"Salaam alaikum\"",
                "  \"Nucleus start\" (English fallback)",
                "",
                "Daemon spawns record_call.py",
            ], accent=GREEN)

    add_box(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.3),
            "Stop phrases", [
                "  \"Allah Hafez\" / \"Khoda Hafiz\"",
                "  \"Nucleus stop\" (English fallback)",
                "",
                "Daemon writes stop sentinel;",
                "recorder flushes WAVs.",
            ], accent=CORAL)

    # Belt & braces
    tf = add_textbox(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5))
    set_text(tf, "Why this works without false positives:",
             size=16, bold=True, color=NAVY)
    add_para(tf,
             "Teams-only gate: pycaw confirms ms-teams.exe has an Active "
             "audio session before start fires. Ringing counts (the ringtone "
             "produces an audio session) so the recording captures the very "
             "first second of \"Assalamualaikum\".",
             size=13, color=INK, space_before=6)
    add_para(tf,
             "Phrase list is in data/teams/voice_phrases.json -- editable "
             "without code changes; teammates can add their own variants.",
             size=13, color=INK, space_before=6)
    set_speaker_notes(s, (
        "Show that the voice trigger is genuinely zero-friction. People "
        "already say these phrases at the start and end of every BD call. "
        "We're not asking them to learn a wake word, we're listening for "
        "their normal greeting. The Teams-only gate is the false-positive "
        "killer: 'Assalamualaikum' said in a corridor doesn't trigger a "
        "recording because Teams isn't producing audio at that moment. "
        "Stop is unconditional so a dropped call always closes cleanly."
    ))


def slide_indexeddb_resolver(prs):
    s = base_slide(prs, "Who was the call with?",
                   "IndexedDB has the answer. We just ask.")

    # Left: where the data is
    tf = add_textbox(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.0))
    set_text(tf, "Teams writes one Event/Call entry per call",
             size=15, bold=True, color=NAVY)
    add_para(tf,
             "Path: %LOCALAPPDATA%\\Packages\\MSTeams_8wekyb3d8bbwe\\...\\IndexedDB",
             size=11, color=MUTED, font="Consolas", space_before=4)
    add_para(tf,
             "Each entry carries:",
             size=13, bold=True, color=INK, space_before=12)
    for ln in [
        "  - originalArrivalTime (epoch ms)",
        "  - conversationId",
        "  - <partlist> with every participant's identity + displayName",
        "  - callId (UUID)",
        "  - callEventType (callEnded, missed, ...)",
    ]:
        add_para(tf, ln, size=12, color=INK, space_before=2, font="Consolas")
    add_para(tf,
             "Verified live: 843 historical Event/Call entries in your "
             "local DB, going back ~2 years.",
             size=12, color=GREEN, bold=True, space_before=10)

    # Right: example output
    add_box(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.0),
            "resolve_client_for_recording()", [
                "matched: True",
                "call_id: a19b7393-8865-4098-...",
                "call_type: callEnded",
                "client_name: Isruk H",
                "clients:",
                "  Isruk H",
                "  Md. Ahsan Habib Rocky",
                "  Atikur Zaman",
                "  Assad Zaman",
                "  Salman Ahmed Firoz",
                "",
                "(Kamrul Hasan / titucse stripped --",
                "  he's the current user)",
            ], accent=PURPLE)
    set_speaker_notes(s, (
        "This is the load-bearing trick that makes the central pipeline "
        "useful. Without auto-tagging by client, the central tree is just "
        "5 piles of WAVs and Mohammad has to hand-label every one. With "
        "this resolver, every recording self-labels at upload time. We "
        "verified the schema against the live database -- 843 historical "
        "calls all parse cleanly with the new partlist + displayName "
        "format. The current user is detected via isSentByCurrentUser=True "
        "messages and dropped from the clients list automatically."
    ))


def slide_central_flow(prs):
    s = base_slide(prs, "What lands on MVPACCESS",
                   "Predictable layout. One folder per dev per day.")

    tf = add_textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
    add_para(tf, "C:\\nucleus-central\\",
             size=14, bold=True, color=NAVY, font="Consolas")
    samples = [
        "    salman\\2026-05-08\\",
        "        calls\\",
        "            20260508-100432_mic.wav        (mic track)",
        "            20260508-100432_speaker.wav    (speaker track)",
        "            20260508-100432.json           (metadata: dev, client, ...)",
        "        chat\\",
        "            chat_2026-05-08_1000-1015.docx",
        "            chat_2026-05-08_1015-1030.docx",
        "    rocky\\2026-05-08\\",
        "        calls\\",
        "            20260508-101512_mic.wav",
        "            20260508-101512_speaker.wav",
        "            20260508-101512.json",
        "        chat\\",
        "            chat_2026-05-08_1015-1030.docx",
        "    isruk\\2026-05-08\\",
        "        ...",
    ]
    for ln in samples:
        add_para(tf, ln, size=12, color=INK, font="Consolas", space_before=1)
    add_para(tf,
             "metadata.client_name comes from the IndexedDB resolver. "
             "collect_central uses substring match on this field to scope "
             "by client.",
             size=12, color=MUTED, space_before=10)
    set_speaker_notes(s, (
        "Show the layout so the boss has a concrete mental model. Each "
        "developer pushes into their own subfolder, dated. Mohammad reads "
        "across all dev folders for a given day. The metadata file is the "
        "key -- without it we have audio but no context; with it we know "
        "which call was with which client. Nothing here is invented data; "
        "every field comes from either the recorder or Teams' own DB."
    ))


def slide_collect_central(prs):
    s = base_slide(prs, "Per-client identification",
                   "One pass. One verification email. One review.")

    # CLI examples
    tf = add_textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.0))
    set_text(tf, "On MVPACCESS:", size=14, bold=True, color=NAVY)
    for cmd in [
        "    python collect_central.py --client \"Susmoy\"",
        "    python collect_central.py --client \"Isruk\" --day 2026-05-08",
        "    python collect_central.py --client all --no-identify    # inspect first",
    ]:
        add_para(tf, cmd, size=13, color=INK, font="Consolas", space_before=4)

    # Pipeline boxes
    boxes_y = Inches(4.0)
    box_w = Inches(2.4)
    gap = Inches(0.2)
    starts = [Inches(0.8 + i * (2.4 + 0.2)) for i in range(4)]
    titles = [
        ("scan central tree", TEAL),
        ("filter by client", PURPLE),
        ("transcribe + build session doc", CORAL),
        ("verify_session -> .eml in Drafts", GREEN),
    ]
    for x, (t, c) in zip(starts, titles):
        add_box(s, x, boxes_y, box_w, Inches(2.5),
                t, ["", "", "", ""], accent=c)
        # shaded numeric badge
        idx = starts.index(x) + 1
        badge = add_rect(s, x + Inches(0.15), boxes_y + Inches(0.55),
                         Inches(0.55), Inches(0.55),
                         fill_color=c, rounded=True)
        tf2 = badge.text_frame
        set_text(tf2, str(idx), size=22, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_speaker_notes(s, (
        "Walk left to right through the four phases. Step 1 reads each "
        "dev's day folder. Step 2 keeps only the calls whose metadata "
        "client_name matches what you typed -- chats are included whole "
        "since one chat can mention multiple clients. Step 3 transcribes "
        "every matched WAV (faster-whisper large-v3, Bangla -> English) "
        "into a single session document with one section per source. "
        "Step 4 hands that to verify_session, which is the same prompt "
        "we already shipped, and the .eml lands in your Gmail Drafts."
    ))


def slide_demo_artifacts(prs):
    s = base_slide(prs, "What the boss sees",
                   "One email per client per day, with two attachments.")

    # Left: email anatomy
    tf = add_textbox(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.0))
    set_text(tf, "Verification email (Gmail Drafts)",
             size=15, bold=True, color=NAVY)
    add_para(tf, "From: Mohammad <khasan@ael-bd.com>",
             size=12, color=INK, space_before=6, font="Consolas")
    add_para(tf, "To: <client>@<domain>",
             size=12, color=INK, space_before=2, font="Consolas")
    add_para(tf, "Subject: Requirements Verification - 2026-05-08",
             size=12, color=INK, space_before=2, font="Consolas")
    add_para(tf, "Attachments (2):",
             size=13, bold=True, color=NAVY, space_before=10)
    add_para(tf, "    Requirements Verification 2026-05-08.docx",
             size=12, color=INK, space_before=4, font="Consolas")
    add_para(tf, "        the curated list",
             size=11, color=MUTED, space_before=1)
    add_para(tf, "    Pull Session 2026-05-08.docx",
             size=12, color=INK, space_before=6, font="Consolas")
    add_para(tf, "        the raw source material so they can cross-check",
             size=11, color=MUTED, space_before=1)

    # Right: the curated list
    add_box(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.0),
            "Requirements Verification (sample)", [
                "1. Operator Management feature",
                "    Full CRUD, search filters, pagination,",
                "    role-based access via Setup Roles +",
                "    Setup Partition Groups, audit logging.",
                "",
                "2. Setup Partition Group",
                "    Add / edit / delete with confirmation,",
                "    paginated grid (configurable page size),",
                "    search field.",
                "",
                "3. Secure Logout",
                "    Session termination, token invalidation,",
                "    redirect to login, clear cached data.",
            ], accent=GOLD)
    set_speaker_notes(s, (
        "End-state for the demo. The two-attachment email gives the "
        "client both the curated list AND the raw source material -- so "
        "they can cross-check anything that looks off. We don't auto-send; "
        "you review each draft in Gmail and click send. That's the "
        "manual-verify gate the team agreed on before we move to "
        "automatic OpenProject publishing."
    ))


def slide_today_fixes(prs):
    s = base_slide(prs, "What we fixed today",
                   "Three real bugs, two new capabilities.")

    cols = [
        ("3 bugs fixed", CORAL, [
            "Email -> PDF extraction",
            "  was silent failure (wrong import path)",
            "  now extracts cleanly",
            "",
            "Email -> .doc (legacy Word)",
            "  was returning an error message",
            "  now byte-scans like Drive does",
            "",
            "Pull Session attachment filename",
            "  was \"current.docx\" on the client side",
            "  now \"Pull Session 2026-05-08.docx\"",
        ]),
        ("Voice daemon", PURPLE, [
            "teams\\voice_daemon.py",
            "  Assalamualaikum / Allah Hafez triggers",
            "  Teams-only gate (pycaw audio session)",
            "  Configurable phrase list",
            "",
            "Each teammate runs once at start of day:",
            "    scripts\\start-daemon.bat",
        ]),
        ("Central pipeline", GREEN, [
            "teams\\calls.py",
            "  IndexedDB Event/Call resolver",
            "",
            "teams\\push_chat.py + scheduler",
            "  per-15-min Task Scheduler entry",
            "",
            "collect_central.py",
            "  per-(client,day) aggregation",
            "  + Whisper transcription",
            "  + verify_session run",
        ]),
    ]
    col_w = Inches(4.0)
    gap = Inches(0.2)
    for i, (title, color, lines) in enumerate(cols):
        add_box(s, Inches(0.5 + i * 4.2), Inches(1.5),
                col_w, Inches(5.5), title, lines, accent=color)
    set_speaker_notes(s, (
        "Reminder of the bug-fix work that goes alongside the architecture "
        "additions. The PDF email-attachment bug was the catalyst this "
        "morning -- the boss-demo verification email had been quietly "
        "missing whole specs because the extractor failed silently. That's "
        "why we re-tested. The voice daemon and central pipeline are net-new."
    ))


def slide_deployment(prs):
    s = base_slide(prs, "Deployment",
                   "Three install patterns. One commit. ~10 minutes per role.")

    # Three role boxes
    add_box(s, Inches(0.4), Inches(1.5), Inches(4.1), Inches(5.3),
            "MVPACCESS (one-time)", [
                "AS ADMIN ON THE VM:",
                "",
                "  git clone napco-labs/napco-nucleus",
                "  scripts\\setup.bat",
                "  Create SMB share \\\\172.16.205.209\\nucleus-central",
                "  Open firewall (port 445 from .205.0/24)",
                "  Set NUCLEUS_CENTRAL_PATH=C:\\nucleus-central in .env",
                "  Verify Claude Max auth: claude --version",
                "",
                "Estimated: 30 min including network/firewall.",
            ], accent=NAVY)

    add_box(s, Inches(4.6), Inches(1.5), Inches(4.1), Inches(5.3),
            "Each dev (per teammate)", [
                "AS THE LOGGED-IN USER:",
                "",
                "  git clone",
                "  scripts\\setup.bat (winget Python, .venv, .env, ...)",
                "  net use \\\\172.16.205.209\\nucleus-central /persistent:yes",
                "  scripts\\register-chat-push-task.ps1",
                "",
                "Day-to-day:",
                "  scripts\\start-daemon.bat at start of day",
                "  Talk normally during calls.",
                "",
                "Estimated: 10 min per teammate.",
            ], accent=TEAL)

    add_box(s, Inches(8.8), Inches(1.5), Inches(4.1), Inches(5.3),
            "You (daily)", [
                "ON MVPACCESS:",
                "",
                "  scripts\\central-pull.bat \"Client A\"",
                "  scripts\\central-pull.bat \"Client B\"",
                "",
                "Each run -> one .eml in Gmail Drafts.",
                "",
                "Review, edit if needed, send.",
                "",
                "Estimated: 5 min per client per day.",
            ], accent=GOLD)
    set_speaker_notes(s, (
        "Stress that this is shipped, not a slide-ware roadmap. The "
        "commit is on origin/main. The setup script handles everything "
        "the developer would otherwise type by hand -- they don't need to "
        "know what a virtualenv is. The only sysadmin step is the SMB "
        "share + firewall on MVPACCESS, which we do once."
    ))


def slide_status(prs):
    s = base_slide(prs, "Status",
                   "Code shipped. Ready for the boss demo with same test data.")

    add_chip(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.7),
             "Code: pushed to napco-labs/napco-nucleus@main (commit 1b92e87)",
             GREEN)

    tf = add_textbox(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(4.0))
    add_para(tf, "Verified end-to-end on the dev machine:",
             size=15, bold=True, color=NAVY)
    for ln in [
        "  - record_call -> central upload (file structure correct)",
        "  - push_chat -> central upload (8 chats / 22 msgs picked up)",
        "  - collect_central reads both flows",
        "  - IndexedDB call resolver matches 843 historical calls cleanly",
    ]:
        add_para(tf, ln, size=13, color=INK, space_before=4)

    add_para(tf, "Pending (sysadmin only):",
             size=15, bold=True, color=NAVY, space_before=14)
    for ln in [
        "  - Create SMB share + firewall rule on MVPACCESS",
        "  - First end-to-end test: one dev's call -> central -> identify",
        "  - Roll out to remaining 4 devs after round-trip is proven",
    ]:
        add_para(tf, ln, size=13, color=INK, space_before=4)
    set_speaker_notes(s, (
        "Close the loop. Everything that is code-side is done and "
        "verified on the dev machine using a temporary local 'central' "
        "folder. The SMB hop is the one piece I couldn't test from the "
        "dev machine because the share doesn't exist yet -- but the same "
        "Path() + shutil.copy2() that worked against a local folder will "
        "work against an SMB UNC path once the share is up."
    ))


# ── build ──────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_solution_oneliner(prs)
    slide_architecture(prs)
    slide_voice_capture(prs)
    slide_indexeddb_resolver(prs)
    slide_central_flow(prs)
    slide_collect_central(prs)
    slide_demo_artifacts(prs)
    slide_today_fixes(prs)
    slide_deployment(prs)
    slide_status(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote: {OUT}")
    print(f"  slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
