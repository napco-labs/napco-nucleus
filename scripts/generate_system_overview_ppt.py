"""NAPCO Nucleus - end-to-end system-overview deck.

Stakeholder-facing explainer aimed at non-technical viewers
(executives, sponsors, boss demo). Each slide is readable in
about 5 seconds; technical detail lives in speaker notes.

Narrative arc: PROBLEM -> WHAT WE BUILT -> HOW IT WORKS -> WHAT'S NEXT.
No personal names appear in slide-visible text. Jargon is restricted
to speaker notes.

Self-contained: every helper that the deck needs is in this file.

Run:
    python scripts\\generate_system_overview_ppt.py
Output:
    docs\\NAPCO-Nucleus-System-Overview.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT = ROOT / "docs" / "NAPCO-Nucleus-System-Overview.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette - match docs/NAPCO-Nucleus-Presentation.pptx
NAVY   = RGBColor(0x1F, 0x4E, 0x79)
TEAL   = RGBColor(0x2E, 0x8A, 0x8A)
CORAL  = RGBColor(0xE0, 0x78, 0x56)
GREEN  = RGBColor(0x4A, 0x7A, 0x4A)
GOLD   = RGBColor(0xC9, 0x96, 0x2B)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)
INK    = RGBColor(0x22, 0x22, 0x22)
MUTED  = RGBColor(0x6B, 0x77, 0x85)
SOFT   = RGBColor(0xF5, 0xF7, 0xFA)
RULE   = RGBColor(0xD5, 0xDC, 0xE5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# -- helpers --------------------------------------------------------

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def no_line(shape):
    shape.line.fill.background()


def line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_text(tf, text, size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_para(tf, text, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, space_before=0, font="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    return tf


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None,
             line_width=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color is None:
        no_fill(rect)
    else:
        solid(rect, fill_color)
    if line_color is None:
        no_line(rect)
    else:
        line(rect, line_color, line_width)
    rect.text_frame.margin_left = Inches(0.12)
    rect.text_frame.margin_right = Inches(0.12)
    rect.text_frame.margin_top = Inches(0.08)
    rect.text_frame.margin_bottom = Inches(0.08)
    return rect


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width_pt=2.0):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn


def base_slide(prs, title_text, subtitle_text=None,
               title_size=36, subtitle_size=16):
    """Standard slide with bigger, more breathable headings."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # Top accent bar
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill_color=NAVY)
    # Title
    tf = add_textbox(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.9))
    set_text(tf, title_text, size=title_size, bold=True, color=NAVY)
    if subtitle_text:
        add_para(tf, subtitle_text, size=subtitle_size, color=MUTED,
                 space_before=4)
    # Footer rule
    add_rect(s, Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.02),
             fill_color=RULE)
    return s


def set_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    p.text = text


def add_chip(slide, x, y, w, h, text, color, font_size=14):
    """Pill with white text, used for labeled roles / channels."""
    rect = add_rect(slide, x, y, w, h, fill_color=color, rounded=True)
    tf = rect.text_frame
    set_text(tf, text, size=font_size, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return rect


def add_box(slide, x, y, w, h, title, lines, accent=NAVY,
            title_size=16, body_size=13):
    """Title-bar + content box for grouped concepts."""
    container = add_rect(slide, x, y, w, h, fill_color=WHITE, line_color=RULE)
    add_rect(slide, x, y, Inches(0.08), h, fill_color=accent)
    tf = add_textbox(slide, x + Inches(0.22), y + Inches(0.12),
                     w - Inches(0.32), h - Inches(0.24))
    set_text(tf, title, size=title_size, bold=True, color=accent)
    for ln in lines:
        add_para(tf, ln, size=body_size, color=INK, space_before=3)
    return container


# -- slides ---------------------------------------------------------

def slide_title(prs):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=NAVY)

    tf = add_textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6))
    set_text(tf, "NAPCO Nucleus", size=64, bold=True, color=WHITE)

    tf2 = add_textbox(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.4))
    set_text(tf2,
             "Turning client conversations into verified requirements,",
             size=26, color=RGBColor(0xC8, 0xD2, 0xE0))
    add_para(tf2, "automatically.",
             size=26, color=RGBColor(0xC8, 0xD2, 0xE0))

    tf3 = add_textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
    set_text(tf3, "2026-05-14   |   Adaptive Enterprise / NAPCO labs",
             size=12, color=RGBColor(0x9C, 0xAB, 0xBE))

    set_speaker_notes(s, (
        "Set the frame. Nucleus is the system that takes everything the "
        "team says, types, mails or shares with a client over the course "
        "of a day and -- with one human review step -- produces a single "
        "verification email back to the client confirming what we heard. "
        "It runs hands-off; developers do not change their workflow. "
        "The central host moved from the Windows MVPACCESS box to a "
        "Linux docker stack on .123 on 2026-05-14. The deck explains "
        "the whole pipeline end-to-end."
    ))


def slide_problem(prs):
    s = base_slide(prs, "The problem",
                   "Client needs arrive everywhere at once. Without help, things slip.")

    pains = [
        ("Scattered", CORAL,
         "Clients raise needs across calls, chats,\nemails, and shared files."),
        ("Fragmented", PURPLE,
         "Each developer hears only their slice.\nNo one sees the whole picture."),
        ("Lost", GOLD,
         "Requirements get missed, half-captured,\nor re-asked weeks later."),
    ]
    cw = Inches(3.9)
    ch = Inches(3.0)
    gap = Inches(0.25)
    total = 3 * 3.9 + 2 * 0.25
    x0 = (13.333 - total) / 2
    y = Inches(2.0)
    for i, (title, color, body) in enumerate(pains):
        x = Inches(x0 + i * (3.9 + 0.25))
        add_box(s, x, y, cw, ch, title, body.split("\n"),
                accent=color, title_size=22, body_size=14)

    tf = add_textbox(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.2))
    set_text(tf, "Nothing was centralized. Nothing was reconciled.",
             size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_para(tf,
             "The cost was invisible -- until a deliverable shipped without "
             "what the client thought they had asked for.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER, space_before=8)

    set_speaker_notes(s, (
        "Set up the pain. Client requirements arrive through four "
        "different surfaces: Teams calls, Teams chats, Gmail threads, "
        "and the shared Google Drive folder. Each developer on our side "
        "talks to the client through some subset of those, but never "
        "all of them. Before Nucleus there was no place where all "
        "conversations with a client landed in the same shape, on the "
        "same day, attributable back to source. Half a feature would "
        "be mentioned on a call with one developer, the rest typed "
        "into an email reply to another -- and unless someone happened "
        "to be in both conversations, the request fell into the gap. "
        "This deck is the story of how we closed that gap."
    ))


def slide_what_we_built(prs):
    s = base_slide(prs, "What we built",
                   "Four inputs. One AI pass. One verified email out.")

    # Four input chips on the left
    inputs = [
        ("Calls", PURPLE),
        ("Chats", TEAL),
        ("Emails", CORAL),
        ("Files", GOLD),
    ]
    in_x = Inches(0.7)
    in_w = Inches(2.4)
    in_h = Inches(0.75)
    in_y0 = 1.9
    in_gap = 0.35
    for i, (name, color) in enumerate(inputs):
        y = Inches(in_y0 + i * (0.75 + in_gap))
        add_chip(s, in_x, y, in_w, in_h, name, color, font_size=20)

    # Central AI box
    ai_x = Inches(5.4)
    ai_y = Inches(2.6)
    ai_w = Inches(2.6)
    ai_h = Inches(1.8)
    add_rect(s, ai_x, ai_y, ai_w, ai_h, fill_color=NAVY, line_color=NAVY)
    tf_ai = add_textbox(s, ai_x, ai_y, ai_w, ai_h, anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf_ai, "AI", size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_para(tf_ai, "reads and reconciles", size=14,
             color=RGBColor(0xC8, 0xD2, 0xE0), align=PP_ALIGN.CENTER,
             space_before=4)

    # Arrows from each input into AI box
    for i in range(4):
        y = in_y0 + i * (0.75 + in_gap) + 0.375
        add_arrow(s, in_x + in_w, Inches(y),
                  ai_x, ai_y + Inches(0.9),
                  color=MUTED, width_pt=2)

    # Output chip on the right
    out_x = Inches(10.0)
    out_y = Inches(2.9)
    out_w = Inches(2.6)
    out_h = Inches(1.3)
    add_chip(s, out_x, out_y, out_w, out_h,
             "Verification\nemail", GREEN, font_size=22)
    add_arrow(s, ai_x + ai_w, ai_y + Inches(0.9),
              out_x, out_y + Inches(0.65),
              color=NAVY, width_pt=3)

    tf = add_textbox(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.0))
    set_text(tf,
             "Everything a client said to anyone, anywhere, in one day --",
             size=16, color=INK, align=PP_ALIGN.CENTER)
    add_para(tf,
             "summarised back to them for sign-off.",
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             space_before=4)

    set_speaker_notes(s, (
        "One-sentence answer to 'what does this thing do?'. Four capture "
        "surfaces feed one daily AI pass, which produces one verification "
        "email per client. The AI in question is Claude (Anthropic), "
        "running via the Claude Agent SDK on the central Linux host. "
        "Capture happens continuously through the day; the AI pass fires "
        "once at 23:45 Bangladesh time and writes a Gmail draft for "
        "human review. After this slide, the viewer should be able to "
        "explain Nucleus in one sentence."
    ))


def slide_four_channels(prs):
    s = base_slide(prs, "Four ways we capture",
                   "Whatever the client touches, we pick up.")

    channels = [
        ("Calls", PURPLE,
         "Recorded automatically\nwhen a Teams call starts."),
        ("Chats", TEAL,
         "Collected from each\nteammate's Teams chats."),
        ("Emails", CORAL,
         "Pulled from the team's\nGmail throughout the day."),
        ("Shared files", GOLD,
         "Watched in the shared\nGoogle Drive folder."),
    ]
    cw = Inches(2.9)
    ch = Inches(3.6)
    gap = Inches(0.2)
    total = 4 * 2.9 + 3 * 0.2
    x0 = (13.333 - total) / 2
    y = Inches(2.0)
    for i, (title, color, body) in enumerate(channels):
        x = Inches(x0 + i * (2.9 + 0.2))
        add_box(s, x, y, cw, ch, title, body.split("\n"),
                accent=color, title_size=22, body_size=14)

    tf = add_textbox(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.0))
    set_text(tf, "Capture is silent. No one has to remember to log anything.",
             size=16, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

    set_speaker_notes(s, (
        "Four capture surfaces, each running in the background. Teams "
        "calls are picked up by a small voice daemon on each developer "
        "PC -- it watches the audio session via pycaw, and the instant "
        "Teams becomes active it records the mic and the speaker as "
        "two separate WAV tracks. Teams chats are collected by a "
        "scheduled task on each developer PC in three windows per day "
        "(morning, transition, evening). Email is pulled by a container "
        "on the central host via Gmail IMAP every 15 minutes -- threads "
        "are converted to structured .docx with attachments fanned out. "
        "Drive is watched by another container on the central host, "
        "also every 15 minutes; voice notes uploaded to the shared "
        "folder get re-transcribed through the same pipeline as calls."
    ))


def slide_call_to_text(prs):
    s = base_slide(prs, "How calls become text",
                   "Fast cloud transcription. Local backup if cloud fails.")

    # Step strip
    steps = [
        ("1", "Call starts",
         "Recording begins\nautomatically.", PURPLE),
        ("2", "Call ends",
         "Recording stops.\nFiles uploaded.", TEAL),
        ("3", "Transcribed",
         "Text ready within\na few minutes.", GREEN),
    ]
    sw = Inches(3.6)
    sh = Inches(3.4)
    gap = Inches(0.3)
    total = 3 * 3.6 + 2 * 0.3
    x0 = (13.333 - total) / 2
    y = Inches(2.0)
    for i, (num, title, body, color) in enumerate(steps):
        x = Inches(x0 + i * (3.6 + 0.3))
        container = add_rect(s, x, y, sw, sh,
                             fill_color=WHITE, line_color=RULE)
        add_rect(s, x, y, sw, Inches(0.5), fill_color=color)
        tf_n = add_textbox(s, x + Inches(0.2), y + Inches(0.05),
                           Inches(0.6), Inches(0.4),
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_n, num, size=18, bold=True, color=WHITE)
        tf_t = add_textbox(s, x + Inches(0.2), y + Inches(0.7),
                           sw - Inches(0.4), Inches(0.7))
        set_text(tf_t, title, size=22, bold=True, color=color)
        tf_b = add_textbox(s, x + Inches(0.2), y + Inches(1.7),
                           sw - Inches(0.4), Inches(1.6))
        set_text(tf_b, body.split("\n")[0], size=15, color=INK)
        for ln in body.split("\n")[1:]:
            add_para(tf_b, ln, size=15, color=INK, space_before=4)
        if i < 2:
            x_arr = Inches(x0 + i * (3.6 + 0.3) + 3.6)
            add_arrow(s, x_arr, y + Inches(1.7),
                      x_arr + Inches(0.3), y + Inches(1.7),
                      color=NAVY, width_pt=3)

    tf = add_textbox(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(1.2))
    set_text(tf, "Cloud-first for speed. Local fallback when cloud is unavailable.",
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_para(tf,
             "Bangla speech is translated to English in the same step -- "
             "no separate translation pass.",
             size=13, color=MUTED, align=PP_ALIGN.CENTER, space_before=6)

    set_speaker_notes(s, (
        "Transcription pipeline in plain terms. Primary path is the Groq "
        "API hitting whisper-large-v3 on their LPU hardware; we use the "
        "translation endpoint so Bangla audio comes back as English text "
        "in one round trip. A typical 3-minute call comes back in about "
        "30 seconds. Groq's free tier gives us 8 hours of audio per "
        "day, which comfortably covers the team's actual call volume. "
        "Fallback path is faster-whisper int8 quantized on CPU on the "
        "Linux host -- same model lineage, slower (~10 min for the "
        "same call) but no external dependency. Fallback fires "
        "automatically if the API key is missing, the network is down, "
        "we hit the rate limit, or the file is over the 25 MB upload "
        "cap. On normal days the fallback model is never loaded, "
        "saving ~3 GB of memory."
    ))


def slide_central_server(prs):
    s = base_slide(prs, "What our central server does",
                   "One Linux box. Six quiet background jobs.")

    roles = [
        ("Stores recordings", TEAL,
         "Holds every call, chat,\nemail, and shared file."),
        ("Listens for new content", CORAL,
         "Notices the moment\nsomething new arrives."),
        ("Turns audio into text", PURPLE,
         "Sends recordings out\nfor fast transcription."),
        ("Pulls email and files", GOLD,
         "Checks Gmail and Drive\nevery 15 minutes."),
        ("Runs the AI", NAVY,
         "Each night, reviews the\nwhole day's material."),
        ("Drafts the email", GREEN,
         "Writes the verification\nemail for your review."),
    ]
    cw = Inches(4.05)
    ch = Inches(2.15)
    gap = Inches(0.15)
    x0 = Inches(0.5)
    y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(roles):
        row, col = divmod(i, 3)
        x = x0 + col * (cw + gap)
        y = y0 + Inches(row * 2.3)
        add_box(s, x, y, cw, ch, title, body.split("\n"),
                accent=color, title_size=17, body_size=13)

    tf = add_textbox(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.6))
    set_text(tf,
             "All six jobs run continuously. No one logs in to babysit them.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)

    set_speaker_notes(s, (
        "The central server is the Ubuntu 24.04 box at 172.16.205.123. "
        "Implementation detail: it's a single docker-compose stack of "
        "six containers -- nucleus-samba (serves the shared folder over "
        "SMB on port 445), nucleus-transcribe (the call-to-text loop), "
        "nucleus-stage-email (Gmail IMAP every 15 min), "
        "nucleus-stage-drive (Google Drive watcher every 15 min), "
        "nucleus-daily-draft (the 23:45 BD agent run), and "
        "nucleus-gha-runner (self-hosted GitHub Actions runner). The "
        "repo is bind-mounted read-only so `git pull` on the host hot-"
        "reloads every worker without an image rebuild. State that "
        "must survive `compose down` lives in a named docker volume "
        "(the sqlite memory database) and on the host filesystem (the "
        "shared central tree). This box took over from the old Windows "
        "MVPACCESS agent host on 2026-05-14."
    ))


def slide_what_ai_does(prs):
    s = base_slide(prs, "What the AI does",
                   "Reads everything from the day. Drafts one email per client.")

    # Big sequential cards
    cards = [
        ("Reads", NAVY,
         "Every call transcript, chat,\nemail, and file from the day."),
        ("Identifies", PURPLE,
         "Picks out real client requirements.\nIgnores noise and chatter."),
        ("Drafts", GOLD,
         "Writes a verification email,\nattaches the source material."),
    ]
    cw = Inches(3.9)
    ch = Inches(3.4)
    gap = Inches(0.25)
    total = 3 * 3.9 + 2 * 0.25
    x0 = (13.333 - total) / 2
    y = Inches(1.9)
    for i, (title, color, body) in enumerate(cards):
        x = Inches(x0 + i * (3.9 + 0.25))
        add_box(s, x, y, cw, ch, title, body.split("\n"),
                accent=color, title_size=24, body_size=15)
        if i < 2:
            x_arr = Inches(x0 + i * (3.9 + 0.25) + 3.9)
            add_arrow(s, x_arr, y + Inches(1.7),
                      x_arr + Inches(0.25), y + Inches(1.7),
                      color=NAVY, width_pt=3)

    tf = add_textbox(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.4))
    set_text(tf, "One pass per day. One email per client.",
             size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_para(tf,
             "Already-confirmed items from earlier days are skipped, "
             "so the client only ever sees what's new.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER, space_before=6)

    set_speaker_notes(s, (
        "The AI is Anthropic's Claude, running through the Claude Agent "
        "SDK in the nucleus-daily-draft container on the central host. "
        "It fires once per day at 23:45 BD. The input is a 'Pull Session' "
        "document -- about 16,000 characters for a typical day -- that "
        "stitches every call transcript, chat window, email thread, and "
        "Drive file together, grouped by client and ordered in time. "
        "Pre-filters strip system notifications, recorder test snippets, "
        "and peer-to-peer dev chatter before Claude sees it. The output "
        "is two artefacts per client: a curated Requirements "
        "Verification .docx with one entry per requirement (each with "
        "a source pointer back to the call timestamp / chat / email) "
        "and a .eml draft in Gmail Drafts ready for human review. "
        "Marginal cost per run is $0 -- the Claude pass runs on the "
        "team's existing Max-tier subscription rather than per-token "
        "API billing."
    ))


def slide_human_in_loop(prs):
    s = base_slide(prs, "Human in the loop",
                   "You always review before it sends. By design.")

    # Big centered emphasis with three short cards below
    tf = add_textbox(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.4))
    set_text(tf, "Nothing goes to the client without your click.",
             size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    cards = [
        ("Review", NAVY,
         "Open Gmail Drafts.\nRead the list of\nidentified items."),
        ("Edit if needed", GREEN,
         "Fix anything that looks\nwrong. Add what's missing.\nThen click send."),
        ("Reply loops back", PURPLE,
         "When the client replies,\nthe system learns.\nConfirmed items don't\nresurface tomorrow."),
    ]
    cw = Inches(3.9)
    ch = Inches(2.9)
    gap = Inches(0.25)
    total = 3 * 3.9 + 2 * 0.25
    x0 = (13.333 - total) / 2
    y = Inches(3.6)
    for i, (title, color, body) in enumerate(cards):
        x = Inches(x0 + i * (3.9 + 0.25))
        add_box(s, x, y, cw, ch, title, body.split("\n"),
                accent=color, title_size=20, body_size=14)
        if i < 2:
            x_arr = Inches(x0 + i * (3.9 + 0.25) + 3.9)
            add_arrow(s, x_arr, y + Inches(1.4),
                      x_arr + Inches(0.25), y + Inches(1.4),
                      color=NAVY, width_pt=2.5)

    tf2 = add_textbox(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.5))
    set_text(tf2, "The system drafts. The human ships.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)

    set_speaker_notes(s, (
        "This is the only manual step in the whole pipeline, and the "
        "deliberate design choice is that the system never sends mail to "
        "a client without a human pressing send. Open Gmail Drafts -- "
        "the .eml is already there with the curated .docx and the raw "
        "Pull Session source doc attached. Skim, cross-check anything "
        "weird against the source, edit if needed, send. When the client "
        "replies, tools.poll_replies on the central host matches the "
        "reply by thread id, parses line-by-line accept/reject/clarify, "
        "and writes the deltas into the sqlite memory database. "
        "Tomorrow's daily run reads that memory before drafting, so "
        "confirmed items do not get re-asked. This slide matters for "
        "trust-building -- emphasize that the human gate is the policy, "
        "not a temporary precaution."
    ))


def slide_day_in_life(prs):
    s = base_slide(prs, "A day in the life",
                   "Same client mentioned across channels. One draft at the end.")

    # Vertical timeline
    events = [
        ("9 AM",  PURPLE, "Call",    "Client asks for a feature on a call."),
        ("11 AM", CORAL,  "Email",   "Same client follows up with more detail by email."),
        ("3 PM",  GOLD,   "File",    "Client uploads a reference document to the shared folder."),
        ("6 PM",  TEAL,   "Chat",    "Another teammate gets a clarifying chat message."),
        ("11:45 PM", GREEN, "Draft", "The AI produces one verification email covering all four."),
    ]
    line_x = Inches(2.5)
    add_rect(s, line_x, Inches(1.7), Inches(0.04), Inches(5.0),
             fill_color=RULE)
    y0 = 1.6
    step = 1.0
    for i, (time_label, color, kind, body) in enumerate(events):
        y = Inches(y0 + i * step)
        # Time on the left
        tf_t = add_textbox(s, Inches(0.7), y, Inches(1.7), Inches(0.5))
        set_text(tf_t, time_label, size=18, bold=True, color=color,
                 align=PP_ALIGN.RIGHT)
        # Dot on the line
        add_chip(s, line_x - Inches(0.1), y + Inches(0.08),
                 Inches(0.25), Inches(0.25), "", color)
        # Kind + body on the right
        kind_chip_w = Inches(1.3)
        add_chip(s, Inches(2.8), y + Inches(0.02),
                 kind_chip_w, Inches(0.45), kind, color, font_size=13)
        tf_b = add_textbox(s, Inches(2.8) + kind_chip_w + Inches(0.2),
                           y + Inches(0.05), Inches(7.5), Inches(0.5))
        set_text(tf_b, body, size=15, color=INK)

    tf = add_textbox(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.5))
    set_text(tf, "No one saw all four. The system did.",
             size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    set_speaker_notes(s, (
        "Concrete illustrative scenario. The point is that no single "
        "person was in the room (or thread, or call) for all four "
        "touchpoints. A client requirement that spans channels is the "
        "common case, not the edge case, and that's exactly the kind "
        "of thing humans lose track of. The AI pass at 23:45 BD doesn't "
        "care which channel said what -- it stitches the day together "
        "and surfaces the full picture in one draft. The times in this "
        "slide are illustrative; the only real time on the slide is "
        "23:45 BD, which is the actual daily-draft fire time."
    ))


def slide_what_it_costs(prs):
    s = base_slide(prs, "What it costs",
                   "No new line items. No per-run charge.")

    # Giant $0 dominates the slide
    tf_big = add_textbox(s, Inches(0.7), Inches(1.5), Inches(11.9),
                         Inches(2.6), anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf_big, "$0", size=200, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)

    # Subtitle / explainer directly under the $0
    tf_sub = add_textbox(s, Inches(0.7), Inches(4.2), Inches(11.9),
                         Inches(0.7))
    set_text(tf_sub,
             "No per-run charge. No per-developer license. "
             "No new infrastructure cost.",
             size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Quiet supporting strip: 2x2 grid of "where each cost would go"
    rows = [
        ("AI / requirement drafting", "Existing Claude subscription"),
        ("Audio transcription",       "Free tier (Groq)"),
        ("Central server",            "Shared with OpenProject (no new cost)"),
        ("GitHub runner + storage",   "Self-hosted (no charge)"),
    ]
    cell_w = Inches(5.9)
    cell_h = Inches(0.7)
    gap_x = Inches(0.2)
    gap_y = Inches(0.1)
    total_w = 2 * 5.9 + 0.2
    x0 = (13.333 - total_w) / 2
    y0 = 5.3
    for i, (label, value) in enumerate(rows):
        row, col = divmod(i, 2)
        x = Inches(x0 + col * (5.9 + 0.2))
        y = Inches(y0 + row * (0.7 + 0.1))
        add_rect(s, x, y, cell_w, cell_h,
                 fill_color=SOFT, line_color=RULE)
        # Label on the left half, value on the right half
        tf_l = add_textbox(s, x + Inches(0.2), y, Inches(2.6), cell_h,
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_l, label, size=12, color=MUTED, align=PP_ALIGN.LEFT)
        tf_v = add_textbox(s, x + Inches(2.85), y,
                           cell_w - Inches(3.0), cell_h,
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_v, value, size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.LEFT)

    set_speaker_notes(s, (
        "Cost story for the budget conversation. The incremental cost of "
        "running Nucleus per daily pass is effectively zero -- there is "
        "no new line item on anyone's bill. The AI side runs against an "
        "Anthropic Claude Max-tier subscription, which is a flat monthly "
        "subscription we already have in place; it's not metered per "
        "token, so the daily verify_session pass doesn't show up as a "
        "per-run charge. Audio transcription goes through Groq's free "
        "tier -- 8 hours of audio per day, which is well above our "
        "actual call volume. The central server is 172.16.205.123, the "
        "Linux box that already hosts our OpenProject instance; Nucleus "
        "runs as a docker-compose stack alongside it on the same host, "
        "so there is no new VM, no new cloud subscription, no new "
        "hardware. The GitHub Actions runner that builds and deploys "
        "the stack is self-hosted on the same .123 box, so we don't "
        "pay GitHub for runner minutes or artifact storage. Heads up "
        "on an older number that may surface: a prior version of this "
        "deck quoted '~6 cents per run' as the Claude cost. That "
        "figure came from a snapshot where Nucleus was billed against "
        "the regular Claude API with per-token pricing. Since we moved "
        "to Max-tier subscription auth, that per-run charge no longer "
        "applies -- the subscription is flat and the marginal cost of "
        "one more daily run is $0."
    ))


def slide_where_we_are(prs):
    s = base_slide(prs, "Where we are today",
                   "Live since 2026-05-14.")

    # Top status row
    statuses = [
        ("Live", "since 2026-05-14", GREEN),
        ("6", "background jobs running", NAVY),
        ("4", "capture channels active", TEAL),
        ("1", "verification email last night", GOLD),
    ]
    box_w = Inches(2.95)
    gap = Inches(0.2)
    total = 4 * 2.95 + 3 * 0.2
    x0 = (13.333 - total) / 2
    y = Inches(1.7)
    for i, (big, label, color) in enumerate(statuses):
        x = Inches(x0 + i * (2.95 + 0.2))
        container = add_rect(s, x, y, box_w, Inches(2.0),
                             fill_color=WHITE, line_color=RULE)
        add_rect(s, x, y, box_w, Inches(0.14), fill_color=color)
        tf_big = add_textbox(s, x, y + Inches(0.4),
                             box_w, Inches(1.0))
        set_text(tf_big, big, size=38, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        tf_lbl = add_textbox(s, x, y + Inches(1.4),
                             box_w, Inches(0.5))
        set_text(tf_lbl, label, size=13, color=MUTED,
                 align=PP_ALIGN.CENTER)

    # Healthy / ready strip
    add_box(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.4),
            "Yesterday, end to end", [
                "The day's calls, chats, emails, and files were captured.",
                "The AI pass ran on schedule.",
                "A draft verification email landed in Gmail Drafts.",
                "Memory was updated so today won't re-ask anything.",
            ], accent=GREEN, title_size=18, body_size=14)

    set_speaker_notes(s, (
        "Status snapshot as of 2026-05-14, the cutover day. Everything "
        "that was supposed to be running on the Linux central host IS "
        "running. The six containers are nucleus-samba, "
        "nucleus-transcribe, nucleus-stage-email, nucleus-stage-drive, "
        "nucleus-daily-draft, and nucleus-gha-runner. The four "
        "channels are calls, chats, email, and Drive. The 23:45 BD "
        "fire produced a verification email last night for the one "
        "active client. Run `python -m tools.healthcheck` on the host "
        "for the live version of this snapshot at any time."
    ))


def slide_whats_next(prs):
    s = base_slide(prs, "What's next",
                   "Three priorities for the coming weeks.")

    nexts = [
        ("Bring on more clients", CORAL,
         "The pipeline scales\nper client at no\nadditional cost."),
        ("Onboard remaining\ndevelopers", TEAL,
         "Add each teammate's PC\nto the capture network.\nAround 10 minutes each."),
        ("Tighten the loop", GOLD,
         "Track edits the human\nmakes to drafts. Reduce\nthem over time."),
    ]
    cw = Inches(3.9)
    ch = Inches(3.4)
    gap = Inches(0.25)
    total = 3 * 3.9 + 2 * 0.25
    x0 = (13.333 - total) / 2
    y = Inches(2.0)
    for i, (title, color, body) in enumerate(nexts):
        x = Inches(x0 + i * (3.9 + 0.25))
        add_box(s, x, y, cw, ch, title, body.split("\n"),
                accent=color, title_size=20, body_size=15)

    tf = add_textbox(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.2))
    set_text(tf, "The foundation is built. Now we extend it.",
             size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_para(tf,
             "Each next step builds on the same daily pass -- no "
             "re-architecture needed.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER, space_before=6)

    set_speaker_notes(s, (
        "Three near-term priorities, ordered by what unblocks what. "
        "First, extend the client roster -- today's roster is one "
        "client; the pipeline is already per-client (collect_central "
        "takes --client as an argument), so adding another is "
        "configuration, not code (Drive folder + email alias + "
        "metadata.client_name resolution on calls). Second, onboard "
        "the remaining developers onto the voice daemon -- 10 minutes "
        "per teammate using setup.bat plus the SMB mount plus "
        "start-daemon.bat. Third, track the human-edit deltas on the "
        "AI's drafts; once accuracy is high enough on routine "
        "confirmations we can flip a feature flag and let the system "
        "auto-send those, while keeping the human gate on new asks. "
        "Beyond these: OpenProject auto-publishing of confirmed "
        "requirements, calendar-aware capture, per-developer Samba "
        "accounts, Prometheus monitoring, and backups for the "
        "central tree and sqlite memory volume."
    ))


def slide_questions(prs):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=NAVY)

    tf = add_textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.0),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, "Questions?", size=88, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    tf2 = add_textbox(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.8),
                      anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf2, "Thank you.",
             size=24, color=RGBColor(0xC8, 0xD2, 0xE0),
             align=PP_ALIGN.CENTER)

    tf3 = add_textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
    set_text(tf3, "Adaptive Enterprise / NAPCO labs",
             size=12, color=RGBColor(0x9C, 0xAB, 0xBE),
             align=PP_ALIGN.CENTER)

    set_speaker_notes(s, (
        "Open the floor. If there's time, offer a live walkthrough -- "
        "show the verification email sitting in Gmail Drafts, or run "
        "`python -m tools.healthcheck` on the central host to display "
        "the live status of all six containers."
    ))


# -- build ----------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_what_we_built(prs)
    slide_four_channels(prs)
    slide_call_to_text(prs)
    slide_central_server(prs)
    slide_what_ai_does(prs)
    slide_human_in_loop(prs)
    slide_day_in_life(prs)
    slide_what_it_costs(prs)
    slide_where_we_are(prs)
    slide_whats_next(prs)
    slide_questions(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote: {OUT}")
    print(f"  slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
