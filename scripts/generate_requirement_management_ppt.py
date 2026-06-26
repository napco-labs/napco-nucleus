"""NAPCO Nucleus — Requirement Management deck.

18 slides, 16:9. Tells the audience how client requirements get
captured, identified, and turned into a verification email. Minimal
text per slide so Titu drives the narration; speaker notes carry
the talking points.

Audience-targeted blocks:
- Slides 2-7 (problem, solution, end-to-end journey, tools, security,
  costs) — boss / non-tech viewers. Self-contained — can be presented
  standalone.
- Slides 12-13 (architecture, where AI helps) — engineering team.
  Pairs the where-the-work-happens diagram with the AI-contribution
  breakdown so devs see both the topology and the labor split.

Run:
    python scripts\\generate_requirement_management_ppt.py
Output:
    docs\\NAPCO-Nucleus-Requirement-Management.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT = ROOT / "docs" / "NAPCO-Nucleus-Requirement-Management.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette — matches the Central-Architecture deck
NAVY   = RGBColor(0x1F, 0x4E, 0x79)
TEAL   = RGBColor(0x2E, 0x8A, 0x8A)
CORAL  = RGBColor(0xE0, 0x78, 0x56)
GREEN  = RGBColor(0x4A, 0x7A, 0x4A)
GOLD   = RGBColor(0xC9, 0x96, 0x2B)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)
INK    = RGBColor(0x22, 0x22, 0x22)
MUTED  = RGBColor(0x6B, 0x77, 0x85)
SOFT   = RGBColor(0xF5, 0xF7, 0xFA)
RULE   = RGBColor(0xD5, 0xDC, 0xE5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# ── helpers ────────────────────────────────────────────────────────

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def no_line(shape):
    shape.line.fill.background()


def line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_text(tf, text, size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_para(tf, text, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, space_before=0, font="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    return tf


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None,
             line_width=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color is None:
        no_fill(rect)
    else:
        solid(rect, fill_color)
    if line_color is None:
        no_line(rect)
    else:
        line(rect, line_color, line_width)
    rect.text_frame.margin_left = Inches(0.1)
    rect.text_frame.margin_right = Inches(0.1)
    rect.text_frame.margin_top = Inches(0.05)
    rect.text_frame.margin_bottom = Inches(0.05)
    return rect


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width_pt=2.0):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn


def base_slide(prs, title_text, subtitle_text=None):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill_color=NAVY)
    tf = add_textbox(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    set_text(tf, title_text, size=28, bold=True, color=NAVY)
    if subtitle_text:
        add_para(tf, subtitle_text, size=14, color=MUTED, space_before=2)
    add_rect(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.02),
             fill_color=RULE)
    return s


def set_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    p.text = text


def add_chip(slide, x, y, w, h, text, color):
    rect = add_rect(slide, x, y, w, h, fill_color=color, rounded=True)
    tf = rect.text_frame
    set_text(tf, text, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return rect


def add_box(slide, x, y, w, h, title, lines, accent=NAVY):
    container = add_rect(slide, x, y, w, h, fill_color=WHITE, line_color=RULE)
    add_rect(slide, x, y, Inches(0.08), h, fill_color=accent)
    tf = add_textbox(slide, x + Inches(0.18), y + Inches(0.12),
                     w - Inches(0.28), h - Inches(0.24))
    set_text(tf, title, size=14, bold=True, color=accent)
    for ln in lines:
        add_para(tf, ln, size=11, color=INK, space_before=2)
    return container


# ── slides ─────────────────────────────────────────────────────────

def slide_title(prs):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=NAVY)
    tf = add_textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.5))
    set_text(tf, "NAPCO Nucleus", size=58, bold=True, color=WHITE)
    add_para(tf, "Requirement Management Workflow",
             size=28, color=WHITE)
    tf2 = add_textbox(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.5))
    set_text(tf2, "Client conversations -> one verification email, "
                  "end-to-end.",
             size=18, color=RGBColor(0xC8, 0xD2, 0xE0))
    tf3 = add_textbox(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    set_text(tf3, "Mohammad Kamrul Hasan (Titu)   |   "
                  "Adaptive Enterprise / NAPCO labs",
             size=12, color=RGBColor(0x9C, 0xAB, 0xBE))
    set_speaker_notes(s, (
        "Open with the framing: requirements come at us from multiple "
        "channels and multiple people. Today's deck shows how NAPCO "
        "Nucleus stitches all of that back together into one verification "
        "email that I review and send to the client myself. The team's "
        "involvement is small but important -- I'll explain exactly what "
        "I need from each of you."
    ))


def slide_problem(prs):
    s = base_slide(prs, "The problem",
                   "A single client requirement is rarely in one place.")
    add_chip(s, Inches(0.8), Inches(1.6), Inches(2.4), Inches(0.6),
             "Client", CORAL)
    add_chip(s, Inches(0.8), Inches(2.55), Inches(2.4), Inches(0.6),
             "Dev 1", NAVY)
    add_chip(s, Inches(0.8), Inches(3.30), Inches(2.4), Inches(0.6),
             "Dev 2", NAVY)
    add_chip(s, Inches(0.8), Inches(4.05), Inches(2.4), Inches(0.6),
             "Dev 3", NAVY)

    channels = [
        ("Teams chat",  CORAL,  "\"Make the login flow log failed attempts\""),
        ("Teams call",  TEAL,   "\"...and also keep audit timestamps in UTC\""),
        ("Email",       GOLD,   "PDF attached: \"Requirement spec v2\""),
        ("Drive",       PURPLE, "Spreadsheet of edge cases dropped in folder"),
    ]
    for i, (name, color, quote) in enumerate(channels):
        y = Inches(1.6 + i * 1.20)
        add_chip(s, Inches(4.2), y, Inches(2.0), Inches(0.6), name, color)
        tf = add_textbox(s, Inches(6.5), y, Inches(6.4), Inches(0.6),
                         anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, quote, size=12, color=INK)

    add_rect(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
             fill_color=SOFT, line_color=RULE)
    tf = add_textbox(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.4),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, "Nobody has all the pieces. Stitching them by hand is "
                 "slow and error-prone.",
             size=12, color=MUTED)

    set_speaker_notes(s, (
        "Today, a single requirement might mention something in a Teams "
        "chat with Dev 2, get clarified in a call with Dev 1, and arrive "
        "fully formed as a PDF in my inbox. By the time I notice all "
        "three exist I've already lost a day. The cost isn't just time "
        "-- it's that pieces get dropped silently."
    ))


def slide_solution(prs):
    s = base_slide(prs, "The solution",
                   "One pipeline, four channels in, one verification email out.")
    boxes = [
        ("CAPTURE",   "Every dev machine + agent host pulls "
                      "chat, calls, email, and Drive into one document.",
                      NAVY),
        ("IDENTIFY",  "Claude reads the document and extracts a "
                      "deduped, structured requirement list.",
                      TEAL),
        ("VERIFY",    "Titu reviews the list, edits if needed, "
                      "and clicks Send when satisfied.",
                      GREEN),
        ("DELIVER",   "Client receives an email asking them to "
                      "confirm. Nothing leaves Gmail without Titu.",
                      CORAL),
    ]
    box_w = Inches(2.85)
    box_h = Inches(3.6)
    gap = Inches(0.25)
    start_x = Inches(0.7)
    y = Inches(1.7)
    for i, (title, body, color) in enumerate(boxes):
        x = start_x + (box_w + gap) * i
        add_box(s, x, y, box_w, box_h, title, [body], accent=color)
        if i < len(boxes) - 1:
            ax1 = x + box_w + Inches(0.02)
            ax2 = ax1 + gap - Inches(0.04)
            add_arrow(s, ax1, y + box_h / 2, ax2, y + box_h / 2,
                      color=MUTED, width_pt=2.2)

    tf = add_textbox(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.8),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, "Human-in-the-loop at the Verify step. Nothing is auto-sent.",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    set_speaker_notes(s, (
        "Four stages. Capture is the only one running in the background. "
        "Identify is one LLM call, on my command. Verify is me reading "
        "the document. Deliver is me clicking Send. The team is in "
        "Capture only -- the rest is on me."
    ))


def slide_channels(prs):
    s = base_slide(prs, "The four input channels",
                   "All four stage to central automatically. No operator action.")
    channels = [
        ("Teams chat",
         "Reads each dev's local Teams. No API, no token, no extra login.",
         "Pushed throughout the day; more often in the evening when clients are online.",
         CORAL),
        ("Teams calls",
         "Records mic + speaker as separate tracks while you talk.",
         "Starts automatically when you join a Teams call. Around the clock.",
         TEAL),
        ("Email",
         "Reads my Gmail inbox on the central server.",
         "Picked up every 15 min, day and night. Never reads the same mail twice.",
         GOLD),
        ("Google Drive",
         "Reads one shared folder we set up for the project.",
         "Picked up every 15 min, day and night.",
         PURPLE),
    ]
    w = Inches(5.8)
    h = Inches(2.55)
    positions = [
        (Inches(0.6), Inches(1.5)),
        (Inches(6.95), Inches(1.5)),
        (Inches(0.6), Inches(4.25)),
        (Inches(6.95), Inches(4.25)),
    ]
    for (title, l1, l2, color), (x, y) in zip(channels, positions):
        add_box(s, x, y, w, h, title, [l1, "", l2], accent=color)

    set_speaker_notes(s, (
        "All four channels stage to central in the background -- nothing "
        "to remember. Chat push runs on three Dhaka-time schedules: every "
        "two hours through the workday, once at 17:30 to bridge the gap, "
        "then every thirty minutes in the evening when most client "
        "messages come in. Teams calls record themselves the moment Teams "
        "opens an audio session -- no phrase to say, no button to press. "
        "Email and Drive run around the clock from the central server. "
        "Dev machines never hold Gmail or Drive credentials -- those "
        "live on the central server alone. The Requirement Management "
        "workflow reads from the central store when I'm ready to draft."
    ))


def slide_dev_setup(prs):
    s = base_slide(prs, "What every developer does -- once",
                   "Five steps. Ten minutes. No secrets.")
    steps = [
        ("1", "Clone the repo",
         "git clone https://github.com/napco-labs/napco-nucleus"),
        ("2", "Run scripts\\setup.bat",
         "Python, venv, deps, .env. No App Password needed."),
        ("3", "Confirm one line in .env",
         "NUCLEUS_CENTRAL_PATH=\\\\172.16.205.123\\nucleus-central"),
        ("4", "Register the chat-push schedule",
         ".\\scripts\\register-chat-push-task.ps1  (admin PowerShell)"),
        ("5", "Install the voice daemon (starts itself at every logon)",
         "Double-click scripts\\install-voice-daemon.bat"),
    ]
    y = Inches(1.55)
    row_h = Inches(0.95)
    for num, head, body in steps:
        add_chip(s, Inches(0.7), y, Inches(0.7), Inches(0.7), num, NAVY)
        tf = add_textbox(s, Inches(1.6), y, Inches(11.0), Inches(0.7),
                         anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, head, size=16, bold=True, color=NAVY)
        add_para(tf, body, size=12, color=MUTED, space_before=2)
        y = y + row_h

    add_rect(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.45),
             fill_color=SOFT, line_color=RULE)
    tf = add_textbox(s, Inches(0.8), Inches(6.58), Inches(11.7), Inches(0.4),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, "Full walk-through: docs/Setup_Guide.pdf in the repo.",
             size=12, color=MUTED)

    set_speaker_notes(s, (
        "The Developer Setup PDF in docs/ goes through every step in "
        "detail. The two things to remember: run the chat-push "
        "registration in admin PowerShell so the three scheduled tasks "
        "land, and use install-voice-daemon.bat -- it registers a "
        "Scheduled Task that re-launches the daemon at every logon and "
        "restarts it on crash. Recordings happen automatically the "
        "moment Teams opens an audio session; no buttons, no phrases."
    ))


def slide_daily(prs):
    s = base_slide(prs, "What every developer does -- daily",
                   "Mostly nothing. The system does the work.")

    rows = [
        ("Get your chat / attachments to central",
         "Nothing -- it pushes itself throughout the day"),
        ("Record a Teams call",
         "Nothing -- it records the moment you join a call"),
        ("Include a file from Teams chat",
         "Click \"Download\" on the attachment so it lands in ~/Downloads"),
        ("Update to the latest code",
         "Double-click scripts\\update.bat after a git pull notice"),
    ]
    y = Inches(1.7)
    for what, how in rows:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0),
                 fill_color=WHITE, line_color=RULE)
        tf1 = add_textbox(s, Inches(0.9), y + Inches(0.15),
                          Inches(5.0), Inches(0.7),
                          anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf1, what, size=14, bold=True, color=NAVY)
        tf2 = add_textbox(s, Inches(6.1), y + Inches(0.15),
                          Inches(6.6), Inches(0.7),
                          anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf2, how, size=13, color=INK)
        y = y + Inches(1.15)

    set_speaker_notes(s, (
        "I want everyone to internalise: you don't run any command "
        "during the day. Your machine pushes chat on its three BD-time "
        "schedules; the voice daemon auto-records the second Teams "
        "opens an audio session -- nothing to say, nothing to press; "
        "downloads land in ~/Downloads automatically when you click. "
        "That's it."
    ))


def slide_titu_command(prs):
    s = base_slide(prs, "What I do",
                   "One command, end-to-end, on demand.")
    add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.05),
             fill_color=NAVY)
    tf = add_textbox(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.8),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, "GHA: Run workflow \"Requirement Management\"   "
                 "or   scripts\\requirement-management.bat",
             size=18, bold=True, color=WHITE, font="Consolas")

    steps = [
        ("1", "Push my local Teams chat to central (force-flush)"),
        ("2", "Walk central: chat + calls + already-staged email + Drive"),
        ("3", "Transcribe every new call (Whisper chunked, parallel)"),
        ("4", "Bangla -> English happens inline at identify time"),
        ("5", "Claude extracts requirements -> drafts verification email"),
        ("6", "Draft lands in [Gmail]/Drafts for me to review and send"),
    ]
    y = Inches(3.05)
    for num, body in steps:
        add_chip(s, Inches(0.8), y, Inches(0.55), Inches(0.55),
                 num, TEAL)
        tf = add_textbox(s, Inches(1.55), y, Inches(11.0), Inches(0.55),
                         anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, body, size=14, color=INK)
        y = y + Inches(0.6)

    set_speaker_notes(s, (
        "One workflow. Two triggers, same code path: the GitHub Actions "
        "\"Requirement Management\" workflow on the self-hosted runner, "
        "or scripts\\requirement-management.bat on my own machine. "
        "Email + Drive are already on central (auto-staged every 15 min), "
        "so the workflow just walks what's there, transcribes new calls, "
        "and asks Claude to identify + draft. Output is one verification "
        "email in [Gmail]/Drafts -- I read it, edit it, send it, or skip it."
    ))


def slide_architecture(prs):
    s = base_slide(prs, "Where the work happens",
                   "Dev machines push. The shared server pulls and aggregates. Claude reads.")

    # Three columns
    col_w = Inches(3.85)
    col_h = Inches(4.6)
    y = Inches(1.55)
    xs = [Inches(0.6), Inches(4.75), Inches(8.9)]

    add_box(s, xs[0], y, col_w, col_h, "Dev machines (x N)", [
        "Read each dev's local Teams data",
        "Push chat + attachments throughout the day",
        "Record every Teams call automatically",
        "Write to \\\\172.16.205.123\\nucleus-central",
        "No secrets. No API keys.",
    ], accent=NAVY)

    add_box(s, xs[1], y, col_w, col_h, "Shared central server", [
        "Always on, one Linux box for the team",
        "Pulls email + Drive in the background",
        "Transcribes every new call within minutes",
        "Falls back to a local model if needed",
        "Aggregates everything into one document",
    ], accent=TEAL)

    add_box(s, xs[2], y, col_w, col_h, "Claude identifies + drafts", [
        "Claude reads the daily document",
        "Runs every night at 23:45 (Dhaka time)",
        "Also runs on demand when I trigger it",
        "Writes the Requirements Verification doc",
        "Drops the draft email in my Gmail Drafts",
    ], accent=CORAL)

    add_arrow(s, xs[0] + col_w + Inches(0.02), y + col_h / 2,
              xs[1] - Inches(0.04), y + col_h / 2,
              color=MUTED, width_pt=2.5)
    add_arrow(s, xs[1] + col_w + Inches(0.02), y + col_h / 2,
              xs[2] - Inches(0.04), y + col_h / 2,
              color=MUTED, width_pt=2.5)

    set_speaker_notes(s, (
        "Three responsibilities. Dev machines push -- only push -- to "
        "the central Samba share on .123. The Linux central host runs "
        "a docker-compose stack that pulls email + Drive on its own "
        "cadence and transcribes new calls within ~2 minutes of arrival. "
        "The LLM identify + draft step runs once daily at 23:45 BD "
        "inside a container, and also on demand whenever I trigger it. "
        "Each layer holds the smallest permissions it needs -- dev PCs "
        "have no Gmail or Drive credentials at all."
    ))


def slide_journey(prs):
    s = base_slide(prs, "End-to-end: one requirement's journey",
                   "Client said it. Verification email lands the next morning.")
    steps = [
        ("1", "Client says it",        "Teams call, chat, email, or Drive",     "live",          CORAL),
        ("2", "We capture",            "Voice daemon + chat-push, automatic",   "seconds",       NAVY),
        ("3", "Lands on central",      "Uploaded to the shared server",         "minutes",       NAVY),
        ("4", "Transcribed",           "Whisper turns audio into text",         "~2 min later",  TEAL),
        ("5", "Claude drafts",         "Reads the day, writes the email",       "23:45 (Dhaka)", PURPLE),
        ("6", "I review and send",     "Open the draft, click Send",            "next morning",  GREEN),
    ]
    # Six tiles left-to-right
    start_x = Inches(0.4)
    y_box = Inches(1.65)
    box_w = Inches(2.02)
    box_h = Inches(3.6)
    gap = Inches(0.13)
    chip_d = Inches(0.85)
    chip_y = y_box + Inches(0.25)

    for i, (num, head, body, when, color) in enumerate(steps):
        x = start_x + (box_w + gap) * i
        # Card
        add_rect(s, x, y_box, box_w, box_h,
                 fill_color=WHITE, line_color=RULE)
        add_rect(s, x, y_box, box_w, Inches(0.10), fill_color=color)
        # Big number chip centered horizontally
        chip_x = x + (box_w - chip_d) / 2
        add_chip(s, chip_x, chip_y, chip_d, chip_d, num, color)
        # Step name (bold)
        tf_h = add_textbox(s, x + Inches(0.1), y_box + Inches(1.30),
                           box_w - Inches(0.2), Inches(0.6),
                           anchor=MSO_ANCHOR.TOP)
        set_text(tf_h, head, size=15, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        # Body (one short line)
        tf_b = add_textbox(s, x + Inches(0.1), y_box + Inches(1.90),
                           box_w - Inches(0.2), Inches(1.0),
                           anchor=MSO_ANCHOR.TOP)
        set_text(tf_b, body, size=11, color=MUTED,
                 align=PP_ALIGN.CENTER)
        # Timing pill at the bottom of the card
        pill_w = Inches(1.5)
        pill_h = Inches(0.36)
        pill_x = x + (box_w - pill_w) / 2
        pill_y = y_box + box_h - pill_h - Inches(0.15)
        add_chip(s, pill_x, pill_y, pill_w, pill_h, when, color)

        # Arrow to next box (skip on the last)
        if i < len(steps) - 1:
            ay = y_box + box_h / 2
            ax1 = x + box_w + Inches(0.01)
            ax2 = ax1 + gap - Inches(0.02)
            add_arrow(s, ax1, ay, ax2, ay, color=MUTED, width_pt=1.8)

    # Footer
    tf = add_textbox(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.5),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, "Capture is automatic. Drafting is automatic. "
                 "Sending is mine.",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    set_speaker_notes(s, (
        "One concrete example. Say a client mentions during a Teams "
        "call that they want password reset by SMS. The moment the "
        "call starts, the voice daemon on my laptop records the "
        "conversation -- both my mic and Teams' speaker stream. "
        "Within minutes of the call ending, both audio tracks are on "
        "the shared central server. The transcription service "
        "converts them to text in under two minutes. That night at "
        "23:45 Dhaka time, Claude reads everything captured that day "
        "-- this call, plus any chat / email / Drive files -- and "
        "writes a verification email listing the SMS-password-reset "
        "requirement alongside everything else from the day. Next "
        "morning I open my Gmail Drafts, review the requirements "
        "doc, edit if needed, and send. End to end, the human work "
        "is the call itself, the morning review, and the send. The "
        "system handles the steps in the middle."
    ))


def slide_tools(prs):
    s = base_slide(prs, "Tools we use",
                   "Off-the-shelf, well-known, easy to support.")
    boxes = [
        ("AI and language",
         ["Claude (Anthropic) -- reads the daily document and picks out requirements",
          "Whisper on Groq -- turns call audio into searchable text"],
         CORAL),
        ("Communication",
         ["MS Teams -- where calls and chat happen",
          "Gmail -- the verification email lands in my drafts; client replies come back here"],
         TEAL),
        ("Storage and sharing",
         ["Google Drive -- shared client files",
          "A shared drive on the central server -- everything else (transcripts, drafts, audit trail)"],
         GOLD),
        ("Server and code",
         ["Ubuntu Linux + Docker -- one always-on box hosts the whole pipeline",
          "Python + GitHub -- our code base and source of truth"],
         NAVY),
    ]
    # 2 x 2 grid
    box_w = Inches(6.0)
    box_h = Inches(2.55)
    positions = [
        (Inches(0.6), Inches(1.5)),
        (Inches(6.75), Inches(1.5)),
        (Inches(0.6), Inches(4.25)),
        (Inches(6.75), Inches(4.25)),
    ]
    for (title, bullets, color), (x, y) in zip(boxes, positions):
        add_box(s, x, y, box_w, box_h, title, bullets, accent=color)

    set_speaker_notes(s, (
        "All mainstream choices, nothing exotic. Claude from Anthropic "
        "handles the requirements identification. Groq runs OpenAI's "
        "Whisper model for transcription -- same model the open-source "
        "community uses, just on faster hardware. MS Teams and Gmail "
        "are already AEL's communication backbone, so we read what's "
        "there. Google Drive is the shared-files convention we already "
        "use with clients. The central server is one Ubuntu box "
        "running Docker, which keeps each service isolated and easy "
        "to restart. The code is Python on GitHub -- standard, "
        "auditable, no proprietary lock-in. If we needed to move the "
        "central server tomorrow, the docker-compose file plus one git "
        "clone gets us running on a fresh Linux box."
    ))


def slide_security(prs):
    s = base_slide(prs, "How secure is it?",
                   "What we capture, where it goes, what stays inside.")
    items = [
        ("Recording is MS Teams calls only",
         "Zoom, Google Meet, and any other meeting app are out of scope. "
         "Nothing records outside an active Teams call."),
        ("Dev machines never hold secrets",
         "Gmail and Drive credentials live only on the central server. "
         "Your laptop has no API keys."),
        ("Nothing leaves my mailbox without my click",
         "Every client email is a draft I review and send. "
         "There is no auto-send anywhere in the system."),
        ("External services see only what they must",
         "Call audio briefly goes to a transcription service; the result "
         "comes straight back. Everything else stays on AEL's network."),
    ]
    y = Inches(1.65)
    for what, why in items:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.15),
                 fill_color=WHITE, line_color=RULE)
        add_rect(s, Inches(0.6), y, Inches(0.08), Inches(1.15),
                 fill_color=GREEN)
        tf = add_textbox(s, Inches(0.9), y + Inches(0.15),
                         Inches(11.5), Inches(0.85))
        set_text(tf, what, size=15, bold=True, color=NAVY)
        add_para(tf, why, size=12, color=MUTED, space_before=2)
        y = y + Inches(1.30)

    set_speaker_notes(s, (
        "Worth being explicit about safety because clients and management "
        "will ask. Four guarantees. One: only MS Teams audio sessions "
        "trigger recording -- the process-name check inside the voice "
        "daemon confirms ms-teams.exe / teams.exe / msteams.exe and "
        "nothing else, and the bypass flag is off on every dev PC. Two: "
        "each dev's machine holds no Gmail or Drive credentials -- those "
        "live only on the central server, so there's one host to secure "
        "instead of N. Three: the auto-send boundary stays at me; every "
        "email the client sees, I clicked Send on. Four: audio is sent "
        "to Groq for transcription over encrypted HTTPS and is not "
        "retained for AI training; the requirements identification "
        "step also uses encrypted Claude calls. Everything else -- "
        "chat, email, Drive files, transcripts, the final draft -- "
        "stays on AEL's internal network."
    ))


def slide_costs(prs):
    s = base_slide(prs, "What does it cost to run?",
                   "Zero per run. Everything we use is already paid for or free.")

    # Hero $0
    hero_y = Inches(1.4)
    tf = add_textbox(s, Inches(0.6), hero_y, Inches(12.1), Inches(1.8),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, "$0", size=160, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)
    tf2 = add_textbox(s, Inches(0.6), hero_y + Inches(1.85),
                      Inches(12.1), Inches(0.5),
                      anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf2, "incremental cost per requirement-management run",
             size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # 4 reason tiles
    tiles = [
        ("Claude", "Flat-rate subscription. All our usage included.", CORAL),
        ("Groq", "Free tier. Covers our daily audio volume with room to spare.", TEAL),
        ("Gmail + Drive", "Existing AEL accounts. No new licenses.", GOLD),
        ("Linux server", "Already running for AEL workloads. No new hardware.", NAVY),
    ]
    tile_w = Inches(2.95)
    tile_h = Inches(1.85)
    gap = Inches(0.13)
    start_x = Inches(0.6)
    tile_y = Inches(5.05)
    for i, (name, why, color) in enumerate(tiles):
        x = start_x + (tile_w + gap) * i
        add_rect(s, x, tile_y, tile_w, tile_h,
                 fill_color=WHITE, line_color=RULE)
        add_rect(s, x, tile_y, tile_w, Inches(0.10), fill_color=color)
        tf_name = add_textbox(s, x + Inches(0.15), tile_y + Inches(0.25),
                              tile_w - Inches(0.3), Inches(0.5),
                              anchor=MSO_ANCHOR.TOP)
        set_text(tf_name, name, size=16, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        tf_why = add_textbox(s, x + Inches(0.15), tile_y + Inches(0.85),
                             tile_w - Inches(0.3), Inches(0.95),
                             anchor=MSO_ANCHOR.TOP)
        set_text(tf_why, why, size=11, color=MUTED,
                 align=PP_ALIGN.CENTER)

    set_speaker_notes(s, (
        "Cost question always comes up first, so I lead with the "
        "answer: zero. The model is flat-rate -- Anthropic's Claude "
        "subscription bills monthly regardless of how much I run it, "
        "and one Requirement Management run is a tiny fraction of "
        "that budget. Call transcription rides on Groq's free tier, "
        "which currently covers about eight hours of audio a day; "
        "we're using a small fraction of that. Gmail and Google "
        "Drive are AEL's existing communication accounts -- this "
        "system uses what's already paid for. The central Linux "
        "server was already running before Nucleus existed. So the "
        "marginal cost of adding this pipeline, and the marginal "
        "cost of each daily run, is zero. The only meaningful spend "
        "is the time I spend reviewing the verification email -- "
        "and that's the part that should stay human."
    ))


def slide_ai_helps(prs):
    s = base_slide(prs, "Where AI does the work",
                   "Boring work to the machine. Judgment stays with us.")
    quadrants = [
        ("LISTENING", TEAL,
         "Turns hours of call audio into searchable text in minutes.",
         "Instead of: replaying calls and typing transcripts."),
        ("READING", CORAL,
         "Reads the whole day's content -- chat, calls, email, Drive -- in one pass.",
         "Instead of: scanning four channels in four places."),
        ("SORTING", GOLD,
         "Separates real client requirements from chit-chat; deduplicates across channels.",
         "Instead of: maintaining a spreadsheet by hand."),
        ("DRAFTING", PURPLE,
         "Writes the verification email in clean English, ready to send.",
         "Instead of: composing from scratch every time."),
    ]
    box_w = Inches(6.05)
    box_h = Inches(2.55)
    positions = [
        (Inches(0.6),  Inches(1.5)),
        (Inches(6.75), Inches(1.5)),
        (Inches(0.6),  Inches(4.25)),
        (Inches(6.75), Inches(4.25)),
    ]
    for (verb, color, ai_does, replaces), (x, y) in zip(quadrants, positions):
        add_rect(s, x, y, box_w, box_h,
                 fill_color=WHITE, line_color=RULE)
        # Side accent bar
        add_rect(s, x, y, Inches(0.10), box_h, fill_color=color)
        # Big verb at top
        tf_v = add_textbox(s, x + Inches(0.25), y + Inches(0.15),
                           box_w - Inches(0.4), Inches(0.7),
                           anchor=MSO_ANCHOR.TOP)
        set_text(tf_v, verb, size=32, bold=True, color=color)
        # AI does
        tf_a = add_textbox(s, x + Inches(0.25), y + Inches(0.95),
                           box_w - Inches(0.4), Inches(0.65),
                           anchor=MSO_ANCHOR.TOP)
        set_text(tf_a, ai_does, size=13, color=INK)
        # Replaces
        tf_r = add_textbox(s, x + Inches(0.25), y + Inches(1.65),
                           box_w - Inches(0.4), Inches(0.65),
                           anchor=MSO_ANCHOR.TOP)
        set_text(tf_r, replaces, size=12, color=MUTED)

    set_speaker_notes(s, (
        "Four jobs AI takes off the team's plate. The principle is "
        "simple: tedious-but-rule-based work goes to the model; "
        "judgment stays with us. Listening, reading, sorting, "
        "drafting -- each one is hours of manual work compressed to "
        "minutes of compute. What stays human: deciding which "
        "requirements are right, what makes it into the final email, "
        "and when to send. The team's hours go into building what "
        "clients asked for, instead of chasing what they asked for."
    ))


def slide_outputs(prs):
    s = base_slide(prs, "What you get back",
                   "Three artifacts, every run.")
    out_w = Inches(3.85)
    out_h = Inches(4.6)
    y = Inches(1.55)
    xs = [Inches(0.6), Inches(4.75), Inches(8.9)]

    add_box(s, xs[0], y, out_w, out_h, "Raw session document", [
        "Word .docx",
        "Every email body, attachment text,",
        "chat message, call transcript",
        "Audit trail -- read this if",
        "the LLM output is surprising",
    ], accent=NAVY)

    add_box(s, xs[1], y, out_w, out_h, "Requirements verification doc", [
        "Word .docx",
        "Numbered, deduped requirement list",
        "Ready to attach to the client email",
        "Re-runs dedupe against previous",
        "runs via requirements_seen memory",
    ], accent=TEAL)

    add_box(s, xs[2], y, out_w, out_h, "Draft email in [Gmail]/Drafts", [
        "Standard Gmail draft",
        "Verification doc attached",
        "Pre-filled subject + body",
        "I review -> edit -> send",
        "Nothing is auto-sent",
    ], accent=CORAL)

    set_speaker_notes(s, (
        "Three artifacts. The raw doc is the audit trail. The "
        "requirements doc is what the client will see. The draft is "
        "ready to send. If I find an issue, I can re-run with a "
        "different window or fix the source data and try again."
    ))


def slide_out_of_scope(prs):
    s = base_slide(prs, "What it does NOT do",
                   "Setting boundaries clearly.")
    items = [
        ("Auto-send emails",
         "Every outbound email leaves my mailbox with my explicit click. This is the hard boundary."),
        ("Auto-edit the requirements doc",
         "The draft document is final until I edit it. Claude does not rewrite it on its own."),
        ("Push to OpenProject",
         "OpenProject ingest is manual today; auto only when we hit >=80% monthly accuracy."),
        ("Replace client conversations",
         "This is a capture aid. Real client engagement still happens with humans."),
    ]
    y = Inches(1.65)
    for what, why in items:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.15),
                 fill_color=WHITE, line_color=RULE)
        add_rect(s, Inches(0.6), y, Inches(0.08), Inches(1.15),
                 fill_color=CORAL)
        tf = add_textbox(s, Inches(0.9), y + Inches(0.15),
                         Inches(11.5), Inches(0.85))
        set_text(tf, what, size=15, bold=True, color=NAVY)
        add_para(tf, why, size=12, color=MUTED, space_before=2)
        y = y + Inches(1.30)

    set_speaker_notes(s, (
        "Worth being explicit about boundaries so expectations are right. "
        "Capture is fully automated -- chat, calls, email, and Drive all "
        "stage themselves. Identify + draft also runs on its own once "
        "daily at 23:45 BD. The one thing that stays manual is the send "
        "click: every email that goes to a client leaves my mailbox "
        "because I clicked Send. That's the human-in-the-loop boundary, "
        "and it stays there until the requirements doc is >=80% accurate "
        "month after month. This is a capture aid that I drive at the "
        "final step."
    ))


def slide_roadmap(prs):
    s = base_slide(prs, "Roadmap",
                   "Human-in-the-loop today. Slower-and-correct first.")
    rows = [
        ("Now",
         "Each Requirement Management run produces a draft I review. "
         "Calibration loop, conflict detection, and run telemetry are in.",
         "Manual, predictable, auditable. No incremental cost per run.",
         NAVY),
        ("Next",
         "Tighter dedup across re-runs + Teams attachment auto-resolve "
         "from the in-Teams preview cache",
         "Less manual cleanup before sending.",
         TEAL),
        ("Later", "Auto-publish to OpenProject -- only after >=80% "
                  "monthly accuracy on the requirements doc",
         "Only when we trust the output. Quality gate, not deadline gate.",
         GOLD),
        ("Aspirational", "Auto-send the verification email to clients",
         "Only after the auto-publish step is rock solid.",
         CORAL),
    ]
    y = Inches(1.55)
    for when, what, why, color in rows:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.20),
                 fill_color=WHITE, line_color=RULE)
        add_chip(s, Inches(0.8), y + Inches(0.32), Inches(1.5),
                 Inches(0.55), when, color)
        tf = add_textbox(s, Inches(2.55), y + Inches(0.15),
                         Inches(10.0), Inches(0.95))
        set_text(tf, what, size=14, bold=True, color=NAVY)
        add_para(tf, why, size=11, color=MUTED, space_before=2)
        y = y + Inches(1.32)

    set_speaker_notes(s, (
        "I'm not chasing automation for its own sake. The order is: "
        "stabilise the manual flow, then layer dedup and convenience, "
        "then layer auto-publish behind an accuracy gate, then -- only "
        "if it's rock solid -- consider auto-send. None of this happens "
        "without measurable wins on the earlier step."
    ))


def slide_what_we_need(prs):
    s = base_slide(prs, "What I need from each of you",
                   "Three small things.")
    items = [
        ("1", "Run setup.bat once",
         "Ten minutes. No secrets. The PDF in docs/ walks you through it."),
        ("2", "Leave Teams open + signed in",
         "The chat ingest reads Teams' own local cache."),
        ("3", "Click \"Download\" on chat files that matter",
         "If a teammate shares a PDF in chat, click Download so we can extract it. "
         "Otherwise only the URL gets captured."),
    ]
    y = Inches(1.8)
    for num, head, body in items:
        add_chip(s, Inches(0.7), y, Inches(0.75), Inches(0.75), num, GREEN)
        tf = add_textbox(s, Inches(1.65), y, Inches(11.0), Inches(0.75),
                         anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, head, size=18, bold=True, color=NAVY)
        add_para(tf, body, size=12, color=MUTED, space_before=2)
        y = y + Inches(1.55)

    set_speaker_notes(s, (
        "Three things, kept deliberately small. Setup once. Keep Teams "
        "open. Click Download. That's the whole ask. Everything else -- "
        "the pulling, the transcribing, the identifying, the drafting "
        "-- is on the system and on me."
    ))


def slide_close(prs):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=NAVY)
    tf = add_textbox(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.5))
    set_text(tf, "Questions?", size=64, bold=True, color=WHITE)
    tf2 = add_textbox(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.5))
    set_text(tf2, "Then let's run a live one.",
             size=24, color=RGBColor(0xC8, 0xD2, 0xE0))
    tf3 = add_textbox(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    set_text(tf3, "github.com/napco-labs/napco-nucleus",
             size=12, color=RGBColor(0x9C, 0xAB, 0xBE))
    set_speaker_notes(s, (
        "Wrap on questions, then offer a live run -- do_it_now.py with "
        "the seeded test data so the team can see the verification doc "
        "and the draft email appear."
    ))


# ── main ───────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Boss-facing block (1-7)
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_journey(prs)
    slide_tools(prs)
    slide_security(prs)
    slide_costs(prs)
    # Team-facing block
    slide_channels(prs)
    slide_dev_setup(prs)
    slide_daily(prs)
    slide_titu_command(prs)
    slide_architecture(prs)
    slide_ai_helps(prs)
    slide_outputs(prs)
    slide_out_of_scope(prs)
    slide_roadmap(prs)
    slide_what_we_need(prs)
    slide_close(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
