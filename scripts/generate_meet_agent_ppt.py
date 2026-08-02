"""NAPCO Nucleus - briefing deck for extending the agent to Google Meet.

Companion to generate_teams_agent_ppt.py and deliberately identical in visual
language (same palette, same header, same card and step shapes) so the two
decks read as one set.

Audience: the development team and team leads. Explains that the Meet side is
a PLAN and not yet live, what already exists and can be reused unchanged, what
still has to be built, the two possible routes, and what has to be decided.

Rules followed (Titu's standing preferences):
  * plain operational language, no em dashes
  * jargon is explained where it appears, never assumed
  * no personal names in slide-visible text
  * no capability is described as working when it is not built yet
  * the developer level detail lives in the speaker notes, not on the slides

Run:
    py -3 scripts\\generate_meet_agent_ppt.py
Output:
    C:\\Users\\khasan\\Downloads\\NAPCO-Nucleus-Meet-Agent.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(r"C:\Users\khasan\Downloads\NAPCO-Nucleus-Meet-Agent.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x2E, 0x8A, 0x8A)
CORAL = RGBColor(0xE0, 0x78, 0x56)
GREEN = RGBColor(0x4A, 0x7A, 0x4A)
GOLD = RGBColor(0xC9, 0x96, 0x2B)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x77, 0x85)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
RULE = RGBColor(0xD5, 0xDC, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line=None,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.25)
    s.shadow.inherit = False
    return s


def _text(slide, x, y, w, h, text, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(space)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return box


def _header(slide, title, kicker=None):
    _rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill=NAVY,
          shape=MSO_SHAPE.RECTANGLE)
    _text(slide, Inches(0.7), Inches(0.16), Inches(11.9), Inches(0.5),
          title, size=29, bold=True, color=WHITE)
    if kicker:
        _text(slide, Inches(0.72), Inches(0.71), Inches(11.9), Inches(0.34),
              kicker, size=13, color=RGBColor(0xC9, 0xDA, 0xEA))


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _card(slide, x, y, w, h, heading, body, accent, hsize=16, bsize=13):
    _rect(slide, x, y, w, h, fill=SOFT, line=RULE)
    _rect(slide, x, y, Inches(0.07), h, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5),
          Inches(0.4), heading, size=hsize, bold=True, color=NAVY)
    _text(slide, x + Inches(0.28), y + Inches(0.66), w - Inches(0.5),
          h - Inches(0.8), body, size=bsize, color=INK, space=3)


def _arrow(slide, x, y, w=Inches(0.42), h=Inches(0.34), color=MUTED):
    return _rect(slide, x, y, w, h, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)


def _bullet_row(slide, x, y, w, text, accent=TEAL, size=15):
    """One bullet line: a small square marker then the text."""
    _rect(slide, x, y + Inches(0.09), Inches(0.14), Inches(0.14), fill=accent,
          shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.34), y, w - Inches(0.34), Inches(0.42),
          text, size=size, color=INK, space=0)


def _pill(slide, x, y, w, h, label, fill, tsize=12):
    """Small status pill, used for the planned / already working markers."""
    _rect(slide, x, y, w, h, fill=fill)
    _text(slide, x, y + Inches(0.05), w, h, label, size=tsize, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)


# ----------------------------------------------------------------- slides ---
def slide_title(prs):
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    _rect(s, 0, Inches(4.55), SLIDE_W, Inches(0.06), fill=GOLD,
          shape=MSO_SHAPE.RECTANGLE)
    _text(s, Inches(1.1), Inches(2.15), Inches(11), Inches(1.1),
          "Napco Nucleus on Google Meet", size=48, bold=True, color=WHITE)
    _text(s, Inches(1.12), Inches(3.15), Inches(11.2), Inches(1.2),
          "The same virtual colleague we already use in MS Teams, extended to "
          "our Google Meet calls.\n"
          "Same result for you: add it to the meeting and the requirements "
          "write themselves.\n"
          "This is the plan and what it needs. The Meet side is not live yet.",
          size=19, color=RGBColor(0xC9, 0xDA, 0xEA), space=4)
    _text(s, Inches(1.12), Inches(4.9), Inches(11), Inches(0.6),
          "Plan briefing  |  August 2026",
          size=16, color=RGBColor(0x9F, 0xB8, 0xCE))
    _notes(s, "Set expectations in the first sentence: Teams is live and "
              "proven, Meet is a plan with a known amount of work left. "
              "Nothing on these slides claims Meet works today.")


def slide_what_it_does(prs):
    s = _blank(prs)
    _header(s, "What it will do on Meet",
            "Exactly what it already does on Teams, same output to the team")
    left = [
        "Joins the Meet call as a participant",
        "Records the discussion",
        "Transcribes it automatically, in Bangla and English",
        "Picks out the requirements we agreed during the call",
    ]
    right = [
        "Writes them into a requirement document",
        "Emails the summary to the team",
        "Keeps Meet recordings in their own separate folder",
        "Uses the same server and the same review step as Teams",
    ]
    y = Inches(1.85)
    for i, item in enumerate(left):
        _bullet_row(s, Inches(0.75), y + i * Inches(0.72), Inches(5.6), item)
    for i, item in enumerate(right):
        _bullet_row(s, Inches(6.95), y + i * Inches(0.72), Inches(5.6), item)

    _rect(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.0),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.58), Inches(11.2), Inches(0.6),
          "Client meetings on Meet stop being the gap in our requirement "
          "record.",
          size=18, bold=True, color=NAVY)
    _notes(s, "The separate folder point matters to the team: Meet material "
              "is kept apart from Teams material on the central server, so "
              "the two sources never get mixed up in a document.")


def slide_how_it_works(prs):
    s = _blank(prs)
    _header(s, "How it will work",
            "Only the first step is new, the rest is the system we already run")
    _rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.5),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(1.82), Inches(11.2), Inches(1.2),
          "One dedicated machine joins the Meet call and captures the audio, "
          "the same way the Teams agent does today. From the moment that audio "
          "reaches the central server, every later step is the system that is "
          "already running and already proven on Teams calls. We are adding a "
          "new way in, not a second pipeline.",
          size=15, color=INK, space=4)

    steps = [
        ("1", "Joins", "Opens the meeting link\nand joins the call.", CORAL,
         "TO BUILD"),
        ("2", "Captures", "Audio reaches the server\nwhile the call is running.",
         GOLD, "PARTLY THERE"),
        ("3", "Transcribes", "Speech becomes text,\nBangla and English.", GREEN,
         "ALREADY WORKS"),
        ("4", "Delivers", "Requirement document\nand summary to the team.",
         GREEN, "ALREADY WORKS"),
    ]
    x = Inches(0.7)
    for num, head, body, col, status in steps:
        _rect(s, x, Inches(3.35), Inches(2.75), Inches(2.45), fill=WHITE,
              line=RULE)
        _rect(s, x, Inches(3.35), Inches(2.75), Inches(0.5), fill=TEAL,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, x + Inches(0.2), Inches(3.43), Inches(2.4), Inches(0.4),
              f"{num}.  {head}", size=15, bold=True, color=WHITE)
        _text(s, x + Inches(0.25), Inches(4.05), Inches(2.3), Inches(1.1),
              body, size=13, color=INK, space=3)
        _pill(s, x + Inches(0.25), Inches(5.15), Inches(1.75), Inches(0.34),
              status, col)
        if num != "4":
            _arrow(s, x + Inches(2.85), Inches(4.25))
        x += Inches(3.05)
    _notes(s, "Developer detail, not for the slide: capture is a system audio "
              "loopback recording, so it does not care whether the sound came "
              "from Teams or a browser. The trigger that starts it is what is "
              "Teams specific today, and that is the piece being replaced.")


def slide_reuse(prs):
    s = _blank(prs)
    _header(s, "What we reuse and what we build",
            "Most of the chain does not need to be touched at all")
    items = [
        ("Reused with no change",
         "The central server, the transcription, the requirement extraction, "
         "the document and the email all work on whatever lands on the server. "
         "They do not need to know the call came from Meet.", GREEN),
        ("Small changes",
         "Recording has to start when a Meet call starts instead of a Teams "
         "call, and Meet recordings have to be filed in their own folder on "
         "the server.", GOLD),
        ("Genuinely new work",
         "Joining the meeting at the right time from the calendar, and "
         "handling the fact that Meet does not tell us who was on the call.",
         CORAL),
    ]
    x = Inches(0.7)
    for head, body, col in items:
        _card(s, x, Inches(1.9), Inches(3.9), Inches(2.75), head, body, col)
        x += Inches(4.15)
    _text(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.9),
          "About two thirds of the work is already done and running in "
          "production for Teams.",
          size=19, bold=True, color=NAVY)
    _notes(s, "Developer detail: the server side walks every folder on the "
              "share and processes whatever it finds, so a new Meet folder is "
              "picked up with no code change. The known trap is the recording "
              "allowlist, which decides at the end of a call whether to keep "
              "the audio based on who was on it. Meet gives us no participant "
              "list, so that check has to be handled explicitly or it will "
              "throw good recordings away.")


def slide_routes(prs):
    s = _blank(prs)
    _header(s, "Two ways to do it",
            "One is much less work if our Google plan supports it")
    _card(s, Inches(0.7), Inches(1.75), Inches(5.85), Inches(3.0),
          "Option A: use Google's own recording",
          "If our Google Workspace plan includes Meet recording, the meeting "
          "is recorded by Google itself and the recording and transcript are "
          "saved to Drive. We already have the piece that picks files up from "
          "Drive. Nothing joins the call, nothing is captured by us, and Meet "
          "shows its normal recording notice to everyone.", GREEN,
          hsize=18, bsize=14)
    _card(s, Inches(6.9), Inches(1.75), Inches(5.7), Inches(3.0),
          "Option B: a dedicated machine joins",
          "A dedicated machine opens the meeting link and joins the call, then "
          "records the audio the same way the Teams agent does. This works on "
          "any plan and gives us full control, but it is more moving parts and "
          "someone has to let it into the meeting.", GOLD,
          hsize=18, bsize=14)

    _rect(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.3),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.24), Inches(11.2), Inches(0.4),
          "Recommendation", size=16, bold=True, color=NAVY)
    _text(s, Inches(1.05), Inches(5.66), Inches(11.2), Inches(0.6),
          "Check the Google plan first. If recording is included, Option A "
          "removes most of the work and is the cleaner answer on privacy. "
          "Option B is the fallback and stays available either way.",
          size=14, color=INK, space=3)
    _notes(s, "The plan check is a ten minute job in the Google admin console "
              "under billing and subscriptions. Do it before any code is "
              "written, because Option A deletes three of the four build "
              "items. Option B has one known blocker to solve first: on a "
              "remote desktop session the sound is redirected away from the "
              "machine, so it records silence. It has to run at the physical "
              "console.")


def slide_decisions(prs):
    s = _blank(prs)
    _header(s, "What we need to decide",
            "Four answers and the build can start")
    rows = [
        "Which Google plan we are on, so we know if Google can record for us",
        "Which machine hosts the Meet agent, and that it is not the Teams one",
        "How the agent gets invited, so it only attends meetings it should",
        "Whether Meet summaries go out with the Teams ones or separately",
    ]
    y = Inches(1.9)
    for i, item in enumerate(rows):
        _bullet_row(s, Inches(0.75), y + i * Inches(0.78), Inches(11.7), item,
                    accent=TEAL, size=17)

    _rect(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.55),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.37), Inches(11.2), Inches(0.4),
          "One thing to say plainly", size=16, bold=True, color=CORAL)
    _text(s, Inches(1.05), Inches(5.81), Inches(11.2), Inches(0.8),
          "If we capture the audio ourselves, Meet will not show its usual "
          "recording notice, exactly as on Teams today. If we use Google's own "
          "recording, everyone sees the notice. Please raise any concern "
          "directly.",
          size=14, color=INK, space=3)
    _notes(s, "Say the recording line out loud, do not leave it only on the "
              "slide. It is the same disclosure that already goes out on every "
              "summary email. Close by asking for the plan check, because "
              "everything else depends on that answer.")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_what_it_does(prs)
    slide_how_it_works(prs)
    slide_reuse(prs)
    slide_routes(prs)
    slide_decisions(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"[OK] {len(prs.slides)} slides -> {OUT}")


if __name__ == "__main__":
    build()
