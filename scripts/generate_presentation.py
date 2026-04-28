"""Generate the NAPCO Nucleus presentation deck (.pptx).

Targeted for a 1-hour talk on 2026-04-30 (Thursday).
~29 slides, 16:9, with speaker notes on every slide.

Run:
    py -3 scripts/generate_presentation.py
Output:
    docs/NAPCO-Nucleus-Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT_PATH = ROOT / "docs" / "NAPCO-Nucleus-Presentation.pptx"

# 16:9 slide size
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette (mirrors the technical guide)
NAVY    = RGBColor(0x1F, 0x4E, 0x79)
TEAL    = RGBColor(0x2E, 0x8A, 0x8A)
CORAL   = RGBColor(0xE0, 0x78, 0x56)
GREEN   = RGBColor(0x4A, 0x7A, 0x4A)
GOLD    = RGBColor(0xC9, 0x96, 0x2B)
PURPLE  = RGBColor(0x6A, 0x4C, 0x93)
INK     = RGBColor(0x22, 0x22, 0x22)
MUTED   = RGBColor(0x6B, 0x77, 0x85)
SOFT    = RGBColor(0xF5, 0xF7, 0xFA)
RULE    = RGBColor(0xD5, 0xDC, 0xE5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x2A, 0x38)


def _no_fill(shape):
    shape.fill.background()


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _set_text(tf, text, size=18, bold=False, color=INK, font="Calibri",
              align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(tf, text, size=18, bold=False, color=INK, font="Calibri",
              align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _bullets(tf, items, size=16, color=INK):
    """Each item is a (bold_lead, body) tuple, or just body string."""
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            p.space_before = Pt(6)
        p.level = 0
        bullet_run = p.add_run()
        bullet_run.text = "• "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = NAVY
        bullet_run.font.bold = True
        if isinstance(item, tuple):
            lead, body = item
            r1 = p.add_run()
            r1.text = lead + " "
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = body
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color


# ── slide chrome ────────────────────────────────────────────────────

def _add_chrome(slide, page_num, total, footer_text):
    # top navy bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.10))
    _solid(bar, NAVY)
    _no_line(bar)

    # bottom rule + footer
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.18), Inches(12.333), Emu(9000))
    _solid(rule, RULE)
    _no_line(rule)

    foot_l = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.22), Inches(8), Inches(0.25))
    _set_text(foot_l.text_frame, footer_text, size=9, color=MUTED)

    foot_r = slide.shapes.add_textbox(
        Inches(11.5), Inches(7.22), Inches(1.4), Inches(0.25))
    _set_text(foot_r.text_frame, f"{page_num} / {total}",
              size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.30), Inches(12.333), Inches(0.7))
    _set_text(box.text_frame, text, size=28, bold=True, color=NAVY)
    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.95), Inches(12.333), Inches(0.4))
        _set_text(sub.text_frame, subtitle, size=14, color=MUTED)


def _section_band(slide, x, y, w, color, text):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.36))
    _solid(band, color)
    _no_line(band)
    band.text_frame.margin_left = Inches(0.15)
    band.text_frame.margin_right = Inches(0.15)
    band.text_frame.margin_top = Inches(0.04)
    band.text_frame.margin_bottom = Inches(0.04)
    _set_text(band.text_frame, text, size=12, bold=True, color=WHITE)


def _card(slide, x, y, w, h, header, color, body_lines=None,
          body_size=14):
    """Header strip + white body. body_lines is list of strings or
    (bold_lead, body) tuples."""
    head_h = Inches(0.36)
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, head_h)
    _solid(head, color)
    _no_line(head)
    head.text_frame.margin_left = Inches(0.15)
    head.text_frame.margin_top = Inches(0.05)
    _set_text(head.text_frame, header, size=12, bold=True, color=WHITE)

    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, y + head_h, w, h - head_h)
    _solid(body, WHITE)
    _line(body, RULE, 0.5)
    body.text_frame.margin_left = Inches(0.18)
    body.text_frame.margin_right = Inches(0.18)
    body.text_frame.margin_top = Inches(0.12)
    body.text_frame.margin_bottom = Inches(0.12)
    body.text_frame.word_wrap = True
    if body_lines:
        _bullets(body.text_frame, body_lines, size=body_size)
    return body


def _metric_tile(slide, x, y, w, h, value, label, color):
    tile = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(tile, SOFT)
    _no_line(tile)
    # accent top stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
    _solid(stripe, color)
    _no_line(stripe)
    # value
    val = slide.shapes.add_textbox(x, y + Inches(0.30),
                                    w, Inches(0.85))
    _set_text(val.text_frame, value, size=40, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER)
    # label
    lbl = slide.shapes.add_textbox(x, y + h - Inches(0.5),
                                    w, Inches(0.4))
    _set_text(lbl.text_frame, label, size=11, color=MUTED,
              align=PP_ALIGN.CENTER)


def _box(slide, x, y, w, h, fill, label, sublabel=None,
         label_size=14, sublabel_size=10):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _solid(sh, fill)
    _no_line(sh)
    sh.adjustments[0] = 0.10
    tf = sh.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, label, size=label_size, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)
    if sublabel:
        _add_text(tf, sublabel, size=sublabel_size, color=WHITE,
                  align=PP_ALIGN.CENTER, space_before=2)


def _arrow(slide, x1, y1, x2, y2, color=MUTED, width_pt=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    # add arrowhead
    ln = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")


def _code_block(slide, x, y, w, h, lines, font_size=11):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(box, CODE_BG)
    _no_line(box)
    box.text_frame.margin_left = Inches(0.18)
    box.text_frame.margin_right = Inches(0.18)
    box.text_frame.margin_top = Inches(0.14)
    box.text_frame.margin_bottom = Inches(0.14)
    box.text_frame.word_wrap = True
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Consolas"
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE


def _table(slide, x, y, w, rows, col_widths, header_color=NAVY,
           font_size=12, header_size=12):
    """Build a styled table at (x, y) with given total width w."""
    # Compute total height first via row height heuristic
    n_rows = len(rows)
    # Use a simple uniform height; PowerPoint will autosize text
    base_row_h = Inches(0.42)
    h = base_row_h * n_rows
    tbl_shape = slide.shapes.add_table(n_rows, len(col_widths), x, y, w, h)
    table = tbl_shape.table

    total_w = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = int(w * cw / total_w)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else SOFT
                run.font.size = Pt(font_size)
                run.font.color.rgb = INK
    return tbl_shape


def _notes(slide, text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = text


# ── slide builders ──────────────────────────────────────────────────

FOOTER = "NAPCO Nucleus  |  Mohammad Kamrul Hasan  |  Adaptive Enterprise Limited"


def _new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def slide_title(prs, total):
    s = _new_slide(prs)
    # full bleed navy band on left
    band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), SLIDE_H)
    _solid(band, NAVY)
    _no_line(band)
    # accent rectangle
    acc = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(0.5), Inches(2.5))
    _solid(acc, GOLD)
    _no_line(acc)

    # eyebrow
    eb = s.shapes.add_textbox(Inches(1.0), Inches(0.9),
                               Inches(11), Inches(0.35))
    _set_text(eb.text_frame, "AEL  ·  Engineering Showcase  ·  30 April 2026",
              size=12, bold=True, color=GOLD)

    # title
    t = s.shapes.add_textbox(Inches(1.0), Inches(1.6),
                              Inches(11.5), Inches(2.0))
    tf = t.text_frame
    tf.word_wrap = True
    _set_text(tf, "NAPCO Nucleus", size=54, bold=True, color=NAVY)
    _add_text(tf, "An Autonomous Intelligent Orchestrator",
              size=30, bold=True, color=NAVY, space_before=4)

    # subtitle
    sub = s.shapes.add_textbox(Inches(1.0), Inches(4.1),
                                Inches(11.5), Inches(0.6))
    _set_text(sub.text_frame,
              "Autonomous AI Orchestrator for QA & Project Governance",
              size=18, color=TEAL)

    # author block
    auth = s.shapes.add_textbox(Inches(1.0), Inches(5.2),
                                 Inches(11.5), Inches(1.4))
    af = auth.text_frame
    af.word_wrap = True
    _set_text(af, "Developed & Engineered by", size=11, color=MUTED)
    _add_text(af, "Mohammad Kamrul Hasan", size=22, bold=True, color=NAVY,
              space_before=2)
    _add_text(af, "AI-Augmented QA Architect  |  Adaptive Enterprise Limited",
              size=14, color=INK, space_before=2)

    # bottom strip
    foot = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.30), SLIDE_W, Inches(0.30))
    _solid(foot, NAVY)
    _no_line(foot)
    fb = s.shapes.add_textbox(Inches(0.6), SLIDE_H - Inches(0.28),
                               Inches(12), Inches(0.25))
    _set_text(fb.text_frame,
              "Engineering Showcase  ·  Thursday 30 April 2026  ·  60 minutes",
              size=10, color=WHITE)

    _notes(s,
        "Open by setting the stakes: this is not a slide-deck about a "
        "future plan. NAPCO Nucleus is in production, running 9 workflows "
        "every day. The talk will cover what it is, how it's built, why "
        "it works, what it costs, and what comes next. Speak briefly to "
        "the title — \"Autonomous Intelligent Orchestrator\" is the "
        "honest one-line description, not marketing. The NAPCO Nucleus "
        "Framework is the architectural pattern reused from prior work.")


def slide_agenda(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "Agenda",
               "Sixty minutes, eight sections, one architecture")

    sections = [
        ("01", "The Case",          "Why NN exists, the thesis, the impact",      NAVY),
        ("02", "Architecture",      "Component map, single-agent model, lifecycle", TEAL),
        ("03", "Inside NN",         "Orchestrator, SDK, MCP tools, memory",        CORAL),
        ("04", "Prompts",           "Algorithms in markdown, the 8 prompt files",  GREEN),
        ("05", "Integrations & Stack", "Everything NN talks to, end to end",       GOLD),
        ("06", "CI/CD & Operations","Workflows, runner, guardrails",               PURPLE),
        ("07", "Extending NN",      "Adding workflows and tools — the recipes",    NAVY),
        ("08", "Closing",           "Benefits, limits, roadmap, Q&A",              TEAL),
    ]

    cols = 2
    rows = 4
    x0 = Inches(0.6)
    y0 = Inches(1.7)
    cw = Inches(6.1)
    ch = Inches(1.20)
    gx = Inches(0.13)
    gy = Inches(0.16)
    for i, (num, title, sub, color) in enumerate(sections):
        r = i // cols
        c = i % cols
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        # number block
        nb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, y, Inches(1.0), ch)
        _solid(nb, color)
        _no_line(nb)
        nb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(nb.text_frame, num, size=28, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        # body
        bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(1.0), y, cw - Inches(1.0), ch)
        _solid(bb, WHITE)
        _line(bb, RULE, 0.5)
        bb.text_frame.margin_left = Inches(0.18)
        bb.text_frame.margin_top = Inches(0.12)
        bb.text_frame.word_wrap = True
        _set_text(bb.text_frame, title, size=16, bold=True, color=NAVY)
        _add_text(bb.text_frame, sub, size=11, color=MUTED, space_before=2)

    _notes(s,
        "Walk through agenda once, then commit to the timing. 60 min total: "
        "≈8 min on Case + Architecture overview, ≈25 min on Inside NN + "
        "Prompts + Integrations (the meat), ≈12 min on CI/CD + Extending, "
        "≈10 min Benefits/Limits/Roadmap, leaving ≈5 min for Q&A. If we "
        "run long, the deepest cut is the prompts section — it's the most "
        "interesting but I can summarize and direct people to the docs.")


def slide_problem(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "01 · The Problem",
               "What QA looked like before NAPCO Nucleus")

    # left: before
    b1 = _card(s, Inches(0.6), Inches(1.7), Inches(6.05), Inches(5.2),
               "BEFORE — THE FRAGMENTED REALITY", CORAL,
               body_lines=[
                   ("4 test runs overnight,", "4 separate emails, no consolidation"),
                   ("Requirements lost in chat,", "spoken on calls, never filed"),
                   ("Manual OpenProject issue creation", "with frequent duplicates"),
                   ("Daily reports built by hand", "every morning, ~45 minutes"),
                   ("Test failures triaged by humans", "before classification"),
                   ("Memory in heads, not systems —", "context died with people"),
               ], body_size=14)

    # right: after
    b2 = _card(s, Inches(6.85), Inches(1.7), Inches(6.05), Inches(5.2),
               "AFTER — NAPCO NUCLEUS IN PRODUCTION", GREEN,
               body_lines=[
                   ("Two consolidated emails per day —", "09:00 detailed, 09:30 executive"),
                   ("Requirements captured automatically", "from email + Drive + meetings"),
                   ("OpenProject Work Packages filed twice daily", "with dedup + in-place update on revisions"),
                   ("Reports generated by AI in minutes,", "not by humans for 45"),
                   ("Failures classified by Claude", "with reasoning in markdown prompts"),
                   ("Memory in SQLite committed to git —", "audit trail survives forever"),
               ], body_size=14)

    _notes(s,
        "Frame the problem in human terms. The test team was drowning in "
        "fragmented signal: four overnight runs producing four mailbox "
        "items, requirements that lived in chat threads or call notes "
        "until somebody remembered to file them, OpenProject Work Packages with "
        "frequent duplicates because nobody could remember what was "
        "already filed. NN's mandate was to consolidate, not to replace. "
        "Emphasize: human QA judgment is unchanged. NN takes over the "
        "logistics — what to file, when to file, how to format — so "
        "humans focus on the things that need a human.")


def slide_thesis(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "01 · The Thesis",
               "One sentence that explains every design choice in NN")

    # big quote-style thesis
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.7),
                              Inches(12.13), Inches(2.2))
    _solid(box, NAVY)
    _no_line(box)
    box.text_frame.margin_left = Inches(0.5)
    box.text_frame.margin_right = Inches(0.5)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(box.text_frame,
              "Reasoning lives in prompts. Tools wrap I/O.",
              size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(box.text_frame,
              "Python does no algorithmic work. To change behavior, edit a markdown file. "
              "To add capability, wrap one external system as a tool.",
              size=15, color=GOLD, align=PP_ALIGN.CENTER, space_before=8)

    # 3 supporting pillars
    items = [
        ("ONE TURN PER PROCESS",
         "agent.py boots, runs ONE Claude turn, exits. No daemons, no "
         "long-running graph. Crash-safe by construction.", TEAL),
        ("MEMORY IN GIT",
         "SQLite + FTS5 committed to main after every run. The repo IS "
         "the database. Audit trail and dedup history are version-controlled.",
         CORAL),
        ("CLAUDE MAX, NOT API",
         "No ANTHROPIC_API_KEY. Reasoning runs through the local Claude "
         "Code CLI under a Claude Max subscription. Cost is fixed monthly.",
         PURPLE),
    ]
    cw = Inches(4.05)
    gx = Inches(0.10)
    x0 = Inches(0.6)
    y = Inches(4.20)
    for i, (lead, body, color) in enumerate(items):
        x = x0 + i * (cw + gx)
        _card(s, x, y, cw, Inches(2.7), lead, color,
              body_lines=[("", body)], body_size=13)

    _notes(s,
        "This slide is the centerpiece. If the audience remembers one "
        "thing from the talk, it should be this: reasoning lives in "
        "prompts, tools wrap I/O. Repeat that exact phrase. Then explain "
        "the three pillars: one-turn-per-process means crash safety and "
        "no leaked state; memory-in-git means we get version control "
        "for free; Claude Max means cost ceiling is fixed at the "
        "subscription level, not metered per token. The audience will "
        "ask about cost — be ready with the ~$200/mo number.")


def slide_two_dimensions(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "01 · What NN Does",
               "Two operational dimensions, one agent")

    # PM dimension
    _card(s, Inches(0.6), Inches(1.7), Inches(6.05), Inches(5.2),
          "PROJECT MANAGEMENT DIMENSION", TEAL,
          body_lines=[
              ("Requirement Management —", "13:00 + 01:00 BDT, ingests Email (IMAP) + Google Drive PDFs + meeting audio (Whisper) + Teams channel, publishes to OpenProject Work Packages with dedup + in-place update"),
              ("Daily Report (Detailed) —", "03:30 BDT, 6 sections, full team"),
              ("Daily Report (Summary) —", "03:45 BDT, 7 blocks, leadership"),
              ("", ""),
              ("Reads state from the world,", "produces governance artifacts."),
              ("Where strategy meets execution.", ""),
          ], body_size=14)

    # Test Automation dimension
    _card(s, Inches(6.85), Inches(1.7), Inches(6.05), Inches(5.2),
          "TEST AUTOMATION DIMENSION", CORAL,
          body_lines=[
              ("API Functional Test —", "01:00 BDT, Newman + Postman"),
              ("API Integration Test —", "01:15 BDT, pytest + regression diff"),
              ("API Load Test —", "01:30 BDT, Locust 10 → 10,000 users"),
              ("MVP Access E2E Test —", "02:30 BDT, Playwright full suite"),
              ("MVPAccess CICD —", "03:00 BDT, TFS + MSBuild + IIS deploy"),
              ("", ""),
              ("Pure execution.", "Run, classify, report."),
          ], body_size=14)

    _notes(s,
        "This is where I land the central architectural insight: NN is "
        "not 'just' a test bot or 'just' a PM bot. It's both, because "
        "the same agent.py + Claude Agent SDK + MCP server can swing "
        "between dimensions purely by loading a different prompt. The "
        "two-dimension framing is also how we explain it to "
        "non-engineers: product/leadership care about the PM side, the "
        "test team cares about the automation side, but it's one system "
        "they're both consuming.")


def slide_team_streams(prs, total, page):
    """2x2 grid: how NN delivers value to the team across 4 streams."""
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "01 · Four Streams to the Team",
               "How NAPCO Nucleus multiplies QA + PM throughput")

    # color, status, name, schedule, what it does, output
    streams = [
        (TEAL,   "LIVE",
         "Requirement Management",
         "Twice daily · 13:00 + 01:00 BDT",
         "Ingests requirements from FOUR channels — Email (IMAP), Google Drive PDFs, meeting audio via Groq Whisper, and Teams channel forwards. Splits each requirement into 3-hour tasks. Publishes to OpenProject with dedup + in-place update on revisions",
         "Categorized Work Packages with full audit trail in the Activity tab"),
        (CORAL,  "LIVE",
         "API Testing — Functional · Integration · Load",
         "Daily · 02:00 BDT",
         "Newman + pytest + Locust suites, classifies failures, regression diff against prior runs",
         "Per-run PDFs + per-test rows in test_run_history"),
        (GREEN,  "LIVE",
         "End-to-End Testing",
         "Daily · 02:00 BDT",
         "Full Playwright suite headless on Chromium, screenshots on failure, regression detection",
         "PDF + screenshots + memory entries"),
        (PURPLE, "PENDING",
         "CI/CD — TFS · MSBuild · IIS",
         "Daily · 03:00 BDT (workflow ready)",
         "Pulls source from TFS, MSBuilds the solution, deploys to IIS — gated on 6 IT-managed secrets",
         "Build artifacts + deployment status email"),
    ]

    card_w = Inches(6.05)
    card_h = Inches(2.55)
    margin_x = Inches(0.6)
    top_y = Inches(1.6)
    gap = Inches(0.18)

    for i, (color, status, name, schedule, what, output) in enumerate(streams):
        row, col = divmod(i, 2)
        x = margin_x + col * (card_w + gap)
        y = top_y + row * (card_h + gap)

        # Header band with stream name
        head_h = Inches(0.42)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, head_h)
        _solid(head, color)
        _no_line(head)
        head.text_frame.margin_left = Inches(0.18)
        head.text_frame.margin_top = Inches(0.06)
        head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(head.text_frame, name, size=13, bold=True, color=WHITE)

        # Status pill (top-right)
        pill_w = Inches(0.85)
        pill_h = Inches(0.30)
        pill_x = x + card_w - pill_w - Inches(0.12)
        pill_y = y + Inches(0.06)
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   pill_x, pill_y, pill_w, pill_h)
        _solid(pill, WHITE if status == "LIVE" else GOLD)
        _no_line(pill)
        pill.adjustments[0] = 0.5
        pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(pill.text_frame, status,
                  size=9, bold=True,
                  color=color if status == "LIVE" else INK,
                  align=PP_ALIGN.CENTER)

        # Body
        body = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, y + head_h, card_w, card_h - head_h)
        _solid(body, WHITE)
        _line(body, RULE, 0.5)
        body.text_frame.margin_left = Inches(0.18)
        body.text_frame.margin_right = Inches(0.18)
        body.text_frame.margin_top = Inches(0.10)
        body.text_frame.margin_bottom = Inches(0.10)
        body.text_frame.word_wrap = True

        _bullets(body.text_frame, [
            ("Schedule —", schedule),
            ("What it does —", what),
            ("Output —", output),
        ], size=11)

    _notes(s,
        "This is the boss-facing slide. Anchor every answer to the four "
        "streams. Three are live and producing artifacts you can click "
        "through to today; the fourth (CI/CD) is code-complete and "
        "waiting only on IT secrets — be honest about that, it shows "
        "the gap is operational not technical. The 'twice daily' cadence "
        "for Req Mgmt was tightened on 2026-04-28 from a 5-slot "
        "business-hours poll. Same code, fewer runs, less LLM cost.")


def slide_requirement_pipeline(prs, total, page):
    """Small colorful flow diagram: 4 channels → NN agent → OpenProject."""
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "01 · Requirement Gathering Pipeline",
               "Four input channels, one OpenProject backlog")

    # Row 1: 4 channel boxes (colorful, small)
    channels = [
        (TEAL,   "Email",      "IMAP poll · allowlisted senders"),
        (CORAL,  "Drive PDF",  "Google Drive · pypdf extract"),
        (GREEN,  "Audio File", "Drive · Groq Whisper transcribe"),
        (PURPLE, "Teams",      "Power Automate forwards to inbox"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(1.0)
    gap_x = Inches(0.18)
    total_w = 4 * box_w + 3 * gap_x
    start_x = (SLIDE_W - total_w) // 2
    row1_y = Inches(1.85)

    centers = []
    for i, (color, name, sub) in enumerate(channels):
        x = start_x + i * (box_w + gap_x)
        _box(s, x, row1_y, box_w, box_h, color, name, sub,
             label_size=18, sublabel_size=10)
        centers.append(x + box_w // 2)

    # Center: NAPCO Nucleus Agent
    nn_w = Inches(5.6)
    nn_h = Inches(1.30)
    nn_x = (SLIDE_W - nn_w) // 2
    nn_y = Inches(3.55)
    _box(s, nn_x, nn_y, nn_w, nn_h, NAVY,
         "NAPCO Nucleus Agent",
         "Memory check · Read inbox · Identify requirements · "
         "Dedup vs prior · Classify (new / Bug / updatedRequirements) · "
         "Split into 3-hour tasks",
         label_size=18, sublabel_size=10)

    # Bottom: OpenProject
    op_w = Inches(8.4)
    op_h = Inches(1.20)
    op_x = (SLIDE_W - op_w) // 2
    op_y = Inches(5.45)
    _box(s, op_x, op_y, op_w, op_h, GOLD,
         "OpenProject — mvp-access backlog",
         "Type=Task or Bug · Category=AccessGroup / BadgeHolder / Personnel · "
         "Status=New on create · Subject swap + comment on revision",
         label_size=17, sublabel_size=10)

    # Arrows: each channel → NN agent
    for cx in centers:
        _arrow(s, cx, row1_y + box_h, cx, nn_y, color=MUTED, width_pt=1.5)
    # Arrow: NN → OpenProject
    _arrow(s, nn_x + nn_w // 2, nn_y + nn_h,
              nn_x + nn_w // 2, op_y, color=NAVY, width_pt=2.0)

    # Side caption: the dedup gates
    side_x = Inches(0.5)
    side_y = Inches(3.55)
    side_w = Inches(2.6)
    side_h = Inches(1.30)
    _card(s, side_x, side_y, side_w, side_h,
          "DEDUP GATES", CORAL,
          body_lines=[
              "Layer 1 · Open WP titles in OpenProject",
              "Layer 2 · requirements_seen FTS in SQLite",
              "Both populated on every run (self-heal)",
          ], body_size=10)

    side_x2 = SLIDE_W - side_w - Inches(0.5)
    _card(s, side_x2, side_y, side_w, side_h,
          "ON REVISION", PURPLE,
          body_lines=[
              "Subject swap on existing Work Package",
              "Timestamped comment with new spec",
              "Activity tab preserves full diff",
          ], body_size=10)

    _notes(s,
        "Walk through the diagram top-down: four channels feed the "
        "agent; the agent reasons + classifies; OpenProject is the "
        "single backlog of record. Side cards explain the two "
        "non-obvious mechanics — how dedup works (two stacked layers, "
        "both self-healing) and what 'in-place update' means in "
        "practice (subject swap + comment, OpenProject's Activity tab "
        "is the audit trail).")


def slide_headline_numbers(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "01 · Impact at a Glance",
               "What NN delivers, in numbers")

    metrics = [
        ("9",     "PRODUCTION WORKFLOWS",  NAVY),
        ("31",    "MCP TOOLS",             TEAL),
        ("8",     "PROMPT FILES",          GREEN),
        ("~3.6k", "PYTHON LOC",            GOLD),
        ("2",     "EMAILS PER DAY",        CORAL),
        ("$200",  "FIXED MONTHLY COST",    PURPLE),
    ]
    cols = 3
    rows = 2
    x0 = Inches(0.6)
    y0 = Inches(1.8)
    cw = Inches(4.05)
    ch = Inches(2.30)
    gx = Inches(0.10)
    gy = Inches(0.20)
    for i, (v, l, c) in enumerate(metrics):
        r = i // cols
        col = i % cols
        x = x0 + col * (cw + gx)
        y = y0 + r * (ch + gy)
        _metric_tile(s, x, y, cw, ch, v, l, c)

    # caption
    cap = s.shapes.add_textbox(Inches(0.6), Inches(6.7),
                                Inches(12), Inches(0.4))
    _set_text(cap.text_frame,
              "One self-hosted runner, one repo, one .env. Zero cloud bill, "
              "zero per-token billing.", size=13, color=MUTED,
              align=PP_ALIGN.CENTER)

    _notes(s,
        "Use the numbers to anchor the conversation. The most surprising "
        "ones for the audience: $200/month fixed (vs. token-billed "
        "alternatives that scale with traffic), 3.6k LOC (small enough "
        "for one person to own), 8 prompt files (the real intellectual "
        "asset of the system). The two emails per day number replaces "
        "the previous 6+ overnight messages reality.")


def slide_arch_overview(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "02 · Architecture — Component Map",
               "agent.py at the center, talking to six things")

    # Drawing area (within content)
    cx = Inches(6.65)
    cy = Inches(4.3)
    aw, ah = Inches(2.6), Inches(0.85)
    _box(s, cx - aw / 2, cy - ah / 2, aw, ah, NAVY,
         "agent.py", "Python orchestrator",
         label_size=18, sublabel_size=11)

    bw, bh = Inches(2.6), Inches(0.80)
    col = Inches(3.4)   # horizontal offset of node center from cx
    row = Inches(1.6)   # vertical offset

    nodes = [
        (-1,  1, TEAL,   ".env / config",       "secrets, paths, CLI override"),
        ( 1,  1, GOLD,   "Claude Agent SDK",    "ClaudeSDKClient + MCP server"),
        (-1,  0, PURPLE, "MCP server",          "31 tools, in-process"),
        ( 1,  0, GREEN,  "Prompts",             "system.md + <task>.md"),
        (-1, -1, CORAL,  "GitHub Actions",      "9 workflows, self-hosted"),
        ( 1, -1, NAVY,   "nucleus_memory.db",   "SQLite + FTS5, in git"),
    ]

    for c, r, color, lab, sub in nodes:
        ncx = cx + c * col
        ncy = cy + r * row
        _box(s, ncx - bw / 2, ncy - bh / 2, bw, bh, color, lab, sub,
             label_size=14, sublabel_size=10)
        # arrow source: edge of node closest to center
        if r == 0:
            sx = ncx + (bw / 2 if c < 0 else -bw / 2)
            sy = ncy
            tx = cx + (-aw / 2 if c > 0 else aw / 2)
            ty = cy
        else:
            sx = ncx + (bw / 2 if c < 0 else -bw / 2)
            sy = ncy + (-bh / 2 if r > 0 else bh / 2)
            tx = cx + (-aw / 4 if c < 0 else aw / 4)
            ty = cy + (ah / 2 if r > 0 else -ah / 2)
        _arrow(s, sx, sy, tx, ty, color=MUTED, width_pt=1.4)

    # caption
    cap = s.shapes.add_textbox(Inches(0.6), Inches(6.7),
                                Inches(12), Inches(0.4))
    _set_text(cap.text_frame,
              "Every arrow is a contract. agent.py composes them; nothing else does.",
              size=12, color=MUTED, align=PP_ALIGN.CENTER)

    _notes(s,
        "Walk through each node clockwise from top-left: .env feeds "
        "secrets and paths; Claude Agent SDK is how Python talks to "
        "Claude; Prompts on the right are the actual algorithms; "
        "nucleus_memory.db is where state persists across runs; "
        "GitHub Actions is the trigger layer; MCP server is the tool "
        "surface Claude calls into. Emphasize: these six external "
        "boundaries are the entire system. There is nothing else.")


def slide_single_agent(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "02 · The Single-Agent Model",
               "One agent.py, many prompts — not a multi-agent graph")

    _card(s, Inches(0.6), Inches(1.7), Inches(12.13), Inches(2.0),
          "WHY NOT A MULTI-AGENT FRAMEWORK?", PURPLE,
          body_lines=[
              ("Industry default —", "LangGraph / CrewAI / AutoGen with planner + worker + critic agents communicating over a graph"),
              ("NN's choice —", "ONE agent.py. Different markdown prompt per workflow. State persists via SQLite, not via inter-agent message passing"),
              ("Why —", "We don't need the orchestration overhead. Claude is good enough to drive a single workflow end-to-end if the prompt is precise"),
          ], body_size=13)

    _card(s, Inches(0.6), Inches(3.85), Inches(6.05), Inches(3.0),
          "WHAT WE GAIN", GREEN,
          body_lines=[
              ("Clean state", "between workflows — no leaked context"),
              ("One ground truth", "per workflow — read the prompt, know the loop"),
              ("Crash-safe", "— process death leaves SQLite consistent"),
              ("Cost-bounded", "— SDK session closes when the prompt loop ends"),
              ("Easy to reason about", "— each run is a function call"),
          ], body_size=13)

    _card(s, Inches(6.85), Inches(3.85), Inches(6.05), Inches(3.0),
          "WHAT WE GIVE UP", CORAL,
          body_lines=[
              ("No inter-agent debate", "— Claude reasons solo, no critic loop"),
              ("No long-running planner", "— each turn is independent"),
              ("Single point of orchestration", "— if agent.py errors out, the whole run dies"),
              ("Honest tradeoff —", "we accept these in exchange for simplicity"),
          ], body_size=13)

    _notes(s,
        "This is where engineers in the audience will push back. Be "
        "ready: 'why not LangGraph?' Answer: we tried, and the "
        "complexity exceeded the gain. For workflows that are "
        "fundamentally sequential — poll inbox, summarize, file issue — "
        "a planner/worker split adds latency and failure modes without "
        "improving outcomes. We measured. Single agent + great prompt "
        "wins.")


def slide_lifecycle(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "02 · Universal Execution Lifecycle",
               "Every workflow follows this exact shape")

    steps = [
        ("1", "Trigger",   "cron / dispatch",       NAVY),
        ("2", "Checkout",  "actions/checkout",      TEAL),
        ("3", "Install",   "pip install",           CORAL),
        ("4", "Boot",      "py agent.py",           GREEN),
        ("5", "SDK open",  "Claude session",        GOLD),
        ("6", "Loop",      "tool calls",            PURPLE),
        ("7", "Persist",   "git commit",            NAVY),
        ("8", "Push",      "to main",               TEAL),
    ]

    n = len(steps)
    avail = Inches(12.13)
    gap = Inches(0.10)
    bw = (avail - gap * (n - 1)) / n
    bh = Inches(1.6)
    y = Inches(2.0)
    x0 = Inches(0.6)
    for i, (num, title, sub, color) in enumerate(steps):
        x = x0 + i * (bw + gap)
        # numbered head
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, y, bw, Inches(0.40))
        _solid(head, color)
        _no_line(head)
        head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(head.text_frame, f"STEP {num}", size=10, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER)
        body = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, y + Inches(0.40), bw, bh - Inches(0.40))
        _solid(body, WHITE)
        _line(body, RULE, 0.5)
        body.text_frame.margin_top = Inches(0.10)
        body.text_frame.word_wrap = True
        body.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(body.text_frame, title, size=14, bold=True, color=NAVY,
                  align=PP_ALIGN.CENTER)
        _add_text(body.text_frame, sub, size=10, color=MUTED,
                  align=PP_ALIGN.CENTER, space_before=2)

    # detail row
    detail_y = Inches(4.0)
    d_card = _card(s, Inches(0.6), detail_y, Inches(12.13), Inches(2.9),
                    "WHAT EACH STEP DOES", NAVY,
                    body_lines=[
                        ("1·Trigger fires —", "GitHub Actions schedule (cron) or workflow_dispatch starts the job. Pinned to the self-hosted Windows runner."),
                        ("2·Repo checkout —", "actions/checkout@v5 clones the repo. Includes nucleus_memory.db and data/ from prior runs — that's how memory crosses runs."),
                        ("3·Install deps —", "pip install -r requirements.txt. Fast on warm runner."),
                        ("4·agent.py boots —", "Loads .env with override=True. Forces UTF-8 stdout to survive Windows cp1252 console."),
                        ("5·SDK opens —", "create_sdk_mcp_server registers all 31 tools. ClaudeSDKClient opens session against local Claude Code CLI (no API key)."),
                        ("6·The loop runs —", "client.query(kickoff). Claude reads the prompt, calls MCP tools, streams results. agent.py prints text blocks."),
                        ("7·State persists —", "Tool calls wrote to SQLite during the run. Final step stages nucleus_memory.db + data/, commits, rebases."),
                        ("8·Push to main —", "git push lands new memory state. Next workflow picks it up at checkout. Concurrency group prevents overlap."),
                    ], body_size=11)

    _notes(s,
        "Pace this slide deliberately — it's the backbone. Each step "
        "takes seconds, except step 6 which is the actual work (1-15 "
        "minutes depending on workflow). The key insight to deliver: "
        "this lifecycle is universal across all 9 workflows. Adding a "
        "new workflow doesn't change the shape, only what's inside "
        "step 6.")


def slide_design_constraints(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "02 · Design Constraints",
               "Five rules that shaped the codebase")

    items = [
        ("ONE AGENT, MANY PROMPTS",
         "Not a multi-agent graph. The same agent.py loads a different prompt per workflow. Behavior changes are markdown diffs, not code changes.",
         NAVY),
        ("TOOLS ARE IDEMPOTENT",
         "Re-running any workflow is safe. UID checkpoints, Drive file-ID tracking, OpenProject title dedup, fuzzy memory dedup — all baked in.",
         TEAL),
        ("MEMORY SURVIVES THE PROCESS",
         "The agent has no in-RAM state. Anything that persists writes to SQLite mid-run. The runner commits the DB to git after every workflow.",
         GREEN),
        ("PROMPTS ARE HUMAN CONTRACTS",
         "Anyone can change the loop in a workflow by editing markdown. PR review for prompt changes is a plain-English diff.",
         GOLD),
        ("TOOLS NEVER REASON",
         "Every algorithmic decision lives in a prompt. Python only does I/O. No tool calls another tool — composition is Claude's job.",
         PURPLE),
    ]
    y = Inches(1.7)
    h = Inches(1.0)
    gap = Inches(0.08)
    for lead, body, color in items:
        # accent block on left
        acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), y, Inches(0.20), h)
        _solid(acc, color)
        _no_line(acc)
        # white card with text
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.85), y, Inches(11.88), h)
        _solid(card, WHITE)
        _line(card, RULE, 0.5)
        card.text_frame.margin_left = Inches(0.18)
        card.text_frame.margin_top = Inches(0.10)
        card.text_frame.margin_right = Inches(0.18)
        card.text_frame.word_wrap = True
        _set_text(card.text_frame, lead, size=14, bold=True, color=color)
        _add_text(card.text_frame, body, size=12, color=INK, space_before=2)
        y += h + gap

    _notes(s,
        "These five constraints are why every architectural decision "
        "feels coherent. If a future change violates one of these, "
        "we don't make it. Useful for the audience to remember when "
        "they're tempted to ask 'can NN do X?' — answer is yes if it "
        "fits these constraints, and 'maybe but reconsider' if it "
        "doesn't.")


def slide_orchestrator(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "03 · The Python Orchestrator (agent.py)",
               "198 lines. Five responsibilities. No business logic.")

    _card(s, Inches(0.6), Inches(1.7), Inches(5.7), Inches(3.0),
          "WHAT agent.py DOES", NAVY,
          body_lines=[
              ("1.", "Load .env with override=True"),
              ("2.", "Register 31 MCP tools"),
              ("3.", "Build the system prompt (system.md + task.md)"),
              ("4.", "Open ONE Claude SDK session"),
              ("5.", "Stream text, exit"),
          ], body_size=14)

    _card(s, Inches(0.6), Inches(4.85), Inches(5.7), Inches(2.05),
          "ADDING A NEW WORKFLOW", GREEN,
          body_lines=[
              ("In agent.py:", "add task name to TASKS set"),
              ("In agent.py:", "add one-line message to TASK_KICKOFF"),
              ("Total Python change:", "two strings."),
          ], body_size=13)

    # code block on right
    code_lines = [
        "TASKS = {",
        '    "requirement-management",',
        '    "daily-report-detailed",',
        '    "daily-report-summary",',
        '    "api-functional-test",',
        '    "api-integration-test",',
        '    "api-load-test",',
        '    "e2e-test",',
        "}",
        "",
        "async def run_agent(task, dry_run):",
        "    server = create_sdk_mcp_server(",
        '        name="napco-nucleus",',
        "        version=\"0.1.0\", tools=ALL_TOOLS)",
        "    options = ClaudeAgentOptions(",
        "        system_prompt=_load_prompt(task),",
        '        mcp_servers={"napco-nucleus": server},',
        "        allowed_tools=allowed,",
        "    )",
        "    async with ClaudeSDKClient(",
        "            options=options) as client:",
        "        await client.query(",
        "            TASK_KICKOFF[task])",
        "        async for msg in client.receive_response():",
        "            print(_extract_text(msg))",
    ]
    _code_block(s, Inches(6.5), Inches(1.7), Inches(6.23), Inches(5.2),
                code_lines, font_size=10)

    _notes(s,
        "The single most surprising number in the whole talk: 198 lines. "
        "That's the entire orchestrator. People expect AI agent code to "
        "be thousands of lines. It isn't, because we offload reasoning "
        "to prompts. Walk through the code briefly — the SDK does the "
        "heavy lifting via create_sdk_mcp_server and ClaudeSDKClient. "
        "Adding a new workflow really is two string additions plus a "
        "new prompt file.")


def slide_sdk(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "03 · Claude Agent SDK Integration",
               "Three imports, the whole agent")

    _card(s, Inches(0.6), Inches(1.7), Inches(12.13), Inches(1.55),
          "THE SDK SURFACE NN USES", GOLD,
          body_lines=[
              ("ClaudeAgentOptions —", "session config (system prompt, allowed tools, mcp_servers)"),
              ("ClaudeSDKClient —", "the session context manager — opens, queries, streams"),
              ("create_sdk_mcp_server —", "wraps Python functions as MCP tools, registered with @tool"),
          ], body_size=13)

    code = [
        "from claude_agent_sdk import tool",
        "",
        "@tool(",
        '    "search_requirements",  # tool name Claude sees',
        '    "Fuzzy search requirements_seen by title. "',
        '    "Returns recent N rows with wp_url if present.",',
        '    {"query": str, "limit": int},  # arg schema',
        ")",
        "async def search_requirements_tool(args):",
        '    rows = memory.search_requirements(',
        '        args["query"], limit=args.get("limit", 5))',
        '    return {"content": [{"type": "text",',
        '            "text": json.dumps(rows, default=str)}]}',
    ]
    _code_block(s, Inches(0.6), Inches(3.45), Inches(12.13), Inches(3.4),
                code, font_size=12)

    _notes(s,
        "This is the only slide where I show real tool registration code. "
        "Highlight: the @tool decorator handles all the JSON Schema and "
        "MCP envelope plumbing. Claude sees the tool name, description, "
        "and arg schema; the Python function is plain async. No magic.")


def slide_no_api_key(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "03 · No ANTHROPIC_API_KEY",
               "Why NN runs through the local Claude Code CLI")

    _card(s, Inches(0.6), Inches(1.7), Inches(12.13), Inches(1.6),
          "THE TRANSPORT MODEL", PURPLE,
          body_lines=[
              ("SDK normally —", "talks to api.anthropic.com using ANTHROPIC_API_KEY"),
              ("NN —", "configures the SDK with cli_path=<local Claude Code CLI binary>"),
              ("Claude Code CLI handles auth —", "via the user's Claude Max subscription session, no token billing"),
          ], body_size=13)

    _card(s, Inches(0.6), Inches(3.50), Inches(6.05), Inches(3.4),
          "CONSEQUENCES — UPSIDE", GREEN,
          body_lines=[
              ("Cost is fixed —", "Claude Max subscription, ~$200/month, regardless of traffic"),
              ("No token-billing surprise —", "no runaway loops draining the budget"),
              ("No secret to rotate —", "auth lives in the CLI session, not in env vars"),
              ("Predictable forecasting —", "annual cost = $200 × 12, full stop"),
          ], body_size=13)

    _card(s, Inches(6.85), Inches(3.50), Inches(6.05), Inches(3.4),
          "CONSEQUENCES — CONSTRAINT", CORAL,
          body_lines=[
              ("Cannot run on cloud CI —", "no Claude Max session available there"),
              ("Self-hosted runner mandatory —", "the runner has the CLI logged in 24/7"),
              ("If CLI logs out — ", "every workflow fails until somebody re-authenticates"),
              ("Honest single-machine bottleneck —", "tradeoff for the cost ceiling"),
          ], body_size=13)

    _notes(s,
        "This is the slide that earns trust with the budget owner. "
        "Anyone who's looked at API-billed AI agents knows the worry: "
        "what if a prompt loops forever? Token bill blows up. NN's "
        "Claude Max model removes that worry. Be honest about the "
        "constraint: it does mean we're locked to one machine that has "
        "to stay logged in. That's a known tradeoff and we accepted it.")


def slide_mcp_overview(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "03 · MCP Tool Surface",
               "31 tools across 6 modules — the I/O layer")

    rows = [
        ["Module", "N", "Responsibility", "Representative tools"],
        ["memory", "5", "SQLite-backed cross-run continuity",
         "log_activity, recall_activity, search_requirements, recall_test_runs, memory_stats"],
        ["requirements", "4", "Email + Drive ingestion → OpenProject",
         "poll_requirement_emails, ingest_drive_files, read_requirement_inbox, publish_tasks_to_backlog"],
        ["tests", "9", "Test execution + health probes",
         "run_api_tests, run_integration_tests, run_load_tests, run_e2e_tests, run_single_e2e_test, ..."],
        ["files", "5", "Sibling-project file IO + Playwright a11y",
         "list_files, read_file, write_file, edit_file, explore_ui"],
        ["git", "3", "Read-only git context across projects",
         "git_diff, git_recent_commits, git_commit_and_push"],
        ["report", "5", "PDF + email + Teams + log tail + cleanup",
         "generate_pdf_report, send_email, send_teams_digest, tail_log, clean_reports"],
    ]
    _table(s, Inches(0.6), Inches(1.7), Inches(12.13), rows,
           col_widths=[14, 6, 32, 48], header_color=PURPLE,
           font_size=11, header_size=12)

    cap = s.shapes.add_textbox(Inches(0.6), Inches(6.6),
                                Inches(12), Inches(0.4))
    _set_text(cap.text_frame,
              "Each module exposes TOOLS (registered functions) and TOOL_NAMES (allowlist). "
              "tools/__init__.py concatenates them into ALL_TOOLS.",
              size=11, color=MUTED, align=PP_ALIGN.CENTER)

    _notes(s,
        "The 31-tool surface is broad on purpose — covers both "
        "dimensions. The split into 6 modules is by responsibility, "
        "not by feature: memory tools are reused across all workflows, "
        "report tools likewise. Notice tests is the largest module (9 "
        "tools) — that's because each test type has slightly different "
        "invocation contracts.")


def slide_tool_contract(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "03 · Tool Contract — The Invariants",
               "Five rules every MCP tool follows")

    items = [
        ("WRAP ONE EXTERNAL SYSTEM PER TOOL",
         "No tool calls another tool. Composition is Claude's job, not Python's.",
         NAVY),
        ("SIDE EFFECTS LOG TO MEMORY",
         "Every tool calls memory.log_activity inside try/except. A memory write never breaks the primary flow.",
         TEAL),
        ("HONOR NAPCO_NUCLEUS_DRY_RUN",
         "Mutating tools (publish_tasks_to_backlog, send_email, git_commit_and_push) check the env var and short-circuit.",
         CORAL),
        ("RETURN THE SDK ENVELOPE",
         'Every tool returns {"content": [{"type": "text", "text": json.dumps(...)}]}. The _text() helper enforces it.',
         GREEN),
        ("LAZY HEAVY IMPORTS",
         "imaplib, googleapiclient, requests for OpenProject, Playwright runner — all imported inside the tool fn. Cold-start stays sub-second.",
         GOLD),
    ]
    y = Inches(1.7)
    h = Inches(1.0)
    gap = Inches(0.08)
    for lead, body, color in items:
        acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), y, Inches(0.20), h)
        _solid(acc, color)
        _no_line(acc)
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.85), y, Inches(11.88), h)
        _solid(card, WHITE)
        _line(card, RULE, 0.5)
        card.text_frame.margin_left = Inches(0.18)
        card.text_frame.margin_top = Inches(0.10)
        card.text_frame.margin_right = Inches(0.18)
        card.text_frame.word_wrap = True
        _set_text(card.text_frame, lead, size=14, bold=True, color=color)
        _add_text(card.text_frame, body, size=12, color=INK, space_before=2)
        y += h + gap

    _notes(s,
        "These are non-negotiable. New tool reviews fail if they "
        "compose other tools, or skip log_activity, or ignore dry-run. "
        "The lazy-import rule is operational hygiene: in a one-turn "
        "process, cold start matters. We don't pay imaplib's cost on "
        "a workflow that doesn't poll email.")


def slide_memory(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "03 · The Memory Layer",
               "SQLite + FTS5, committed to git after every run")

    _card(s, Inches(0.6), Inches(1.7), Inches(12.13), Inches(1.4),
          "WHY SQLITE IN GIT?", NAVY,
          body_lines=[
              ("", "Cloning the project on a new machine instantly recovers prior context. Every requirement ever filed, every test run, every IMAP UID checkpoint, every Drive file ID processed. No external state store. The repo IS the database."),
          ], body_size=13)

    rows = [
        ["Table", "Purpose"],
        ["activity_logs",
         "Every meaningful action with task_name, result string, technical_details JSON. Drives the daily report."],
        ["requirements_seen",
         "One row per requirement processed. Normalized title + FTS5 fuzzy match. Powers two-layer dedup."],
        ["test_run_history",
         "One row per suite-run: pass/fail counts, duration, PDF path, regression set. Trend graphs."],
        ["email_checkpoints",
         "IMAP UIDVALIDITY + last UID. Idempotency key for poll_requirement_emails."],
        ["drive_processed",
         "Set of Drive file IDs already ingested. Idempotency key for ingest_drive_files."],
    ]
    _table(s, Inches(0.6), Inches(3.30), Inches(12.13), rows,
           col_widths=[22, 78], header_color=NAVY,
           font_size=11, header_size=12)

    _notes(s,
        "Memory in git is the most controversial design choice. "
        "Critics will say: SQLite isn't meant for that. Counter: it's "
        "small (currently ~5MB), it's binary but well-compressed, and "
        "the version control gives us audit + dedup history for free. "
        "If we needed scale we'd switch to Postgres — for now we don't.")


def slide_prompts_intro(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "04 · Prompts as Algorithms",
               "Behavior lives in markdown, not in Python")

    _card(s, Inches(0.6), Inches(1.7), Inches(12.13), Inches(2.5),
          "STRUCTURE — TWO FILES PER WORKFLOW", GREEN,
          body_lines=[
              ("system.md —", "shared backbone. Defines who Claude is, project layout, the two operational dimensions, core principles, tone rules"),
              ("<task>.md —", "the actual loop. Numbered steps, required tool calls, decision rules, output format"),
              ("Concatenated at runtime —", "Claude receives one combined system prompt, then the kickoff message"),
              ("Code review for behavior changes —", "is a markdown diff. PR-reviewable in plain English. No deploy needed."),
          ], body_size=14)

    _card(s, Inches(0.6), Inches(4.35), Inches(12.13), Inches(2.55),
          "WHY THIS WORKS", TEAL,
          body_lines=[
              ("Behavior changes are markdown PRs —", "not refactors. Adjust how tasks split, how failures classify, what goes in the executive summary, all by editing one file."),
              ("Constraints encoded in prompts —", "'memory check-in mandatory', 'title under 70 chars', 'NO em-dashes'. Things that would be brittle validation code become natural-language rules."),
              ("One ground truth per workflow —", "to know what API Functional Test does, you read api_functional_test.md. The Python tools are leaf ops only."),
          ], body_size=13)

    _notes(s,
        "Pause here. This is the deepest insight in NN's architecture. "
        "Most teams put algorithmic logic in code and treat prompts as "
        "config. NN inverts that. Prompts ARE the algorithms. Code is "
        "config — it just exposes capabilities. When you really "
        "internalize this, every architectural choice in NN follows.")


def slide_prompt_files(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "04 · The 8 Prompt Files",
               "One per workflow, plus the shared system backbone")

    rows = [
        ["File", "Workflow it drives"],
        ["system.md", "Concatenated to every task. The shared backbone."],
        ["api_functional_test.md",
         "01:00 BDT — Newman + Postman collection, classify failures"],
        ["api_integration_test.md",
         "01:15 BDT — pytest integration suite with regression diff"],
        ["api_load_test.md",
         "01:30 BDT — Locust multi-tier 10–10,000 users"],
        ["e2e_test.md",
         "02:30 BDT — Playwright full suite + failure screenshots"],
        ["requirement_management.md",
         "13:00 + 01:00 BDT — Email + Drive PDFs + meeting audio (Whisper) + Teams forwards → 3-hour task split → OpenProject Work Packages with dedup + in-place update"],
        ["daily_report_detailed.md",
         "03:30 BDT — 6-section detailed PDF emailed to the full team"],
        ["daily_report_summary.md",
         "03:45 BDT — 7-block executive dashboard emailed to leadership"],
    ]
    _table(s, Inches(0.6), Inches(1.7), Inches(12.13), rows,
           col_widths=[26, 74], header_color=GREEN,
           font_size=11, header_size=12)

    _notes(s,
        "Most prompts are 100-200 lines of markdown. system.md is "
        "longer (~300 lines) because it defines the shared rules. The "
        "real intellectual asset of NN is these 8 files. Code can be "
        "rewritten in a week; the prompts took months of iteration to "
        "get right.")


def slide_integrations(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "05 · Third-Party Integrations",
               "Everything NN talks to, end to end")

    rows = [
        ["Integration", "Library", "Used for"],
        ["IMAP", "stdlib imaplib",
         "Pull allowlisted requirement emails. Idempotent via UIDVALIDITY + since-UID checkpoint."],
        ["Google Drive", "google-api-python-client",
         "List + download new audio/video and PDFs. Service-account auth."],
        ["Groq Whisper", "requests (REST)",
         "Audio + video transcription. Cheap, fast."],
        ["pypdf", "pypdf",
         "Local PDF text extraction. No external call."],
        ["OpenProject", "requests (REST v3 HAL+JSON)",
         "list/create/update Work Packages + comments. API token, HTTP Basic auth."],
        ["Microsoft Teams", "requests",
         "Incoming webhook for the optional one-line digest."],
        ["Gmail SMTP", "stdlib smtplib",
         "Outbound email of daily report PDFs. App-password auth."],
        ["Newman", "subprocess", "API functional tests via Postman collection."],
        ["pytest", "subprocess", "API integration tests."],
        ["Locust", "subprocess", "Load testing 10 → 10,000 users in tiers."],
        ["Playwright", "subprocess", "E2E tests + failure screenshots."],
        ["TFS + MSBuild", "subprocess", "MVPAccess CICD: pull, build, deploy."],
        ["GitHub Actions", "(scheduler)", "Cron + manual triggers, 9 workflows."],
        ["Self-hosted runner", "(machine)", "Windows VM with Claude Code CLI 24/7."],
    ]
    _table(s, Inches(0.6), Inches(1.7), Inches(12.13), rows,
           col_widths=[18, 22, 60], header_color=NAVY,
           font_size=10, header_size=11)

    _notes(s,
        "The integration list is honest: NN is glue. Every meaningful "
        "external system has one tool. No tool wraps multiple "
        "systems. If we add a new client (Slack, Notion, etc.) it's "
        "one new tool, ~50 lines of Python.")


def slide_tech_stack(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "05 · Tech Stack at a Glance",
               "Everything NN depends on")

    items = [
        ("REASONING", "Claude Agent SDK · Claude Code CLI · Claude Max subscription", NAVY),
        ("RUNTIME",   "Python 3.13 · GitHub Actions · Self-hosted Windows VM", TEAL),
        ("MEMORY",    "SQLite + FTS5 · committed to git", CORAL),
        ("REPORTING", "ReportLab (PDFs) · Gmail SMTP · MS Teams webhook", GREEN),
        ("TESTS IT ORCHESTRATES", "Newman · pytest · Locust · Playwright", GOLD),
        ("EXTERNAL APIs", "IMAP · Google Drive · Groq Whisper · OpenProject · TFS · MSBuild · IIS", PURPLE),
        ("PYTHON DEPS", "claude-agent-sdk · python-dotenv · requests · google-api-python-client · pypdf", NAVY),
    ]
    y = Inches(1.7)
    h = Inches(0.72)
    gap = Inches(0.08)
    for lead, body, color in items:
        acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), y, Inches(2.4), h)
        _solid(acc, color)
        _no_line(acc)
        acc.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(acc.text_frame, lead, size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        body_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(3.05), y, Inches(9.68), h)
        _solid(body_box, WHITE)
        _line(body_box, RULE, 0.5)
        body_box.text_frame.margin_left = Inches(0.18)
        body_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        body_box.text_frame.word_wrap = True
        _set_text(body_box.text_frame, body, size=12, color=INK)
        y += h + gap

    _notes(s,
        "Use this as the wrap-up for the inside-NN section. Nothing "
        "exotic — every component is well-understood, well-supported. "
        "Claude Agent SDK is the only piece newer than 1 year. The "
        "rest is boring tech, deliberately.")


def slide_cicd(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "06 · CI/CD Execution Model",
               "Every workflow follows the same five-step YAML")

    code = [
        "name: <Display Name>",
        "on:",
        "  schedule:",
        "    - cron: '<utc-cron>'",
        "  workflow_dispatch:",
        "    inputs:",
        "      dry_run: { type: choice, options: ['false','true'] }",
        "",
        "permissions: { contents: write }",
        "concurrency: { group: <task>, cancel-in-progress: false }",
        "",
        "jobs:",
        "  run:",
        "    runs-on: [self-hosted, Windows]",
        "    timeout-minutes: 25",
        "    env: { ... only the secrets this task needs ... }",
        "    steps:",
        "      - uses: actions/checkout@v5",
        "      - run: py -3 -m pip install -r requirements.txt",
        "      - run: py -3 agent.py --task <task-name>",
        "      - name: Commit state changes",
        "        if: always()",
        "        run: |",
        "          git add nucleus_memory.db data/...",
        "          git commit -m \"<task>: $(Get-Date -Format ...)\"",
        "          git pull --rebase origin main && git push",
    ]
    _code_block(s, Inches(0.6), Inches(1.7), Inches(7.5), Inches(5.2),
                code, font_size=10)

    _card(s, Inches(8.30), Inches(1.7), Inches(4.43), Inches(5.2),
          "OPERATIONAL GUARDRAILS", PURPLE,
          body_lines=[
              ("concurrency groups —", "runs queue, never overlap"),
              ("timeout-minutes —", "25 min cap, longer for load test"),
              ("contents: write —", "needed because we commit memory"),
              ("self-hosted only —", "Cloud runners cannot run NN"),
              ("dry_run input —", "manual testing without side effects"),
          ], body_size=12)

    _notes(s,
        "The YAML is intentionally boring. Anyone who knows GitHub "
        "Actions can read a NN workflow without learning anything new. "
        "That's a feature. Concurrency groups are the most important "
        "guardrail — they're how we guarantee memory commits never "
        "race.")


def slide_workflows(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "06 · The 9 Production Workflows",
               "All 9 currently scheduled, running daily")

    rows = [
        ["#", "Workflow", "Schedule (BDT)", "What it does"],
        ["1", "API Functional Test", "01:00",
         "Newman + Postman collection across the API surface"],
        ["2", "API Integration Test", "01:15",
         "pytest integration suite with regression diff vs. prior runs"],
        ["3", "API Load Test", "01:30",
         "Locust multi-tier 10 → 10,000 users with cooldowns"],
        ["4", "MVP Access E2E Test", "02:30",
         "Playwright full suite. Failure screenshots embedded in PDF"],
        ["5", "MVPAccess CICD", "03:00",
         "TFS pull, MSBuild Release, IIS deploy via UNC, health check"],
        ["6", "Requirement Management", "13:00 + 01:00",
         "IMAP. 3-hour task split. OpenProject Work Packages, dedup + in-place update"],
        ["7", "Daily Report (Detailed)", "03:30",
         "Reads memory + 4 test PDFs. 6-section PDF to full team"],
        ["8", "Daily Report (Summary)", "03:45",
         "7-block executive dashboard to leadership"],
        ["9", "Probe Runner Filesystem", "Manual",
         "Diagnostic — inspects runner state during triage"],
    ]
    _table(s, Inches(0.6), Inches(1.7), Inches(12.13), rows,
           col_widths=[5, 24, 22, 49], header_color=NAVY,
           font_size=10, header_size=11)

    _notes(s,
        "Walk through the timeline: everything fires inside a single "
        "3-hour window, 01:00–04:00 BDT. The four test workflows go "
        "first (01:00 functional, 01:15 integration, 01:30 load, 02:30 "
        "E2E) — they queue on the single self-hosted runner. CICD "
        "follows at 03:00 once tests have passed. Requirement "
        "Management runs twice daily at 13:00 + 01:00 BDT to ingest anything that arrived "
        "overnight. Daily reports close the window at 03:30 (detailed) "
        "and 03:45 (summary), so by 4 AM every consumer has fresh "
        "state in their inbox before the workday starts.")


def slide_runner(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "06 · The Self-Hosted Runner",
               "One Windows VM, 24/7, the whole platform")

    _card(s, Inches(0.6), Inches(1.7), Inches(6.05), Inches(5.2),
          "WHAT'S ON THE RUNNER", NAVY,
          body_lines=[
              ("Windows Server VM", "172.16.205.209"),
              ("Python 3.13", "+ all NN deps installed"),
              ("Claude Code CLI", "logged into Claude Max 24/7"),
              ("Newman / pytest / Locust", "test runners pre-installed"),
              ("Playwright + Chromium", "for E2E"),
              ("MSBuild + TFS client", "for CICD workflow"),
              ("GitHub Actions runner", "registered as 'test-runner'"),
              ("UNC network access", "for IIS deploy"),
          ], body_size=13)

    _card(s, Inches(6.85), Inches(1.7), Inches(6.05), Inches(5.2),
          "WHY ONE MACHINE?", PURPLE,
          body_lines=[
              ("Single Claude Max session —", "can only live in one place"),
              ("Windows-specific test runners —", "MSBuild, TFS, Newman all native"),
              ("Network proximity to TFS / IIS —", "for CICD"),
              ("Cost — ", "one VM is cheaper than scaling out"),
              ("", ""),
              ("Honest tradeoff — ", "single machine = single point of failure. We accept this for the cost ceiling."),
          ], body_size=13)

    _notes(s,
        "The runner is the single most fragile component in NN. If it "
        "goes down, every workflow stops. We accept this because: (a) "
        "the cost saving is meaningful, (b) our usage doesn't justify "
        "redundancy, (c) we can spin up a cold-standby runner in an "
        "afternoon if we need to. Roadmap: cold-standby pre-imaged.")


def slide_add_workflow(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "07 · Adding a New Workflow",
               "Seven-step recipe — usually under an hour")

    steps = [
        ("1", "WRITE THE PROMPT",
         "Create prompts/<new_task>.md. Define the loop: memory check-in (mandatory), the steps, the dedup rule, the output contract.",
         NAVY),
        ("2", "WIRE agent.py",
         "Add task name to TASKS set. Add one-line kickoff message to TASK_KICKOFF dict. No other Python changes if existing tools cover the I/O.",
         TEAL),
        ("3", "ADD A TOOL IF NEEDED",
         "If the task needs new I/O — add a tool (see slide 25). Otherwise skip.",
         CORAL),
        ("4", "AUTHOR THE WORKFLOW YAML",
         "Copy an existing file in .github/workflows/. Update name, cron, env: secrets, --task flag, post-run commit pattern.",
         GREEN),
        ("5", "ADD SECRETS",
         "If the task needs new env vars, add them to GitHub Actions secrets AND to the runner's local .env.",
         GOLD),
        ("6", "DRY-RUN FIRST",
         "Trigger via workflow_dispatch with dry_run=true. Verify the loop logs to memory but does not mutate. Inspect activity_logs.",
         PURPLE),
        ("7", "GO LIVE + MONITOR",
         "Switch to live. Tail logs from the runner. Confirm the post-run commit landed on main with the new state.",
         NAVY),
    ]
    y = Inches(1.7)
    h = Inches(0.72)
    gap = Inches(0.06)
    for num, lead, body, color in steps:
        nb = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.6), y, Inches(0.72), h)
        _solid(nb, color)
        _no_line(nb)
        nb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text(nb.text_frame, num, size=20, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.40), y, Inches(11.33), h)
        _solid(card, WHITE)
        _line(card, RULE, 0.5)
        card.text_frame.margin_left = Inches(0.18)
        card.text_frame.margin_top = Inches(0.06)
        card.text_frame.word_wrap = True
        _set_text(card.text_frame, lead, size=12, bold=True, color=color)
        _add_text(card.text_frame, body, size=10.5, color=INK,
                  space_before=2)
        y += h + gap

    _notes(s,
        "The honest pacing: a brand-new workflow with new tools takes "
        "1-2 days end-to-end. A workflow that uses existing tools is "
        "done in 30 minutes. Most of the time goes into prompt iteration "
        "in step 1 — getting the loop right. Steps 2-7 are mechanical.")


def slide_add_tool(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "07 · Adding a New MCP Tool",
               "Six-step recipe — usually under 30 minutes")

    _card(s, Inches(0.6), Inches(1.7), Inches(5.85), Inches(5.2),
          "STEPS", PURPLE,
          body_lines=[
              ("1.", "Pick the right module (memory / requirements / tests / files / git / report)"),
              ("2.", "Write the function — @tool decorator, lazy imports, try/except, log_activity"),
              ("3.", "Return the SDK envelope via _text() helper from tools/_shared.py"),
              ("4.", "Append to TOOLS + TOOL_NAMES at the bottom of the module"),
              ("5.", "Update the prompt that calls the tool — Claude needs to know it exists"),
              ("6.", "Smoke-test locally with --dry-run, watch the streamed text"),
          ], body_size=13)

    code = [
        "from claude_agent_sdk import tool",
        "import memory",
        "from tools._shared import _text",
        "",
        "@tool(",
        '    \"do_thing\",',
        '    \"What it does, one sentence Claude can act on.\",',
        '    {\"arg1\": str, \"arg2\": int},',
        ")",
        "async def do_thing_tool(args):",
        "    if os.environ.get(\"NAPCO_NUCLEUS_DRY_RUN\") == \"1\":",
        "        return _text({\"dry_run\": True})",
        "    try:",
        "        result = backend.run(",
        "            args[\"arg1\"], args[\"arg2\"])",
        "    except Exception as e:",
        "        memory.log_activity(",
        "            task_name=\"do_thing\",",
        "            result=f\"error:{type(e).__name__}\")",
        "        return _text({\"error\": str(e)})",
        "    memory.log_activity(",
        "        task_name=\"do_thing\",",
        "        result=f\"ok:{result.summary}\")",
        "    return _text(result.to_dict())",
        "",
        "TOOLS = [do_thing_tool]",
        "TOOL_NAMES = [\"do_thing\"]",
    ]
    _code_block(s, Inches(6.65), Inches(1.7), Inches(6.08), Inches(5.2),
                code, font_size=10)

    _notes(s,
        "Show the template if anyone asks. The shape is the same for "
        "every tool — that's deliberate. Code review for new tools is "
        "fast because everything looks the same.")


def slide_benefits(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "08 · Benefits",
               "What this architecture buys")

    items = [
        ("Behavior-via-prompt", "Adjust split rules, classification, summary content — all in markdown PRs", NAVY),
        ("Cost-bounded", "Claude Max subscription is fixed monthly. No per-token surprises.", TEAL),
        ("Single-machine deploy", "One VM, one repo, one .env. No K8s, no message queue.", CORAL),
        ("Memory in git", "Audit trail, dedup history, run history all version-controlled.", GREEN),
        ("Idempotent by construction", "Every ingest path has a checkpoint. Every publish has dedup.", GOLD),
        ("Tool surface stays small", "31 tools across 6 modules — reasoning lives in markdown.", PURPLE),
        ("One Claude turn per process", "No session leak, no context bleed. Process = state boundary.", NAVY),
        ("Standard CI/CD", "Vanilla GitHub Actions YAML. No new orchestrator to learn.", TEAL),
        ("Multi-language input", "Bangla, Malay, mixed-script — Claude reads it. Output always English.", CORAL),
        ("Two consolidated emails per day", "09:00 detail, 09:30 executive. Replaces 6+ overnight fragments.", GREEN),
    ]

    cols = 2
    rows = 5
    x0 = Inches(0.6)
    y0 = Inches(1.7)
    cw = Inches(6.05)
    ch = Inches(1.00)
    gx = Inches(0.10)
    gy = Inches(0.08)
    for i, (lead, body, color) in enumerate(items):
        r = i // cols
        c = i % cols
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(0.20), ch)
        _solid(acc, color)
        _no_line(acc)
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + Inches(0.25), y,
                                   cw - Inches(0.25), ch)
        _solid(card, WHITE)
        _line(card, RULE, 0.5)
        card.text_frame.margin_left = Inches(0.15)
        card.text_frame.margin_top = Inches(0.10)
        card.text_frame.margin_right = Inches(0.15)
        card.text_frame.word_wrap = True
        _set_text(card.text_frame, lead, size=12, bold=True, color=color)
        _add_text(card.text_frame, body, size=10, color=INK, space_before=2)

    _notes(s,
        "Spend 2 minutes here. Pick the 3 most relevant for your "
        "audience: for engineering — behavior-via-prompt, idempotent, "
        "single-machine. For leadership — cost-bounded, two emails per "
        "day. Don't read all ten.")


def slide_demerits(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "08 · Honest Limitations",
               "What this architecture costs")

    items = [
        ("Single-machine bottleneck", "If the runner goes down, every workflow stops. No automatic failover.", CORAL),
        ("Claude Max session is hidden dep", "If the CLI logs out, every workflow fails until somebody re-authenticates.", CORAL),
        ("Non-determinism", "Same input, slightly different output. Acceptable for reports, problematic for invariants.", GOLD),
        ("Prompt drift", "A small markdown edit can change behavior subtly across consumers. No automated regression suite.", GOLD),
        ("Memory growth", "nucleus_memory.db grows monotonically. No retention policy yet.", PURPLE),
        ("Cross-language coupling", "Tools in Python, tests in TypeScript/JS. Subprocess boundary = log-driven debugging.", PURPLE),
        ("Secrets fan-out", "Every secret in two places: GitHub Actions secrets + runner .env. Manual rotation.", NAVY),
        ("Limited concurrency", "One job at a time per workflow group. Single runner can't parallelize.", NAVY),
        ("No multi-tenant story", "Allowlists, paths, OpenProject project ID hard-coded for MVP Access.", TEAL),
        ("Hard-to-test reasoning", "Reasoning lives in prompts. Validating a prompt change is observational.", TEAL),
    ]

    cols = 2
    rows = 5
    x0 = Inches(0.6)
    y0 = Inches(1.7)
    cw = Inches(6.05)
    ch = Inches(1.00)
    gx = Inches(0.10)
    gy = Inches(0.08)
    for i, (lead, body, color) in enumerate(items):
        r = i // cols
        c = i % cols
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(0.20), ch)
        _solid(acc, color)
        _no_line(acc)
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + Inches(0.25), y,
                                   cw - Inches(0.25), ch)
        _solid(card, WHITE)
        _line(card, RULE, 0.5)
        card.text_frame.margin_left = Inches(0.15)
        card.text_frame.margin_top = Inches(0.10)
        card.text_frame.margin_right = Inches(0.15)
        card.text_frame.word_wrap = True
        _set_text(card.text_frame, lead, size=12, bold=True, color=color)
        _add_text(card.text_frame, body, size=10, color=INK, space_before=2)

    _notes(s,
        "This slide builds credibility. The audience knows nothing is "
        "perfect. Listing the limits openly tells them I've thought "
        "about NN clearly. Be especially honest about non-determinism "
        "and prompt drift — those are real and we manage them through "
        "review discipline, not tooling.")


def slide_roadmap(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "08 · Roadmap",
               "What's coming next, in order of value")

    items = [
        ("NEAR TERM (Q2 2026)",
         [
             "Cold-standby pre-imaged runner — disaster recovery in <1h",
             "Memory retention policy + nightly vacuum",
             "Prompt regression harness — golden-output diffs",
             "Slack tool — extend Teams digest pattern",
         ], NAVY),
        ("MID TERM (Q3 2026)",
         [
             "Second-product onboarding — externalize project config",
             "Web dashboard — live workflow status, memory browser",
             "Prompt linter — catch banned patterns pre-merge",
             "Per-workflow latency budgets + alerts",
         ], TEAL),
        ("LONG TERM (Q4 2026+)",
         [
             "Multi-runner failover — geographic redundancy",
             "Cross-product memory — share dedup across products",
             "Voice intake — meeting recordings → requirements directly",
             "Agentic refactor for tasks where multi-agent helps",
         ], CORAL),
    ]
    y0 = Inches(1.7)
    cw = Inches(4.05)
    ch = Inches(5.2)
    gx = Inches(0.10)
    x0 = Inches(0.6)
    for i, (head, lst, color) in enumerate(items):
        x = x0 + i * (cw + gx)
        body_lines = [("•", line) for line in lst]
        _card(s, x, y0, cw, ch, head, color,
              body_lines=body_lines, body_size=12)

    _notes(s,
        "The roadmap is honest about prioritization: near-term items "
        "address the limits we just listed. Mid-term is value expansion. "
        "Long-term is speculation — don't promise dates we can't keep. "
        "Notice: 'agentic refactor' on long-term means we'd revisit the "
        "single-agent decision IF a future workflow truly needs it. We "
        "haven't found one yet.")


def slide_closing(prs, total, page):
    s = _new_slide(prs)
    _add_chrome(s, page, total, FOOTER)
    _add_title(s, "Closing",
               "The thesis, restated — then your questions")

    # big quote
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.7),
                              Inches(12.13), Inches(2.4))
    _solid(box, NAVY)
    _no_line(box)
    box.text_frame.margin_left = Inches(0.5)
    box.text_frame.margin_right = Inches(0.5)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.word_wrap = True
    _set_text(box.text_frame,
              "QA architect plus AI equals senior developer team output.",
              size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(box.text_frame,
              "9 production workflows. 31 tools. ~3.6k Python LOC. Two consolidated emails per day. "
              "Every client requirement filed in OpenProject within two hours of being spoken on a call. "
              "Built and operated by one person — because the reasoning lives in prompts, "
              "not in code that needs a team to maintain.",
              size=14, color=GOLD, align=PP_ALIGN.CENTER, space_before=10)

    # Q&A footer
    qa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(4.6),
                             Inches(12.13), Inches(2.3))
    _solid(qa, SOFT)
    _no_line(qa)
    qa.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    qa.text_frame.margin_left = Inches(0.5)
    qa.text_frame.margin_right = Inches(0.5)
    _set_text(qa.text_frame, "Q & A",
              size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _add_text(qa.text_frame, "Mohammad Kamrul Hasan",
              size=14, color=INK, align=PP_ALIGN.CENTER, space_before=8)
    _add_text(qa.text_frame, "AI-Augmented QA Architect  ·  Adaptive Enterprise Limited",
              size=12, color=MUTED, align=PP_ALIGN.CENTER, space_before=2)
    _add_text(qa.text_frame, "khasan@ael-bd.com",
              size=11, color=TEAL, align=PP_ALIGN.CENTER, space_before=2)

    _notes(s,
        "Land the talk on the thesis line. Pause. Then open Q&A. "
        "Predicted questions: (1) what about token costs / API "
        "billing — covered by Claude Max slide. (2) why not "
        "LangGraph — covered by single-agent slide, repeat the answer. "
        "(3) how do we handle Claude getting it wrong — answer: "
        "human-in-loop is preserved, NN files draft OpenProject Work Packages, "
        "humans approve/edit; reports are read by humans before action. "
        "(4) can this be commercialized — yes, the architecture is "
        "product-agnostic, multi-tenant is on the roadmap.")


# ── build ──────────────────────────────────────────────────────────

def build():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_agenda,
        slide_problem,
        slide_thesis,
        slide_two_dimensions,
        slide_team_streams,
        slide_requirement_pipeline,
        slide_headline_numbers,
        slide_arch_overview,
        slide_single_agent,
        slide_lifecycle,
        slide_design_constraints,
        slide_orchestrator,
        slide_sdk,
        slide_no_api_key,
        slide_mcp_overview,
        slide_tool_contract,
        slide_memory,
        slide_prompts_intro,
        slide_prompt_files,
        slide_integrations,
        slide_tech_stack,
        slide_cicd,
        slide_workflows,
        slide_runner,
        slide_add_workflow,
        slide_add_tool,
        slide_benefits,
        slide_demerits,
        slide_roadmap,
        slide_closing,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        if fn is slide_title:
            fn(prs, total)
        else:
            fn(prs, total, i)

    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}")
