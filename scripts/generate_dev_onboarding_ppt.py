"""NAPCO Nucleus - developer onboarding decks (PRODUCTION reality).

Two TECHNICAL decks for an engineer joining the project or replicating
its pattern. These describe ONLY what runs in production today - the
central-host docker stack on 172.16.205.123 - verified against the
code, not the (stale) README / flow diagram.

Ground-truth as of this rewrite:
  - Orchestrator:  collect_central.py --client all --last-minutes 1440
  - Fires:         once a day, 23:30 BD, nucleus-daily-draft container
  - Transcription: Google STT (Chirp 2) only. No faster-whisper, no Groq.
  - Identify:      agent.py --task verify_session, Claude Agent SDK, Opus 4.7
  - Tracker:       OpenProject Work Packages (GitLab retired 2026-04-28)
  - Delivery:      ONE rollup email/day (mail.daily_rollup) + Verification .docx

Self-contained: helpers mirror generate_system_overview_ppt.py so the
decks look like siblings.

Run:
    py -3 scripts\\generate_dev_onboarding_ppt.py
Output:
    docs\\NAPCO-Nucleus-Architecture-Onboarding.pptx
    docs\\NAPCO-Nucleus-Workflow-Onboarding.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT_ARCH = ROOT / "docs" / "NAPCO-Nucleus-Architecture-Onboarding.pptx"
OUT_FLOW = ROOT / "docs" / "NAPCO-Nucleus-Workflow-Onboarding.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette - match docs/NAPCO-Nucleus-Presentation.pptx
NAVY   = RGBColor(0x1F, 0x4E, 0x79)
TEAL   = RGBColor(0x2E, 0x8A, 0x8A)
CORAL  = RGBColor(0xE0, 0x78, 0x56)
GREEN  = RGBColor(0x4A, 0x7A, 0x4A)
GOLD   = RGBColor(0xC9, 0x96, 0x2B)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)
INK    = RGBColor(0x22, 0x22, 0x22)
MUTED  = RGBColor(0x6B, 0x77, 0x85)
SOFT   = RGBColor(0xF5, 0xF7, 0xFA)
RULE   = RGBColor(0xD5, 0xDC, 0xE5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MONO   = "Consolas"


# -- helpers --------------------------------------------------------

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def no_line(shape):
    shape.line.fill.background()


def line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_text(tf, text, size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_para(tf, text, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, space_before=0, font="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    return tf


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None,
             line_width=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color is None:
        no_fill(rect)
    else:
        solid(rect, fill_color)
    if line_color is None:
        no_line(rect)
    else:
        line(rect, line_color, line_width)
    rect.text_frame.margin_left = Inches(0.12)
    rect.text_frame.margin_right = Inches(0.12)
    rect.text_frame.margin_top = Inches(0.08)
    rect.text_frame.margin_bottom = Inches(0.08)
    return rect


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width_pt=2.0):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn


def base_slide(prs, title_text, subtitle_text=None,
               title_size=34, subtitle_size=15):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill_color=NAVY)
    tf = add_textbox(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    set_text(tf, title_text, size=title_size, bold=True, color=NAVY)
    if subtitle_text:
        add_para(tf, subtitle_text, size=subtitle_size, color=MUTED,
                 space_before=4)
    add_rect(s, Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.02),
             fill_color=RULE)
    return s


def set_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    p.text = text


def add_chip(slide, x, y, w, h, text, color, font_size=14, font="Calibri"):
    rect = add_rect(slide, x, y, w, h, fill_color=color, rounded=True)
    tf = rect.text_frame
    set_text(tf, text, size=font_size, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font=font)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return rect


def add_box(slide, x, y, w, h, title, lines, accent=NAVY,
            title_size=16, body_size=13, body_font="Calibri"):
    container = add_rect(slide, x, y, w, h, fill_color=WHITE, line_color=RULE)
    add_rect(slide, x, y, Inches(0.08), h, fill_color=accent)
    tf = add_textbox(slide, x + Inches(0.22), y + Inches(0.12),
                     w - Inches(0.32), h - Inches(0.24))
    set_text(tf, title, size=title_size, bold=True, color=accent)
    for ln in lines:
        add_para(tf, ln, size=body_size, color=INK, space_before=3,
                 font=body_font)
    return container


def title_slide(prs, big, sub1, sub2, tag):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=NAVY)
    tf = add_textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6))
    set_text(tf, big, size=54, bold=True, color=WHITE)
    tf2 = add_textbox(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(1.4))
    set_text(tf2, sub1, size=24, color=RGBColor(0xC8, 0xD2, 0xE0))
    if sub2:
        add_para(tf2, sub2, size=24, color=RGBColor(0xC8, 0xD2, 0xE0))
    tf3 = add_textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
    set_text(tf3, tag, size=12, color=RGBColor(0x9C, 0xAB, 0xBE))
    return s


def closing_slide(prs, big, sub, pointers):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=NAVY)
    tf = add_textbox(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.2))
    set_text(tf, big, size=44, bold=True, color=WHITE)
    tf2 = add_textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(0.7))
    set_text(tf2, sub, size=18, color=RGBColor(0xC8, 0xD2, 0xE0))
    y = 3.5
    for label, path in pointers:
        tf_l = add_textbox(s, Inches(0.9), Inches(y), Inches(3.2), Inches(0.5))
        set_text(tf_l, label, size=14, bold=True, color=WHITE)
        tf_p = add_textbox(s, Inches(4.2), Inches(y), Inches(8.2), Inches(0.5))
        set_text(tf_p, path, size=14, color=RGBColor(0xC8, 0xD2, 0xE0),
                 font=MONO)
        y += 0.55
    return s


# ===================================================================
# DECK 1 - ARCHITECTURE (production)
# ===================================================================

def arch_mental_model(prs):
    s = base_slide(prs, "The one idea to grok first",
                   "The business logic is NOT in the Python. It's in the prompt.")
    cols = [
        ("Prompt = the brain", NAVY,
         ["prompts/*.md",
          "Decides what is a real",
          "requirement, how to split,",
          "what to dedupe, when to stop.",
          "Edit here to change behavior."]),
        ("Tools = the hands", TEAL,
         ["tools/*.py  (MCP server)",
          "Deterministic wrappers:",
          "Google STT, IMAP, Drive,",
          "OpenProject, SQLite. No",
          "judgement - just do-the-thing."]),
        ("Agent = the loop", CORAL,
         ["agent.py --task verify_session",
          "Loads system + task prompt,",
          "registers the MCP server,",
          "runs ONE Claude turn (Opus 4.7),",
          "exits. That's it."]),
    ]
    cw = Inches(3.9); ch = Inches(3.3); gap = 0.25
    total = 3 * 3.9 + 2 * gap
    x0 = (13.333 - total) / 2
    for i, (t, c, body) in enumerate(cols):
        x = Inches(x0 + i * (3.9 + gap))
        add_box(s, x, Inches(2.0), cw, ch, t, body, accent=c,
                title_size=20, body_size=14)
    tf = add_textbox(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.2))
    set_text(tf, "Claude is the decider. The tools are dumb on purpose.",
             size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_para(tf,
             "Coming from a normal codebase, you'll hunt for logic in "
             "tools/ and not find it. Read the prompt first.",
             size=13, color=MUTED, align=PP_ALIGN.CENTER, space_before=8)
    set_speaker_notes(s, (
        "The single biggest mental-model jump for someone from a "
        "conventional codebase. Nucleus is built on the Claude Agent SDK, "
        "authenticated via the local Claude Code CLI - no ANTHROPIC_API_KEY. "
        "The identify step is agent.py --task verify_session: it loads "
        "prompts/system.md plus the task prompt, registers the in-process "
        "'napco-nucleus' MCP server (the tools/ submodules), runs a single "
        "Claude turn on Opus 4.7 (pinned via NUCLEUS_AGENT_MODEL="
        "claude-opus-4-7, commit 2017d01), and exits. All judgement - what "
        "counts as a requirement, how to split, what to dedupe - lives in "
        "the markdown prompt, not in Python. To change behavior you edit a "
        "prompt and re-run evals/run.py. (A 3-stage pipeline.py variant - "
        "Haiku extract / Sonnet critique / Opus draft - exists but is NOT "
        "the daily-production default.)"
    ))


def arch_topology(prs):
    s = base_slide(prs, "Production topology",
                   "Dev PCs capture and push. One Linux host does everything else.")
    # Dev PCs
    add_box(s, Inches(0.6), Inches(1.7), Inches(3.6), Inches(2.9),
            "Each developer PC (Windows)", [
                "voice daemon -> records the",
                "Teams call (mic + speaker WAV)",
                "push_chat.py -> Teams chat .docx",
                "",
                "Writes to the Samba share only.",
                "No secrets except the share path.",
            ], accent=PURPLE, title_size=15, body_size=12)
    # arrow to share
    add_arrow(s, Inches(4.2), Inches(3.1), Inches(4.9), Inches(3.1),
              color=NAVY, width_pt=3)
    add_chip(s, Inches(4.45), Inches(4.7), Inches(3.4), Inches(0.5),
             "\\\\172.16.205.123\\nucleus-central  (SMB, guest)", MUTED,
             font_size=11, font=MONO)
    # Central host
    add_rect(s, Inches(5.1), Inches(1.55), Inches(7.6), Inches(4.9),
             fill_color=SOFT, line_color=RULE)
    tf = add_textbox(s, Inches(5.3), Inches(1.68), Inches(7.2), Inches(0.6))
    set_text(tf, "Central host  -  172.16.205.123  (Ubuntu, docker-compose)",
             size=14, bold=True, color=NAVY)
    containers = [
        ("nucleus-samba", TEAL, "serves the central share (SMB 445)"),
        ("nucleus-transcribe", PURPLE, "2-min loop: call WAV -> Google STT"),
        ("nucleus-stage-email", CORAL, "15-min: Gmail IMAP -> central"),
        ("nucleus-stage-drive", GOLD, "15-min: Google Drive -> central"),
        ("nucleus-daily-draft", NAVY, "23:30 BD: full pipeline + 1 email"),
        ("nucleus-gha-runner", GREEN, "optional: manual GHA dispatch"),
    ]
    cw = Inches(3.5); chh = Inches(1.15); gx = 0.2; gy = 0.12
    x0 = 5.35; y0 = 2.4
    for i, (name, color, desc) in enumerate(containers):
        row, col = divmod(i, 2)
        x = Inches(x0 + col * (3.5 + gx))
        y = Inches(y0 + row * (1.15 + gy))
        add_box(s, x, y, cw, chh, name, [desc], accent=color,
                title_size=13, body_size=11)
    set_speaker_notes(s, (
        "Production topology since the 2026-05-14 migration. Each "
        "developer's Windows PC runs capture only: a voice daemon (watches "
        "the Teams audio session via pycaw, records mic + speaker as two "
        "separate WAV tracks the instant a call starts) and scheduled chat "
        "pushes. Dev machines hold no secrets beyond the central share "
        "path; they write to the SMB share (guest / no-credential since "
        "commit ff99691), never directly to state. Everything else is one "
        "Ubuntu box running a six-container docker-compose stack "
        "(deploy/linux-central/docker-compose.yml). The repo is "
        "bind-mounted so a git pull on the host hot-reloads workers. An "
        "off-network dev bridges via Google Drive instead of SMB "
        "(drive/offnet_sync.py). The .209 / .195 boxes are retired from "
        "Nucleus duty."
    ))


def arch_containers(prs):
    s = base_slide(prs, "The six containers",
                   "Four capture continuously. One thinks once a day. One is on-demand.")
    rows = [
        ("nucleus-samba", TEAL, "Serves \\\\172.16.205.123\\nucleus-central over SMB. The drop zone every dev writes to."),
        ("nucleus-transcribe", PURPLE, "Every 2 min: finds completed calls on the share, transcribes via Google STT (Chirp 2)."),
        ("nucleus-stage-email", CORAL, "Every 15 min: pulls roster-filtered Gmail (mail.pull_email) into the central tree."),
        ("nucleus-stage-drive", GOLD, "Every 15 min: pulls shared Google Drive files (drive.pull_drive) into the central tree."),
        ("nucleus-daily-draft", NAVY, "Once a day, 23:30 BD: runs collect_central.py -> identify -> OpenProject -> ONE email. 4 GB cap."),
        ("nucleus-gha-runner", GREEN, "Optional (--profile runner): self-hosted runner for manual workflow_dispatch."),
    ]
    y0 = 1.8
    for i, (name, color, desc) in enumerate(rows):
        y = Inches(y0 + i * 0.82)
        add_chip(s, Inches(0.7), y, Inches(3.1), Inches(0.62), name, color,
                 font_size=13, font=MONO)
        tf = add_textbox(s, Inches(4.0), y, Inches(8.6), Inches(0.62),
                         anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, desc, size=13, color=INK)
    set_speaker_notes(s, (
        "The stack from docker-compose.yml. The split that matters: the "
        "first four containers run continuously and only CAPTURE + STAGE - "
        "they move bytes into the central tree (calls, transcripts, email, "
        "Drive). They never call the LLM. nucleus-daily-draft is the only "
        "one that thinks: it fires once a day at 23:30 BD "
        "(DAILY_DRAFT_TARGET_TIME) via daily-draft-loop.sh, runs "
        "collect_central.py --client all --last-minutes 1440, which "
        "aggregates the day, then invokes agent.py --task verify_session "
        "to identify requirements, publishes to OpenProject, and sends one "
        "rollup email. Memory-capped at 4 GB to fence runaway LLM calls. "
        "nucleus-gha-runner is optional and only needed for manual "
        "GitHub Actions dispatch. Note: the per-call event trigger still "
        "gets written by the transcribe loop but daily-draft deliberately "
        "ignores it - one email per day policy (commit 779034d)."
    ))


def arch_memory(prs):
    s = base_slide(prs, "Memory & idempotency",
                   "SQLite + FTS5. Dedup means a re-run never doubles a requirement.")
    add_box(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(3.2),
            "nucleus_memory.db  (SQLite + FTS5)", [
                "requirements_seen   - fuzzy dedup + WP URLs",
                "activity_logs       - every tool action",
                "email_checkpoints   - IMAP UIDVALIDITY + since-UID",
                "drive_processed     - file IDs never re-ingested",
                "requirement_reviews - client accept/reject deltas",
            ], accent=NAVY, title_size=16, body_size=12.5, body_font=MONO)
    add_box(s, Inches(6.85), Inches(1.8), Inches(5.9), Inches(3.2),
            "How state survives", [
                "On the central host the DB lives in a",
                "named docker volume (survives compose down).",
                "",
                "Three dedupe layers stop duplicates:",
                "  1. IMAP since-UID checkpoint",
                "  2. open Work-Package title match",
                "  3. fuzzy match vs requirements_seen",
                "",
                "Confirmed items don't resurface tomorrow.",
            ], accent=GREEN, title_size=16, body_size=12.5)
    tf = add_textbox(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.2))
    set_text(tf,
             "remember_requirement() records each new Work-Package IID + URL "
             "so the next run skips it.",
             size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_speaker_notes(s, (
        "memory.py is a SQLite database with FTS5 full-text search and is "
        "the backbone of idempotency. On the central host it's a named "
        "docker volume so it survives compose down. The dedupe guarantee "
        "has three layers: IMAP uses UIDVALIDITY plus a since-UID "
        "checkpoint so the same email is never re-read; Drive records every "
        "processed file ID; and at publish time "
        "(tools/requirements.py:publish_tasks_to_backlog) each candidate is "
        "matched both against currently-open OpenProject Work-Package "
        "titles and fuzzily against requirements_seen. "
        "remember_requirement() stores the new IID + URL. "
        "requirement_reviews holds the line-by-line client reply deltas so "
        "confirmed items don't get re-asked."
    ))


def arch_external(prs):
    s = base_slide(prs, "External services & the auth model",
                   "No ANTHROPIC_API_KEY. Claude runs through the local CLI.")
    rows = [
        ("Claude Agent SDK", NAVY, "Identify step. Local Claude Code CLI login (Max tier). Opus 4.7. No API key, $0 marginal."),
        ("Google STT", TEAL, "The SOLE call-transcription engine (tools/google_stt.py). Chirp 2 model, best Bangla."),
        ("IMAP / Gmail", CORAL, "Email ingest (mail.pull_email). Roster-filtered allowlist of senders only."),
        ("Google Drive API", GOLD, "Shared-folder watcher (drive.pull_drive). Service-account JSON at repo root."),
        ("OpenProject", PURPLE, "Tracker. Work Packages in project mvp-access (replaced GitLab 2026-04-28)."),
        ("SMTP (Gmail)", GREEN, "Daily rollup send (mail.daily_rollup). From the napco-nucleus@ael-bd.com alias."),
    ]
    y0 = 1.75
    for i, (name, color, desc) in enumerate(rows):
        y = Inches(y0 + i * 0.78)
        add_chip(s, Inches(0.7), y, Inches(2.9), Inches(0.62), name, color,
                 font_size=13)
        tf = add_textbox(s, Inches(3.8), y, Inches(8.7), Inches(0.62),
                         anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, desc, size=13, color=INK)
    set_speaker_notes(s, (
        "The integration surface and the auth rule that trips people up: "
        "do NOT add ANTHROPIC_API_KEY as a secret. Claude is invoked "
        "through the Claude Agent SDK, which authenticates off the local "
        "Claude Code CLI session (Max-tier subscription), so the marginal "
        "cost of a run is zero. Call transcription is Google STT only "
        "(tools/google_stt.py), Chirp 2 model for best Bangla coverage; it "
        "returns Bangla verbatim and Claude translates downstream during "
        "identify. faster-whisper was removed 2026-06-08; Groq is not used "
        "in production. Email and Drive ingest are roster/allowlist gated. "
        "The tracker is OpenProject (openproject_client.py replaced the "
        "deleted gitlab_client.py). Secrets live in .env + "
        "google-credentials.json at the repo root, both gitignored, loaded "
        "with override=True."
    ))


def arch_repo_map(prs):
    s = base_slide(prs, "Repo map - where to look",
                   "These entry points cover the whole production path.")
    rows = [
        ("collect_central.py", "LIVE orchestrator. Aggregates the central tree for one (client, day)."),
        ("agent.py", "--task verify_session. The identify turn (Claude Agent SDK, Opus 4.7)."),
        ("prompts/", "The brains. system.md + task prompts. Read verify_session / requirement first."),
        ("tools/requirements.py", "MCP tools, incl. publish_tasks_to_backlog -> OpenProject."),
        ("tools/google_stt.py", "The sole transcription engine (Chirp 2)."),
        ("openproject_client.py", "Work-Package create/update. Project mvp-access."),
        ("mail/daily_rollup.py", "The once-a-day email + Requirements Verification .docx."),
        ("deploy/linux-central/", "docker-compose.yml + the loop scripts. The whole runtime."),
    ]
    y0 = 1.7
    for i, (path, desc) in enumerate(rows):
        y = Inches(y0 + i * 0.62)
        add_rect(s, Inches(0.7), y, Inches(3.8), Inches(0.5),
                 fill_color=SOFT, line_color=RULE)
        tf_p = add_textbox(s, Inches(0.85), y, Inches(3.6), Inches(0.5),
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_p, path, size=13, bold=True, color=NAVY, font=MONO)
        tf_d = add_textbox(s, Inches(4.7), y, Inches(7.9), Inches(0.5),
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_d, desc, size=12.5, color=INK)
    set_speaker_notes(s, (
        "A reading order grounded in the production path. collect_central.py "
        "is the live orchestrator - start there and follow what it calls. "
        "It pulls email/Drive, walks the central tree for calls + chat + "
        "attachments, transcribes calls via tools/google_stt.py, builds a "
        "unified pull-session doc, then shells out to agent.py --task "
        "verify_session for the identify turn. That turn uses the MCP tools "
        "in tools/requirements.py, whose publish_tasks_to_backlog writes "
        "OpenProject Work Packages via openproject_client.py. Finally "
        "mail/daily_rollup.py sends the single daily email. The whole "
        "runtime is deploy/linux-central/. Ignore README / "
        "docs/requirement-management-flow.md for now - they describe an "
        "older GitLab/Whisper generation and are being rewritten."
    ))


def build_architecture():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    title_slide(prs, "NAPCO Nucleus",
                "Architecture - the production system, as it runs today.",
                None,
                "Developer onboarding  |  Adaptive Enterprise / NAPCO labs")
    set_speaker_notes(prs.slides[0], (
        "Technical onboarding deck for an engineer joining or replicating "
        "Nucleus. Everything here is verified against the code, not the "
        "stale README/flow docs. Goal: by the end they can name every "
        "production moving part and know where to read first. Pair with the "
        "Workflow deck for the run-by-run sequence."
    ))
    arch_mental_model(prs)
    arch_topology(prs)
    arch_containers(prs)
    arch_memory(prs)
    arch_external(prs)
    arch_repo_map(prs)
    closing_slide(prs, "That's the architecture.",
                  "Next: read one daily run, step by step (Workflow deck).",
                  [("Live orchestrator", "collect_central.py --client all --last-minutes 1440"),
                   ("Identify turn", "agent.py --task verify_session   (Opus 4.7)"),
                   ("Runtime", "deploy/linux-central/docker-compose.yml"),
                   ("Tracker", "OpenProject  ->  project mvp-access")])
    OUT_ARCH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_ARCH))
    print(f"Wrote: {OUT_ARCH}  ({len(prs.slides)} slides)")


# ===================================================================
# DECK 2 - WORKFLOW (one production day, end to end)
# ===================================================================

def flow_one_day(prs):
    s = base_slide(prs, "One production day",
                   "Capture all day. Think once, at 23:30 BD. Email once.")
    events = [
        ("all day", PURPLE, "Capture",
         "Dev PCs record Teams calls + push chat to the central share."),
        ("every 2 min", TEAL, "Transcribe",
         "nucleus-transcribe turns finished calls into text (Google STT)."),
        ("every 15 min", CORAL, "Stage",
         "nucleus-stage-email + -drive pull Gmail and Drive into central."),
        ("23:30 BD", NAVY, "Identify",
         "nucleus-daily-draft runs collect_central -> agent verify_session."),
        ("23:30 BD", GREEN, "Publish + email",
         "OpenProject Work Packages created; ONE rollup email sent."),
    ]
    line_x = Inches(3.0)
    add_rect(s, line_x, Inches(1.8), Inches(0.04), Inches(4.6),
             fill_color=RULE)
    y0 = 1.7
    step = 0.98
    for i, (when, color, kind, body) in enumerate(events):
        y = Inches(y0 + i * step)
        tf_t = add_textbox(s, Inches(0.7), y, Inches(2.1), Inches(0.5))
        set_text(tf_t, when, size=16, bold=True, color=color,
                 align=PP_ALIGN.RIGHT)
        add_chip(s, line_x - Inches(0.1), y + Inches(0.05),
                 Inches(0.25), Inches(0.25), "", color)
        add_chip(s, Inches(3.3), y, Inches(2.3), Inches(0.5), kind, color,
                 font_size=14)
        tf_b = add_textbox(s, Inches(5.8), y, Inches(6.9), Inches(0.5),
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_b, body, size=13.5, color=INK)
    tf = add_textbox(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.5))
    set_text(tf,
             "Continuous capture; a single nightly think-and-send. "
             "No per-call emails (policy: one a day).",
             size=13, color=MUTED, align=PP_ALIGN.CENTER)
    set_speaker_notes(s, (
        "The production rhythm. Capture is continuous and dumb - dev PCs "
        "record calls and push chat; the transcribe loop (every 2 min) and "
        "the email/Drive stagers (every 15 min) move everything into the "
        "central tree. None of that calls the LLM. The thinking happens "
        "exactly once a day: at 23:30 BD daily-draft-loop.sh runs "
        "collect_central.py --client all --last-minutes 1440, which "
        "aggregates the day and invokes the identify turn, publishes to "
        "OpenProject, and sends one rollup email. The single-daily-email "
        "policy (commit 779034d, 2026-06-11) replaced an earlier per-call "
        "trigger - the trigger file is still written but deliberately "
        "ignored."
    ))


def flow_steps(prs):
    s = base_slide(prs, "Inside the nightly run",
                   "What collect_central.py actually does, in order.")
    steps = [
        ("1", "Pull email + Drive", "mail.pull_email + drive.pull_drive -> central tree (roster-filtered).", PURPLE),
        ("2", "Walk the central tree", "Gather every dev's calls, chat .docx, and chat attachments for the day.", TEAL),
        ("3", "Transcribe calls", "Google STT (Chirp 2) on mic + speaker WAVs. You = dev, Other = client.", CORAL),
        ("4", "Build pull-session doc", "Stitch calls + chat + email + Drive + attachments into one ordered doc.", GOLD),
        ("5", "Identify (LLM)", "agent.py --task verify_session. Opus 4.7 reads the doc, finds requirements.", NAVY),
        ("6", "Publish + remember", "publish_tasks_to_backlog -> OpenProject WPs; remember_requirement -> memory.", PURPLE),
        ("7", "Daily rollup email", "mail.daily_rollup: one email + Requirements Verification .docx attached.", GREEN),
    ]
    y0 = 1.7
    for i, (num, title, desc, color) in enumerate(steps):
        y = Inches(y0 + i * 0.7)
        add_chip(s, Inches(0.7), y, Inches(0.55), Inches(0.5), num, color,
                 font_size=16)
        tf_t = add_textbox(s, Inches(1.45), y, Inches(3.0), Inches(0.5),
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_t, title, size=15, bold=True, color=color)
        tf_d = add_textbox(s, Inches(4.6), y, Inches(8.0), Inches(0.5),
                           anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_d, desc, size=12, color=INK)
    set_speaker_notes(s, (
        "The real sequence inside collect_central.py. Steps 1-2 are "
        "ingest/aggregation - pull fresh email and Drive, then walk every "
        "developer's folder on the central share for the day's calls, chat "
        "docs, and attachments. Step 3 transcribes the call WAVs via Google "
        "STT (the mic track is labelled You = the AEL dev, the speaker "
        "track is Other = the client, where requirements live). Step 4 "
        "builds the unified pull-session document, ordered in time and "
        "grouped by client. Step 5 is the only LLM step: it shells out to "
        "agent.py --task verify_session, which runs one Opus 4.7 turn over "
        "the doc via the Agent SDK + MCP tools. Step 6 publishes each "
        "requirement as an OpenProject Work Package and records it in "
        "memory for dedup. Step 7 is mail.daily_rollup sending the single "
        "email with the curated .docx."
    ))


def flow_output(prs):
    s = base_slide(prs, "Output & delivery",
                   "Two destinations: the tracker, and one email to the team.")
    add_box(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(4.2),
            "OpenProject Work Packages", [
                "openproject_client.py -> project mvp-access",
                "",
                "Type: Task (requirement) / Bug",
                "Status: New",
                "Category: AccessGroup / BadgeHolder / Personnel",
                "",
                "Dedup before create:",
                "  open-WP title match + requirements_seen",
                "",
                "Each new WP's IID + URL -> memory,",
                "so the next night skips it.",
            ], accent=PURPLE, title_size=16, body_size=12.5)
    add_box(s, Inches(6.85), Inches(1.8), Inches(5.9), Inches(4.2),
            "The daily rollup email", [
                "mail.daily_rollup, once at 23:30 BD",
                "",
                "TO:  NUCLEUS_ROLLUP_TO",
                "CC:  NUCLEUS_ROLLUP_CC (+ reqs-only CC)",
                "From: napco-nucleus@ael-bd.com alias",
                "",
                "Subject: 'NAPCO Nucleus - daily client",
                "  requirements (YYYY-MM-DD)'",
                "  -> '[ACTION NEEDED]' if identify failed",
                "",
                "Body = requirement TITLES only.",
                "Attached: Requirements Verification <date>.docx",
            ], accent=GREEN, title_size=16, body_size=12.5)
    set_speaker_notes(s, (
        "Where requirements land. First, the tracker: "
        "publish_tasks_to_backlog (tools/requirements.py) creates "
        "OpenProject Work Packages via openproject_client.py in project "
        "mvp-access - Type Task for requirements, Bug for bugs, Status New, "
        "categorized AccessGroup / BadgeHolder / Personnel. Every candidate "
        "is deduped against open WP titles and requirements_seen before "
        "create, and the new IID + URL is remembered. Second, delivery: "
        "mail.daily_rollup sends exactly one email at 23:30 BD to "
        "NUCLEUS_ROLLUP_TO with NUCLEUS_ROLLUP_CC (plus a requirements-only "
        "CC list that's used only when requirements were found), from the "
        "napco-nucleus@ael-bd.com verified alias. The body lists "
        "requirement titles only; full descriptions live in the attached "
        "Requirements Verification .docx. If sources were collected but the "
        "identify step failed, the subject flips to [ACTION NEEDED] so a "
        "silent failure can't hide."
    ))


def flow_guardrails(prs):
    s = base_slide(prs, "Guardrails",
                   "What keeps an autonomous nightly run safe.")
    cards = [
        ("Idempotency", NAVY,
         ["IMAP since-UID checkpoint", "Drive file-ID never re-run",
          "2-layer WP dedupe"]),
        ("Dry-run", PURPLE,
         ["dry_run input / env flag", "Every step EXCEPT mutations",
          "No send, no WP create"]),
        ("One email / day", TEAL,
         ["Per-call trigger ignored", "Single 23:30 BD rollup",
          "(commit 779034d)"]),
        ("Roster allowlist", CORAL,
         ["REQUIREMENT_ROSTER filter", "Only roster senders ingested",
          "Random mail dropped"]),
        ("Memory cap", GOLD,
         ["daily-draft capped at 4 GB", "Fences runaway LLM calls",
          "docker-compose limit"]),
        ("Failure shout", GREEN,
         ["Coverage note records what", "was processed; identify-fail",
          "flips subject to ACTION NEEDED"]),
    ]
    cw = Inches(3.95); ch = Inches(2.15); gx = 0.2; gy = 0.18
    x0 = 0.55; y0 = 1.7
    for i, (t, c, body) in enumerate(cards):
        row, col = divmod(i, 3)
        x = Inches(x0 + col * (3.95 + gx))
        y = Inches(y0 + row * (2.15 + gy))
        add_box(s, x, y, cw, ch, t, body, accent=c,
                title_size=16, body_size=12.5)
    set_speaker_notes(s, (
        "Six guardrails that make the nightly run trustworthy. Idempotency "
        "through checkpoints + multi-layer dedupe means a re-run is always "
        "safe. Dry-run (workflow_dispatch dry_run input -> "
        "NAPCO_NUCLEUS_DRY_RUN=1) runs everything except mutations - no "
        "SMTP, no Work-Package create - while still logging. The "
        "one-email-per-day policy means the per-call trigger is "
        "deliberately consumed-and-ignored. The roster allowlist "
        "(napco_config.REQUIREMENT_ROSTER) drops senders not in the working "
        "group. nucleus-daily-draft is memory-capped at 4 GB to fence a "
        "runaway LLM loop. And the coverage note + [ACTION NEEDED] subject "
        "ensure a run that collected sources but failed to identify can't "
        "send a silent 'no requirements' email."
    ))


def flow_run_it(prs):
    s = base_slide(prs, "Run it yourself",
                   "Three ways to trigger production. Start with a dry run.")
    add_box(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.3),
            "Trigger a run", [
                "Manual, via GitHub Actions:",
                "  Actions -> Requirement Management",
                "  -> Run workflow (last_minutes / client",
                "     / day / dry_run)",
                "  -> docker exec nucleus-daily-draft ...",
                "",
                "On-demand from a dev PC:",
                "  py -3 do_it_now.py --client 'X' \\",
                "    --last-minutes 60",
                "",
                "On the host directly:",
                "  python collect_central.py --client all \\",
                "    --last-minutes 1440 [--dry-run]",
            ], accent=NAVY, title_size=16, body_size=12.5, body_font=MONO)
    add_box(s, Inches(6.85), Inches(1.8), Inches(5.9), Inches(4.3),
            "Inspect & verify", [
                "Verify the install:",
                "  py -3 -m tools.healthcheck",
                "",
                "Replay the last run's trace:",
                "  py -3 -m tools.replay_trace --latest",
                "  (data/traces/<date>/<run_id>.jsonl)",
                "",
                "Always dry-run first:",
                "  --dry-run sets NAPCO_NUCLEUS_DRY_RUN=1",
                "  -> no email, no Work-Package create,",
                "     but memory still logs the run.",
            ], accent=GREEN, title_size=16, body_size=12.5, body_font=MONO)
    set_speaker_notes(s, (
        "Hands-on. Three trigger paths, all hitting the same production "
        "code. (1) GitHub Actions workflow_dispatch on "
        "requirement-management.yml - inputs last_minutes / client / day / "
        "dry_run, executed as docker exec nucleus-daily-draft python "
        "collect_central.py on the self-hosted runner. (2) do_it_now.py "
        "from a dev PC, which pushes local chat then SSHes to the host to "
        "run collect_central for one client. (3) Directly on the host. "
        "Always start with --dry-run, which sets NAPCO_NUCLEUS_DRY_RUN=1 so "
        "tools short-circuit every mutation (no SMTP, no WP create) while "
        "still logging to memory. tools.healthcheck verifies config; "
        "tools.replay_trace --latest shows the JSONL trace of the most "
        "recent run (every run writes data/traces/<date>/<run_id>.jsonl)."
    ))


def build_workflow():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    title_slide(prs, "NAPCO Nucleus",
                "Workflow - one production day, end to end.",
                None,
                "Developer onboarding  |  Adaptive Enterprise / NAPCO labs")
    set_speaker_notes(prs.slides[0], (
        "Companion to the Architecture deck. Walks a single production day: "
        "continuous capture, then the once-a-day 23:30 BD think-and-send. "
        "Everything verified against collect_central.py, agent.py, "
        "mail/daily_rollup.py and the docker-compose stack - not the stale "
        "flow diagram."
    ))
    flow_one_day(prs)
    flow_steps(prs)
    flow_output(prs)
    flow_guardrails(prs)
    flow_run_it(prs)
    closing_slide(prs, "You can run it now.",
                  "Start with a dry run, then read the trace.",
                  [("Dry run", "collect_central.py --client all --last-minutes 1440 --dry-run"),
                   ("Health", "py -3 -m tools.healthcheck"),
                   ("Trace", "py -3 -m tools.replay_trace --latest"),
                   ("Runtime", "deploy/linux-central/docker-compose.yml")])
    OUT_FLOW.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FLOW))
    print(f"Wrote: {OUT_FLOW}  ({len(prs.slides)} slides)")


# -- build ----------------------------------------------------------

if __name__ == "__main__":
    build_architecture()
    build_workflow()
