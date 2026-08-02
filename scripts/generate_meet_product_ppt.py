"""Nucleus Meeting Assistant for Google Meet - product deck.

Audience: a stakeholder who has NOT seen anything we have built before. The
deck stands completely on its own: it presents a Google Meet meeting assistant
as a product, and never refers to any other platform or any earlier version.

Product voice, not proposal voice. The arc is: what it is, what it delivers,
how it works, how the requirements are produced, what safeguards are built in,
what it runs on, and how a meeting gets covered day to day. No approval asks,
no delivery stages, no cost figures.

Sending model (Titu, 2026-07-30): the summary goes out automatically by email
to the meeting's own agreed recipient list. There is NO human review or
approval step, so nothing in this deck may claim one. The email carries the
summary in the body and the identified core requirements as an attachment, and
anything left unanswered in the meeting is listed as an open question.

Rules followed (Titu's standing preferences):
  * plain operational language, no em dashes
  * jargon explained where it appears, never assumed
  * no personal names, no machine names, no internal paths on slides
  * the recording disclosure is stated openly, not buried

Run:
    py -3 scripts\\generate_meet_product_ppt.py
Output:
    C:\\Users\\khasan\\Downloads\\Nucleus-Meet-Assistant.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(r"C:\Users\khasan\Downloads\Nucleus-Meet-Assistant.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x2E, 0x8A, 0x8A)
CORAL = RGBColor(0xE0, 0x78, 0x56)
GREEN = RGBColor(0x4A, 0x7A, 0x4A)
GOLD = RGBColor(0xC9, 0x96, 0x2B)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x77, 0x85)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
RULE = RGBColor(0xD5, 0xDC, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line=None,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.25)
    s.shadow.inherit = False
    return s


def _text(slide, x, y, w, h, text, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(space)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return box


def _header(slide, title, kicker=None):
    _rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill=NAVY,
          shape=MSO_SHAPE.RECTANGLE)
    _text(slide, Inches(0.7), Inches(0.16), Inches(11.9), Inches(0.5),
          title, size=29, bold=True, color=WHITE)
    if kicker:
        _text(slide, Inches(0.72), Inches(0.71), Inches(11.9), Inches(0.34),
              kicker, size=13, color=RGBColor(0xC9, 0xDA, 0xEA))


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _card(slide, x, y, w, h, heading, body, accent, hsize=16, bsize=13):
    _rect(slide, x, y, w, h, fill=SOFT, line=RULE)
    _rect(slide, x, y, Inches(0.07), h, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5),
          Inches(0.4), heading, size=hsize, bold=True, color=NAVY)
    _text(slide, x + Inches(0.28), y + Inches(0.66), w - Inches(0.5),
          h - Inches(0.8), body, size=bsize, color=INK, space=3)


def _arrow(slide, x, y, w=Inches(0.42), h=Inches(0.34), color=MUTED):
    return _rect(slide, x, y, w, h, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)


def _bullet_row(slide, x, y, w, text, accent=TEAL, size=15):
    """One bullet line: a small square marker then the text."""
    _rect(slide, x, y + Inches(0.09), Inches(0.14), Inches(0.14), fill=accent,
          shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.34), y, w - Inches(0.34), Inches(0.42),
          text, size=size, color=INK, space=0)


def _stat(slide, x, y, w, h, big, label, accent=TEAL):
    """A short fact tile: one strong line and a caption under it."""
    _rect(slide, x, y, w, h, fill=SOFT, line=RULE)
    _rect(slide, x, y, w, Inches(0.06), fill=accent, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.24), y + Inches(0.32), w - Inches(0.45),
          Inches(0.5), big, size=17, bold=True, color=NAVY)
    _text(slide, x + Inches(0.24), y + Inches(0.92), w - Inches(0.45),
          h - Inches(1.0), label, size=13, color=INK, space=3)


# ----------------------------------------------------------------- slides ---
def slide_title(prs):
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    _rect(s, 0, Inches(4.55), SLIDE_W, Inches(0.06), fill=GOLD,
          shape=MSO_SHAPE.RECTANGLE)
    _text(s, Inches(1.1), Inches(2.1), Inches(11.4), Inches(1.2),
          "Nucleus Meeting Assistant", size=52, bold=True, color=WHITE)
    _text(s, Inches(1.13), Inches(3.05), Inches(11.4), Inches(0.6),
          "A virtual colleague that attends our Google Meet calls",
          size=24, color=WHITE)
    _text(s, Inches(1.13), Inches(3.7), Inches(11.2), Inches(0.9),
          "It listens to the meeting, writes down what was agreed, and sends "
          "the requirements to the team.\n"
          "Nothing to install, nothing to run, no notes to type.",
          size=18, color=RGBColor(0xC9, 0xDA, 0xEA), space=4)
    _text(s, Inches(1.13), Inches(4.9), Inches(11), Inches(0.6),
          "Product overview  |  August 2026",
          size=16, color=RGBColor(0x9F, 0xB8, 0xCE))
    _notes(s, "Frame it in one line before moving on: the meeting happens as "
              "normal, and the written record arrives afterwards without "
              "anyone doing it by hand.")


def slide_what_it_is(prs):
    s = _blank(prs)
    _header(s, "What it is",
            "One participant in the meeting that does the writing")
    _rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.5),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(1.82), Inches(11.2), Inches(1.2),
          "The assistant is a virtual colleague with its own place in the "
          "meeting. It runs on its own dedicated machine, not on anyone's "
          "laptop, so it does not matter whose computer is on, who is "
          "presenting, or who has to leave early. It attends the Google Meet "
          "call, keeps the audio safe, and turns the conversation into a "
          "written requirement record for the team.",
          size=15, color=INK, space=4)

    tiles = [
        ("Attends by itself", "Joins the scheduled Meet call without anyone "
                              "setting it up first.", TEAL),
        ("Bangla and English", "Understands a call that switches between the "
                               "two, which most tools do not.", TEAL),
        ("Runs on our own server", "Recordings and documents stay on our "
                                   "infrastructure, not a third party notes "
                                   "service.", GREEN),
        ("Sends by itself", "The summary reaches the meeting's own recipient "
                            "list without anyone forwarding it.", GREEN),
    ]
    x = Inches(0.7)
    for big, label, col in tiles:
        _stat(s, x, Inches(3.45), Inches(2.9), Inches(2.35), big, label, col)
        x += Inches(3.07)
    _notes(s, "The last two tiles are the ones people ask about first: where "
              "the data lives, and who ends up receiving the summary. The "
              "answer to the second is the list agreed for that meeting, "
              "nobody outside it.")


def slide_delivers(prs):
    s = _blank(prs)
    _header(s, "What it delivers",
            "One email after the meeting, with everything behind it attached")
    items = [
        ("The written transcript",
         "The full conversation as readable text, so a discussion can be "
         "checked word for word instead of argued from memory.", TEAL),
        ("The core requirements",
         "Only the things that were actually asked for, written as clear "
         "requirements and attached to the email as a document.", TEAL),
        ("The open questions",
         "Anything left unclear or unanswered in the meeting, listed "
         "separately so it can be taken back to the client.", CORAL),
        ("The summary email",
         "A short summary of the meeting, emailed on its own to the team "
         "members agreed for that meeting, requirements attached.", GOLD),
    ]
    x = Inches(0.7)
    for head, body, col in items:
        _card(s, x, Inches(1.85), Inches(2.9), Inches(2.9), head, body, col,
              hsize=15, bsize=13)
        x += Inches(3.07)
    _rect(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.1),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.42), Inches(11.2), Inches(0.6),
          "Every meeting also stays on file, so a discussion from weeks ago can "
          "still be checked word for word.",
          size=18, bold=True, color=NAVY)
    _notes(s, "If they only remember one slide, this is the one. The open "
              "questions column is the part people do not expect: it does not "
              "only record what was agreed, it flags what still has no "
              "answer.")


def slide_how_it_works(prs):
    s = _blank(prs)
    _header(s, "How it works",
            "Four steps, none of which need anything from the people on the "
            "call")
    _rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.35),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(1.8), Inches(11.2), Inches(1.05),
          "The audio reaches our server while the call is still going on, not "
          "after it ends, so nothing is lost if a laptop is shut the moment the "
          "meeting finishes. Everything after that happens on its own, with no "
          "one asking for it.",
          size=15, color=INK, space=4)

    steps = [
        ("1", "Joins", "Enters the meeting at the\nscheduled time."),
        ("2", "Captures", "Audio reaches our server\nwhile the call runs."),
        ("3", "Reads", "Speech becomes text, then\nrequirements are pulled out."),
        ("4", "Emails", "Summary to the team, with\nrequirements attached."),
    ]
    x = Inches(0.7)
    for num, head, body in steps:
        _rect(s, x, Inches(3.3), Inches(2.75), Inches(2.2), fill=WHITE,
              line=RULE)
        _rect(s, x, Inches(3.3), Inches(2.75), Inches(0.5), fill=TEAL,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, x + Inches(0.2), Inches(3.38), Inches(2.4), Inches(0.4),
              f"{num}.  {head}", size=15, bold=True, color=WHITE)
        _text(s, x + Inches(0.25), Inches(4.0), Inches(2.3), Inches(1.4),
              body, size=13, color=INK, space=3)
        if num != "4":
            _arrow(s, x + Inches(2.85), Inches(4.2))
        x += Inches(3.05)
    _text(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.5),
          "Meeting records are filed in a folder kept only for Meet calls, "
          "separate from everything else.",
          size=15, bold=True, color=NAVY)
    _notes(s, "Step 2 is the part that sounds small and is not. Audio moving "
              "during the call is why a closed laptop or a dropped connection "
              "does not cost us the meeting.")


def slide_inside(prs):
    s = _blank(prs)
    _header(s, "How the requirements are produced",
            "From recorded speech to a checked list, step by step")
    stages = [
        ("Speech to text",
         "The recording is turned into text, handling Bangla, English, and a "
         "call that moves between them."),
        ("Noise removed",
         "Greetings, side conversation and small talk are dropped so only "
         "meaningful discussion is read."),
        ("Requirements identified",
         "What was actually asked for is separated from what was merely "
         "discussed, and anything left unanswered becomes an open question."),
        ("Checked and split",
         "Each requirement is checked against the transcript and broken into "
         "pieces small enough to plan and assign."),
        ("Emailed to the team",
         "The summary goes out to the members agreed for that meeting, with "
         "the requirements attached, and no one needing to send it."),
    ]
    y = Inches(1.8)
    for i, (head, body) in enumerate(stages):
        col = GREEN if i == len(stages) - 1 else TEAL
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(0.92), fill=SOFT,
              line=RULE)
        _rect(s, Inches(0.7), y, Inches(0.07), Inches(0.92), fill=col,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, Inches(1.0), y + Inches(0.13), Inches(3.1), Inches(0.4),
              f"{i + 1}.  {head}", size=15, bold=True, color=NAVY)
        _text(s, Inches(4.3), y + Inches(0.14), Inches(8.1), Inches(0.7),
              body, size=14, color=INK, space=2)
        y += Inches(1.02)
    _notes(s, "The noise removal step is worth mentioning. Reading everything "
              "including the small talk is what makes most automatic summaries "
              "vague, so it is filtered before the requirements are pulled "
              "out.")


def slide_safeguards(prs):
    s = _blank(prs)
    _header(s, "What is built in",
            "The parts that decide whether a team actually trusts it")
    left = [
        "Summaries go only to the members agreed for that meeting",
        "Recordings and documents stay on our own server",
        "Meeting audio is kept apart from every other source",
        "It only attends the meetings it is invited to",
    ]
    right = [
        "Nothing is installed on anyone's computer",
        "Nothing runs in the background on your machine",
        "It keeps working if a laptop is closed right after the call",
        "Every summary carries a note that the meeting was recorded",
    ]
    y = Inches(1.85)
    for i, item in enumerate(left):
        _bullet_row(s, Inches(0.75), y + i * Inches(0.72), Inches(5.6), item,
                    accent=GREEN)
    for i, item in enumerate(right):
        _bullet_row(s, Inches(6.95), y + i * Inches(0.72), Inches(5.6), item,
                    accent=GREEN)

    _rect(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.55),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.22), Inches(11.2), Inches(0.4),
          "On recording, stated plainly", size=16, bold=True, color=CORAL)
    _text(s, Inches(1.05), Inches(5.66), Inches(11.2), Inches(0.9),
          "Meetings are recorded so they can be turned into text. Where Google "
          "does the recording, everyone on the call sees Google's own notice. "
          "Where the assistant records, that notice does not appear, so "
          "participants are told directly and every summary says so in "
          "writing.",
          size=14, color=INK, space=3)
    _notes(s, "Say the recording paragraph out loud rather than leaving it on "
              "the slide. Being first to raise it is what keeps it from "
              "becoming an objection later.")


def slide_setup(prs):
    s = _blank(prs)
    _header(s, "What it runs on",
            "One machine, our own server, and two ways to capture a meeting")
    _card(s, Inches(0.7), Inches(1.8), Inches(3.9), Inches(2.8),
          "Its own machine",
          "The assistant has one dedicated always on machine. It is not "
          "installed on anyone's laptop and it does not need anyone to be "
          "logged in for it to work.", TEAL, hsize=17, bsize=14)
    _card(s, Inches(4.85), Inches(1.8), Inches(3.9), Inches(2.8),
          "Google records, we process",
          "Where the Google plan includes meeting recording, Google records the "
          "call and saves it to Drive. The assistant collects it from there, "
          "and everyone sees Google's normal recording notice.", GREEN,
          hsize=17, bsize=14)
    _card(s, Inches(9.0), Inches(1.8), Inches(3.6), Inches(2.8),
          "Or it records itself",
          "Where recording is not part of the plan, the assistant records the "
          "meeting itself. This works on any plan and does not depend on "
          "anyone pressing record.", GOLD, hsize=17, bsize=14)

    _rect(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(1.35),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.14), Inches(11.2), Inches(0.4),
          "Either way, the output is the same", size=16, bold=True, color=NAVY)
    _text(s, Inches(1.05), Inches(5.56), Inches(11.2), Inches(0.6),
          "The two capture modes only change where the recording comes from. "
          "The transcript, the requirement document and the summary are "
          "produced the same way in both cases.",
          size=14, color=INK, space=3)
    _notes(s, "Do not turn this into a technical choice for the audience. The "
              "point is that the product covers both cases and the result does "
              "not change.")


def slide_using_it(prs):
    s = _blank(prs)
    _header(s, "Using it day to day",
            "Three things happen, and only the first involves you")
    rows = [
        ("You", "Invite the assistant to the meeting, the same way you would "
                "invite a colleague."),
        ("The assistant", "Joins at the scheduled time, stays for the whole "
                          "call, and prepares the transcript, the requirement "
                          "document and the summary."),
        ("The team", "Receives the summary by email with the requirements "
                     "attached, without anyone writing it or forwarding it."),
    ]
    y = Inches(1.95)
    for i, (who, what) in enumerate(rows):
        col = TEAL if i != 2 else GREEN
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(1.15), fill=SOFT,
              line=RULE)
        _rect(s, Inches(0.7), y, Inches(2.5), Inches(1.15), fill=col)
        _text(s, Inches(0.95), y + Inches(0.34), Inches(2.1), Inches(0.5),
              who, size=17, bold=True, color=WHITE)
        _text(s, Inches(3.45), y + Inches(0.22), Inches(8.9), Inches(0.8),
              what, size=15, color=INK, space=2)
        y += Inches(1.28)

    _text(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.6),
          "Nothing else changes. The meeting runs exactly as it does today.",
          size=19, bold=True, color=NAVY)
    _text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5),
          "Happy to run it live on a real meeting whenever you want to see it.",
          size=16, color=INK)
    _notes(s, "Close here. The offer of a live meeting is more convincing than "
              "any slide, so make it the last thing said.")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_what_it_is(prs)
    slide_delivers(prs)
    slide_how_it_works(prs)
    slide_inside(prs)
    slide_safeguards(prs)
    slide_setup(prs)
    slide_using_it(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"[OK] {len(prs.slides)} slides -> {OUT}")


if __name__ == "__main__":
    build()
