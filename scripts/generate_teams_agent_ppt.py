"""NAPCO Nucleus - team briefing deck for the MS Teams AI agent.

Audience: the development team and team leads. Explains what the agent does,
how it works, the first architecture and why we moved off it, what the new
setup gives everyone, and the ask: start using it, it is available on demand.

Rules followed (Titu's standing preferences):
  * plain operational language, no em dashes
  * jargon is explained where it appears, never assumed
  * no personal names in slide-visible text
  * claims limited to what was actually proven on 2026-07-27

Run:
    py -3 scripts\\generate_teams_agent_ppt.py
Output:
    C:\\Users\\khasan\\Desktop\\NAPCO-Nucleus-Teams-Agent.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(r"C:\Users\khasan\Desktop\NAPCO-Nucleus-Teams-Agent.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x2E, 0x8A, 0x8A)
CORAL = RGBColor(0xE0, 0x78, 0x56)
GREEN = RGBColor(0x4A, 0x7A, 0x4A)
GOLD = RGBColor(0xC9, 0x96, 0x2B)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x77, 0x85)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
RULE = RGBColor(0xD5, 0xDC, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line=None,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.25)
    s.shadow.inherit = False
    return s


def _text(slide, x, y, w, h, text, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(space)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return box


def _header(slide, title, kicker=None):
    _rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill=NAVY,
          shape=MSO_SHAPE.RECTANGLE)
    _text(slide, Inches(0.7), Inches(0.16), Inches(11.9), Inches(0.5),
          title, size=29, bold=True, color=WHITE)
    if kicker:
        _text(slide, Inches(0.72), Inches(0.71), Inches(11.9), Inches(0.34),
              kicker, size=13, color=RGBColor(0xC9, 0xDA, 0xEA))


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _card(slide, x, y, w, h, heading, body, accent, hsize=16, bsize=13):
    _rect(slide, x, y, w, h, fill=SOFT, line=RULE)
    _rect(slide, x, y, Inches(0.07), h, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5),
          Inches(0.4), heading, size=hsize, bold=True, color=NAVY)
    _text(slide, x + Inches(0.28), y + Inches(0.66), w - Inches(0.5),
          h - Inches(0.8), body, size=bsize, color=INK, space=3)


def _arrow(slide, x, y, w=Inches(0.42), h=Inches(0.34), color=MUTED):
    return _rect(slide, x, y, w, h, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)


def _bullet_row(slide, x, y, w, text, accent=TEAL, size=15):
    """One bullet line: a small square marker then the text."""
    _rect(slide, x, y + Inches(0.09), Inches(0.14), Inches(0.14), fill=accent,
          shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.34), y, w - Inches(0.34), Inches(0.42),
          text, size=size, color=INK, space=0)


# ----------------------------------------------------------------- slides ---
def slide_title(prs):
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    _rect(s, 0, Inches(4.55), SLIDE_W, Inches(0.06), fill=GOLD,
          shape=MSO_SHAPE.RECTANGLE)
    _text(s, Inches(1.1), Inches(2.4), Inches(11), Inches(1.1),
          "Napco Nucleus", size=54, bold=True, color=WHITE)
    _text(s, Inches(1.12), Inches(3.4), Inches(11.2), Inches(1.1),
          "Our AI agent now sits in MS Teams and needs no installation on "
          "anyone's PC.\n"
          "Just add this virtual colleague to your call to start using it.\n"
          "It does everything else on its own and updates you by Teams message.",
          size=19, color=RGBColor(0xC9, 0xDA, 0xEA), space=4)
    _text(s, Inches(1.12), Inches(4.9), Inches(11), Inches(0.6),
          "Team briefing  |  August 2026",
          size=16, color=RGBColor(0x9F, 0xB8, 0xCE))
    _notes(s, "Short session. What it does, why the setup changed, what you "
              "get, and the ask at the end. Offer a live demo.")


def slide_what_it_does(prs):
    s = _blank(prs)
    _header(s, "What it does",
            "It attends the call and does the writing that nobody has time for")
    left = [
        "Joins our Teams calls as a member, like any of us",
        "Records the discussion",
        "Transcribes it automatically, in Bangla and English",
        "Picks out the requirements we agreed during the call",
    ]
    right = [
        "Writes them into a requirement document",
        "Emails the summary to the team",
        "Replies to messages in Teams chat",
        "Follows up on what it said it would check",
    ]
    y = Inches(1.85)
    for i, item in enumerate(left):
        _bullet_row(s, Inches(0.75), y + i * Inches(0.72), Inches(5.6), item)
    for i, item in enumerate(right):
        _bullet_row(s, Inches(6.95), y + i * Inches(0.72), Inches(5.6), item)

    _rect(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.0),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.58), Inches(11.2), Inches(0.6),
          "What we discuss verbally no longer depends on someone remembering "
          "to write it down afterwards.",
          size=18, bold=True, color=NAVY)
    _notes(s, "Read the list quickly. The line at the bottom is the point.")


def slide_how_it_works(prs):
    s = _blank(prs)
    _header(s, "How it works", "One machine attends, the server does the rest")
    _rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.5),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(1.82), Inches(11.2), Inches(1.2),
          "The agent sits on one dedicated machine and attends the call as a "
          "participant. While the call is still going on, the audio is already "
          "being copied to our central server, so nothing is lost if a laptop "
          "is closed the moment the call ends. After the call it transcribes "
          "the recording, reads the transcript, identifies what was actually "
          "asked for, and prepares the requirement document and the summary.",
          size=15, color=INK, space=4)

    steps = [
        ("1", "Attends", "Added to the call like any\nother participant."),
        ("2", "Captures", "Audio reaches the server\nwhile the call is running."),
        ("3", "Transcribes", "Speech becomes text,\nBangla and English."),
        ("4", "Delivers", "Requirement document\nand summary to the team."),
    ]
    x = Inches(0.7)
    for num, head, body in steps:
        _rect(s, x, Inches(3.45), Inches(2.75), Inches(2.2), fill=WHITE,
              line=RULE)
        _rect(s, x, Inches(3.45), Inches(2.75), Inches(0.5), fill=TEAL,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, x + Inches(0.2), Inches(3.53), Inches(2.4), Inches(0.4),
              f"{num}.  {head}", size=15, bold=True, color=WHITE)
        _text(s, x + Inches(0.25), Inches(4.15), Inches(2.3), Inches(1.4),
              body, size=13, color=INK, space=3)
        if num != "4":
            _arrow(s, x + Inches(2.85), Inches(4.35))
        x += Inches(3.05)
    _notes(s, "Nothing here needs anything from the team. That is the whole "
              "design goal of the current setup.")


def slide_before(prs):
    s = _blank(prs)
    _header(s, "The first setup, and why we changed it",
            "It depended on software running on every developer's own PC")
    items = [
        ("Installed on each PC",
         "The agent had to be installed on every developer's own machine, "
         "where it captured that person's calls.", CORAL),
        ("Most people said no",
         "Understandably, most of the team was not comfortable installing it "
         "on a personal machine.", CORAL),
        ("It only worked while that PC was on",
         "The machine had to stay switched on and unlocked, so capture "
         "stopped the moment it was not.", CORAL),
    ]
    x = Inches(0.7)
    for head, body, col in items:
        _card(s, x, Inches(1.9), Inches(3.9), Inches(2.6), head, body, col)
        x += Inches(4.15)
    _text(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(0.9),
          "Very few installations happened, so coverage was poor and it "
          "depended entirely on who had it running.",
          size=19, bold=True, color=NAVY)
    _notes(s, "Do not make anyone defensive. The refusal was reasonable and "
              "the new design is the answer to it.")


def slide_benefits(prs):
    s = _blank(prs)
    _header(s, "Why the new setup is better",
            "One machine attends the call, so nothing lands on your PC")
    left = [
        "Nothing is installed on anyone's PC",
        "Nothing runs in the background on your machine",
        "Nobody has to keep anything open or unlocked",
        "Coverage no longer depends on who installed it",
    ]
    right = [
        "Recordings survive if you close the laptop after a call",
        "Transcription accuracy is much better",
        "Requirements are captured per call, not compiled later",
        "It confirms with you before sending anything for you",
    ]
    y = Inches(1.85)
    for i, item in enumerate(left):
        _bullet_row(s, Inches(0.75), y + i * Inches(0.72), Inches(5.6), item,
                    accent=GREEN)
    for i, item in enumerate(right):
        _bullet_row(s, Inches(6.95), y + i * Inches(0.72), Inches(5.6), item,
                    accent=GREEN)

    _rect(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.0),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(5.58), Inches(11.2), Inches(0.6),
          "The only thing asked of you is to add it to the call.",
          size=18, bold=True, color=GREEN)
    _notes(s, "This slide answers the objection from the previous slide, "
              "point by point.")


def slide_request(prs):
    s = _blank(prs)
    _header(s, "The request", "Use it, and use it whenever you need it")
    _card(s, Inches(0.7), Inches(1.85), Inches(5.85), Inches(2.5),
          "Please start using it",
          "Use it in your client calls and internal discussions. The more "
          "calls it attends, the fewer requirements we lose.", TEAL,
          hsize=18, bsize=15)
    _card(s, Inches(6.9), Inches(1.85), Inches(5.7), Inches(2.5),
          "Available on demand",
          "Anyone in the team can use it, any time. Add this virtual colleague "
          "to your call and it does everything else on its own, with no "
          "intervention from you, then updates you by Teams message.",
          GOLD, hsize=18, bsize=15)

    _rect(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(1.55),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(4.92), Inches(11.2), Inches(0.4),
          "One thing to say plainly", size=16, bold=True, color=CORAL)
    _text(s, Inches(1.05), Inches(5.36), Inches(11.2), Inches(0.8),
          "Calls are recorded so they can be transcribed, and Teams does not "
          "show its usual recording indicator for this. If you have any "
          "concern about it, please raise it directly.",
          size=14, color=INK, space=3)
    _notes(s, "Say the recording line out loud. Do not leave it only on the "
              "slide. Offer a live demo to anyone who wants one.")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_what_it_does(prs)
    slide_how_it_works(prs)
    slide_before(prs)
    slide_benefits(prs)
    slide_request(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"[OK] {len(prs.slides)} slides -> {OUT}")


if __name__ == "__main__":
    build()
