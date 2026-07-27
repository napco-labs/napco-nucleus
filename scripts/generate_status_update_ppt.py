"""NAPCO Nucleus - management status update deck (full version).

Audience: Titu's boss and anyone technical he brings along. Covers the business
case AND the architecture, the processing chain, the technology stack, and how
the data is handled.

Rules followed (Titu's standing preferences):
  * plain operational language, no em dashes
  * jargon is explained where it appears, never assumed
  * no personal names in slide-visible text
  * claims limited to what was actually proven on 2026-07-27

Run:
    py -3 scripts\\generate_status_update_ppt.py
Output:
    C:\\Users\\khasan\\Desktop\\NAPCO-Nucleus-Status-Update.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(r"C:\Users\khasan\Desktop\NAPCO-Nucleus-Status-Update.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1F, 0x4E, 0x79)
TEAL = RGBColor(0x2E, 0x8A, 0x8A)
CORAL = RGBColor(0xE0, 0x78, 0x56)
GREEN = RGBColor(0x4A, 0x7A, 0x4A)
GOLD = RGBColor(0xC9, 0x96, 0x2B)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x77, 0x85)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
RULE = RGBColor(0xD5, 0xDC, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line=None,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.25)
    s.shadow.inherit = False
    return s


def _text(slide, x, y, w, h, text, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(space)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return box


def _header(slide, title, kicker=None):
    _rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill=NAVY,
          shape=MSO_SHAPE.RECTANGLE)
    _text(slide, Inches(0.7), Inches(0.16), Inches(11.9), Inches(0.5),
          title, size=29, bold=True, color=WHITE)
    if kicker:
        _text(slide, Inches(0.72), Inches(0.71), Inches(11.9), Inches(0.34),
              kicker, size=13, color=RGBColor(0xC9, 0xDA, 0xEA))


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _card(slide, x, y, w, h, heading, body, accent, hsize=16, bsize=13):
    _rect(slide, x, y, w, h, fill=SOFT, line=RULE)
    _rect(slide, x, y, Inches(0.07), h, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5),
          Inches(0.4), heading, size=hsize, bold=True, color=NAVY)
    _text(slide, x + Inches(0.28), y + Inches(0.66), w - Inches(0.5),
          h - Inches(0.8), body, size=bsize, color=INK, space=3)


def _arrow(slide, x, y, w=Inches(0.42), h=Inches(0.34), color=MUTED):
    a = _rect(slide, x, y, w, h, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)
    return a


# ----------------------------------------------------------------- slides ---
def slide_title(prs):
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    _rect(s, 0, Inches(4.55), SLIDE_W, Inches(0.06), fill=GOLD,
          shape=MSO_SHAPE.RECTANGLE)
    _text(s, Inches(1.1), Inches(2.4), Inches(11), Inches(1.1),
          "NAPCO Nucleus", size=54, bold=True, color=WHITE)
    _text(s, Inches(1.12), Inches(3.45), Inches(11), Inches(0.9),
          "An AI colleague that captures client requirements from our calls",
          size=22, color=RGBColor(0xC9, 0xDA, 0xEA))
    _text(s, Inches(1.12), Inches(4.9), Inches(11), Inches(0.6),
          "Status update, architecture and technology  |  July 2026",
          size=16, color=RGBColor(0x9F, 0xB8, 0xCE))
    _notes(s, "Ten minutes. Business case first, then architecture and stack, "
              "then the ask. Skip to slide 5 if the audience is technical.")


def slide_problem(prs):
    s = _blank(prs)
    _header(s, "The problem we set out to solve",
            "Requirements live in conversations, and conversations are easy to lose")
    items = [
        ("Requirements arrive in calls",
         "Clients describe what they want on a call. Nobody can write everything "
         "down while also running the meeting.", CORAL),
        ("Notes are inconsistent",
         "What gets captured depends on who was on the call and how busy they "
         "were that day.", CORAL),
        ("Things get missed",
         "A missed request becomes rework, a delay, or an unhappy client, and we "
         "usually find out late.", CORAL),
    ]
    x = Inches(0.7)
    for head, body, col in items:
        _card(s, x, Inches(1.9), Inches(3.9), Inches(2.5), head, body, col)
        x += Inches(4.15)
    _text(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.9),
          "The information was never missing. It was simply never written down.",
          size=20, bold=True, color=NAVY)
    _notes(s, "Do not labour this. The boss knows the pain.")


def slide_what(prs):
    s = _blank(prs)
    _header(s, "What we built",
            "A teammate that sits in the call and does the writing")
    _rect(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(1.35),
          fill=SOFT, line=RULE)
    _text(s, Inches(1.05), Inches(1.9), Inches(11.2), Inches(1.0),
          "Napco Nucleus is added to a client call like any other colleague.\n"
          "It listens, writes up what the client asked for, and emails the team.",
          size=20, color=INK, space=6)
    for i, (h, b) in enumerate([
        ("It joins the call",
         "Someone adds it to the meeting. Nothing else changes for the team."),
        ("It does the writing",
         "Every request becomes a small, clear task with a time estimate."),
        ("It tells the team",
         "An email with the requirements document attached, and a chat notice."),
    ]):
        _card(s, Inches(0.7) + i * Inches(4.15), Inches(3.3), Inches(3.9),
              Inches(2.0), h, b, TEAL)
    _text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.6),
          "It also answers questions in Teams chat, in English or Bangla.",
          size=17, bold=True, color=TEAL)
    _notes(s, "Stress that nothing changes in how the team works. Adding it to "
              "a call is the entire ask of them.")


def slide_architecture(prs):
    s = _blank(prs)
    _header(s, "Architecture", "Three parts, all inside our own network")
    boxes = [
        ("The team", "Teams on their own PCs.\nNothing to install.\nThey add the "
         "assistant to a call.", MUTED),
        ("Assistant machine", "A dedicated Windows PC.\nSits in the call, records "
         "both sides,\nand runs the chat assistant.", TEAL),
        ("Central server", "A Linux server in the office.\nStores recordings, "
         "transcribes,\nand extracts the requirements.", NAVY),
    ]
    x = Inches(0.55)
    for head, body, col in boxes:
        _rect(s, x, Inches(1.9), Inches(3.6), Inches(2.5), fill=WHITE, line=RULE)
        _rect(s, x, Inches(1.9), Inches(3.6), Inches(0.5), fill=col,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, x + Inches(0.2), Inches(1.98), Inches(3.2), Inches(0.4),
              head, size=16, bold=True, color=WHITE)
        _text(s, x + Inches(0.25), Inches(2.6), Inches(3.1), Inches(1.6),
              body, size=13, color=INK, space=3)
        x += Inches(4.05)
    _arrow(s, Inches(4.2), Inches(3.0))
    _arrow(s, Inches(8.25), Inches(3.0))

    _rect(s, Inches(0.55), Inches(4.75), Inches(12.2), Inches(1.55),
          fill=SOFT, line=RULE)
    _text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.4),
          "Everything stays on the office network", size=17, bold=True,
          color=GREEN)
    _text(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.8),
          "Recordings and transcripts are stored on our own server. The only "
          "things that leave are the audio sent for transcription and the text "
          "sent for analysis, both to established cloud providers.",
          size=13.5, color=INK, space=3)
    _notes(s, "Assistant machine is a Windows box because it has to drive the "
              "Teams client and capture system audio. Central is Ubuntu running "
              "Docker. The two talk over a standard Windows file share.")


def slide_pipeline(prs):
    s = _blank(prs)
    _header(s, "How a call becomes a requirement", "Seven steps, all automatic")
    steps = [
        ("1", "Capture", "Both sides of the call are recorded as separate tracks."),
        ("2", "Transfer", "Audio streams to the central server while the call runs."),
        ("3", "Transcribe", "Speech becomes text, Bangla and English."),
        ("4", "Read", "Transcripts, chats, email and shared files are read together."),
        ("5", "Identify", "Real client requests are separated from general talk."),
        ("6", "Check", "Anything already captured before is not raised twice."),
        ("7", "Deliver", "A Word document is produced and emailed to the team."),
    ]
    x = Inches(0.35)
    w = Inches(1.72)
    for num, head, body in steps:
        _rect(s, x, Inches(2.0), w, Inches(3.0), fill=WHITE, line=RULE)
        _rect(s, x, Inches(2.0), w, Inches(0.42), fill=NAVY,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, x, Inches(2.05), w, Inches(0.35), "STEP " + num, size=11,
              bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.1), Inches(2.6), w - Inches(0.2), Inches(0.4),
              head, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.12), Inches(3.1), w - Inches(0.24), Inches(1.7),
              body, size=11.5, color=INK, align=PP_ALIGN.CENTER, space=2)
        x += Inches(1.83)
    _text(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.9),
          "Each requirement arrives as a small task of roughly four hours, with "
          "a title, a summary and an estimate.",
          size=16, bold=True, color=TEAL)
    _notes(s, "Steps 4 to 6 are where the AI does the real work. Deduplication "
              "matters: without it the same request would be raised every day "
              "it is mentioned.")


def slide_stack(prs):
    s = _blank(prs)
    _header(s, "Tools and technologies", "What each part is and why we use it")
    rows = [
        ("Microsoft Teams", "Where the calls and chats happen",
         "Already our meeting platform, so nothing new for the team.", NAVY),
        ("Google Speech to Text", "Turns speech into text",
         "Handles Bangla and English, and copes with call audio.", TEAL),
        ("Anthropic Claude", "Reads the text and finds the requirements",
         "Also powers the chat assistant that answers colleagues.", PURPLE),
        ("Python", "The automation itself",
         "Recording, transfer, scheduling and the chat assistant.", GREEN),
        ("Docker on Linux", "Runs the server side",
         "Keeps each service separate and restartable.", CORAL),
        ("GitHub", "Source control and automation",
         "Every change is tracked and reviewable.", GOLD),
    ]
    y = Inches(1.75)
    for name, what, why, col in rows:
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(0.82), fill=SOFT,
              line=RULE)
        _rect(s, Inches(0.7), y, Inches(0.09), Inches(0.82), fill=col,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, Inches(1.0), y + Inches(0.2), Inches(2.9), Inches(0.45),
              name, size=14.5, bold=True, color=NAVY)
        _text(s, Inches(4.0), y + Inches(0.2), Inches(3.5), Inches(0.45),
              what, size=13, color=INK)
        _text(s, Inches(7.6), y + Inches(0.2), Inches(4.8), Inches(0.45),
              why, size=12.5, color=MUTED)
        y += Inches(0.9)
    _notes(s, "If asked about cost: the only recurring charges are speech to "
              "text and the AI model, both usage based and small at our volume. "
              "Everything else is already owned or open source.")


def slide_data(prs):
    s = _blank(prs)
    _header(s, "How we handle the recordings", "Deliberate limits, set up front")
    items = [
        ("Stored on our own server",
         "Recordings and transcripts sit on an office machine, not in anyone's "
         "personal cloud storage.", GREEN),
        ("Only named people",
         "The assistant records calls involving a fixed list of seven colleagues "
         "and nobody else.", GREEN),
        ("Told in writing",
         "Every requirements email carries a plain statement that calls are "
         "recorded and why.", GREEN),
        ("Used for one purpose",
         "The recordings are used to produce the requirements documents. "
         "Nothing else.", GREEN),
    ]
    for i, (h, b, c) in enumerate(items):
        col, row = i % 2, i // 2
        _card(s, Inches(0.7) + col * Inches(6.15),
              Inches(1.9) + row * Inches(1.75),
              Inches(5.9), Inches(1.55), h, b, c, hsize=16, bsize=13)
    _text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.7),
          "Anyone can ask for a call not to be recorded, and it is switched off "
          "for that call.", size=16, bold=True, color=NAVY)
    _notes(s, "Raise this before anyone else does. It shows the limits were "
              "designed in rather than added after a complaint.")


def slide_live(prs):
    s = _blank(prs)
    _header(s, "Where it is today", "Working, on a real call")
    _rect(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(1.15),
          fill=GREEN, line=None)
    _text(s, Inches(1.05), Inches(1.88), Inches(11.2), Inches(0.8),
          "Live and proven end to end on 27 July 2026", size=24, bold=True,
          color=WHITE)
    checks = [
        ("Joins and records", "Both sides of the call captured cleanly."),
        ("Transfers", "Audio reaches the server while the call is running."),
        ("Transcribes", "Bangla and English, automatically."),
        ("Finds requirements", "Turns the conversation into sized tasks."),
        ("Emails the team", "With the requirements document attached."),
        ("Answers questions", "Colleagues ask it for status in Teams chat."),
        ("Follows up", "If it says it will check, it checks and comes back."),
        ("Looks after itself", "Checks daily and repairs common faults."),
    ]
    for i, (h, b) in enumerate(checks):
        col, row = i % 4, i // 4
        x = Inches(0.7) + col * Inches(3.1)
        y = Inches(3.1) + row * Inches(1.3)
        _rect(s, x, y, Inches(2.9), Inches(1.1), fill=SOFT, line=RULE)
        _text(s, x + Inches(0.2), y + Inches(0.14), Inches(2.5), Inches(0.35),
              h, size=13.5, bold=True, color=GREEN)
        _text(s, x + Inches(0.2), y + Inches(0.5), Inches(2.55), Inches(0.55),
              b, size=11.5, color=INK, space=2)
    _notes(s, "Proven today on a real call: recorded, transferred, transcribed "
              "by Google, and the requirement pipeline triggered automatically.")


def slide_value(prs):
    s = _blank(prs)
    _header(s, "What this gives us", "The point is not the technology")
    rows = [
        ("Nothing gets missed",
         "Every client request is written down, whether or not anyone "
         "remembered to note it.", GREEN),
        ("The team stops taking notes",
         "Engineers listen to the client instead of typing during the call.",
         TEAL),
        ("Work arrives ready to start",
         "Requests come through already broken into small, sized tasks.", NAVY),
        ("We have a record",
         "If a client asks what was agreed, we can show them.", GOLD),
    ]
    y = Inches(1.85)
    for head, body, col in rows:
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(1.1), fill=SOFT,
              line=RULE)
        _rect(s, Inches(0.7), y, Inches(0.09), Inches(1.1), fill=col,
              shape=MSO_SHAPE.RECTANGLE)
        _text(s, Inches(1.05), y + Inches(0.16), Inches(4.2), Inches(0.5),
              head, size=17, bold=True, color=NAVY)
        _text(s, Inches(5.4), y + Inches(0.2), Inches(7.0), Inches(0.7),
              body, size=14.5, color=INK)
        y += Inches(1.24)
    _notes(s, "If the boss remembers one slide, make it this one.")


def slide_next(prs):
    s = _blank(prs)
    _header(s, "What comes next", "Three things, in order")
    items = [
        ("1", "Improve the recording quality",
         "Audio from the room side is weaker than it should be, which limits how "
         "accurate the written record can be. Being worked on now.", CORAL),
        ("2", "Roll it out to the whole team",
         "It runs for one team member today. Extending it is mostly a matter of "
         "colleagues adding it to their calls.", TEAL),
        ("3", "Decide on a company Microsoft account",
         "About USD 6 a month moves the assistant off a personal account onto a "
         "company one. Lower risk, and the company owns it.", NAVY),
    ]
    y = Inches(1.9)
    for num, head, body, col in items:
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(1.45), fill=SOFT,
              line=RULE)
        _rect(s, Inches(1.0), y + Inches(0.4), Inches(0.6), Inches(0.6),
              fill=col, shape=MSO_SHAPE.OVAL)
        _text(s, Inches(1.0), y + Inches(0.48), Inches(0.6), Inches(0.45),
              num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(s, Inches(1.9), y + Inches(0.2), Inches(4.0), Inches(0.5),
              head, size=16.5, bold=True, color=NAVY)
        _text(s, Inches(6.0), y + Inches(0.24), Inches(6.4), Inches(1.0),
              body, size=13, color=INK, space=3)
        y += Inches(1.57)
    _notes(s, "Item 3 is the only decision needed. Keep the ask small.")


def slide_ask(prs):
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    _text(s, Inches(1.1), Inches(1.4), Inches(11), Inches(0.8),
          "What we are asking for", size=34, bold=True, color=WHITE)
    _rect(s, Inches(1.1), Inches(2.4), Inches(11), Inches(0.05), fill=GOLD,
          shape=MSO_SHAPE.RECTANGLE)
    asks = [
        "Approval to extend it to the rest of the project team.",
        "Around USD 6 a month for a company Microsoft account.",
        "Agreement that recorded calls stay on our own server.",
    ]
    y = Inches(2.95)
    for a in asks:
        _rect(s, Inches(1.1), y, Inches(0.13), Inches(0.13), fill=GOLD,
              shape=MSO_SHAPE.OVAL)
        _text(s, Inches(1.6), y - Inches(0.14), Inches(10.4), Inches(0.6),
              a, size=20, color=WHITE)
        y += Inches(0.9)
    _text(s, Inches(1.1), Inches(6.0), Inches(11), Inches(0.6),
          "Everything else is already built and running.",
          size=17, color=RGBColor(0xC9, 0xDA, 0xEA))
    _notes(s, "Close here. If pushed on cost, the figure is one licence in "
              "total, not one per person.")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in (slide_title, slide_problem, slide_what, slide_architecture,
               slide_pipeline, slide_stack, slide_data, slide_live,
               slide_value, slide_next, slide_ask):
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("wrote  : %s" % OUT)
    print("slides : %d" % len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
